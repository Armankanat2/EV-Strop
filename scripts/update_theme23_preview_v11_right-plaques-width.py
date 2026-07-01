from pathlib import Path
import shutil

from pptx import Presentation


SRC = Path(
    "presentations/ispring-course/module-01-stropovka-gruzov/live-preview/"
    "S042-S061_theme23_preview_2026-07-01_v10_targeted-layout-fix.pptx"
)
OUT = Path(
    "presentations/ispring-course/module-01-stropovka-gruzov/live-preview/"
    "S042-S061_theme23_preview_2026-07-01_v11_right-plaques-width.pptx"
)


def find_slide(prs: Presentation, code: str):
    for slide in prs.slides:
        texts = [shape.text for shape in slide.shapes if hasattr(shape, "text") and shape.text]
        if code in texts:
            return slide
    raise ValueError(f"Slide {code} not found")


def widen_to_match(slide, upper_idx: int, lower_idx: int) -> None:
    upper = slide.shapes[upper_idx]
    lower = slide.shapes[lower_idx]
    upper.left = lower.left
    upper.width = max(upper.width, lower.width)


def main() -> None:
    OUT.parent.mkdir(parents=True, exist_ok=True)
    shutil.copy2(SRC, OUT)

    prs = Presentation(OUT)

    for code, upper_idx, lower_idx in [
        ("S045", 23, 24),
        ("S046", 23, 24),
        ("S047", 25, 26),
        ("S048", 31, 32),
        ("S049", 23, 24),
        ("S050", 23, 24),
        ("S051", 23, 24),
        ("S052", 17, 18),
        ("S053", 23, 24),
        ("S054", 37, 38),
        ("S055", 23, 24),
        ("S056", 23, 24),
        ("S057", 23, 24),
        ("S058", 19, 20),
        ("S059", 32, 33),
        ("S060", 17, 18),
        ("S061", 17, 18),
    ]:
        slide = find_slide(prs, code)
        widen_to_match(slide, upper_idx, lower_idx)

    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    main()
