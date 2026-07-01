from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


OUT_PATH = Path(
    "presentations/ispring-course/module-01-stropovka-gruzov/live-preview/"
    "S042-S061_theme23_preview_2026-07-01_v06_ev-text_5-6-linear_media-plan_rules.pptx"
)

BLUE = RGBColor(0x00, 0x57, 0xFF)
DEEP_BLUE = RGBColor(0x00, 0x38, 0xB8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAPHITE = RGBColor(0x11, 0x11, 0x11)
DARK_GRAY = RGBColor(0x4B, 0x55, 0x63)
LIGHT_GRAY = RGBColor(0xE5, 0xE7, 0xEB)
PALE_BLUE = RGBColor(0xEA, 0xF0, 0xFF)
SOFT_BLUE = RGBColor(0xF8, 0xFA, 0xFC)
WARN_RED = RGBColor(0xC5, 0x1F, 0x1F)
SOFT_RED = RGBColor(0xFE, 0xF2, 0xF2)
AMBER = RGBColor(0xB4, 0x53, 0x09)
SOFT_AMBER = RGBColor(0xFF, 0xF7, 0xED)
SLATE = RGBColor(0x1F, 0x29, 0x37)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def set_font(run, size, bold=False, color=GRAPHITE):
    run.font.name = "Arial"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_textbox(
    slide,
    x,
    y,
    w,
    h,
    text,
    size=18,
    bold=False,
    color=GRAPHITE,
    align=PP_ALIGN.LEFT,
):
    box = slide.shapes.add_textbox(x, y, w, h)
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Pt(6)
    frame.margin_right = Pt(6)
    frame.margin_top = Pt(4)
    frame.margin_bottom = Pt(4)
    p = frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    set_font(run, size, bold=bold, color=color)
    return box


def add_bullets(
    slide,
    x,
    y,
    w,
    h,
    items,
    size=16,
    color=GRAPHITE,
    bullet_color=BLUE,
    spacing=1.1,
):
    box = slide.shapes.add_textbox(x, y, w, h)
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Pt(8)
    frame.margin_right = Pt(6)
    frame.margin_top = Pt(2)
    frame.margin_bottom = Pt(2)
    for i, item in enumerate(items):
        p = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
        p.text = f"• {item}"
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = spacing
        for run in p.runs:
            set_font(run, size, color=color)
        if p.runs:
            p.runs[0].font.color.rgb = bullet_color
    return box


def add_rect(slide, x, y, w, h, fill=WHITE, line=LIGHT_GRAY, line_width=1.2, rounded=False):
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if rounded else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(line_width)
    return shape


def add_chip(slide, x, y, w, h, text, fill=PALE_BLUE, color=DEEP_BLUE, line=None):
    chip = add_rect(slide, x, y, w, h, fill=fill, line=line or fill, line_width=0.8)
    chip.line.color.rgb = line or fill
    add_textbox(
        slide,
        x + Inches(0.04),
        y + Inches(0.02),
        w - Inches(0.08),
        h - Inches(0.04),
        text,
        size=9,
        bold=True,
        color=color,
        align=PP_ALIGN.CENTER,
    )
    return chip


def add_title_block(slide, slide_id, label, title, subtitle=None, code_fill=DEEP_BLUE, code_color=WHITE):
    add_textbox(slide, Inches(0.55), Inches(0.28), Inches(3.4), Inches(0.22), label, size=10, bold=True, color=DARK_GRAY)
    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.55), Inches(0.58), Inches(1.25), Inches(0.045))
    bar.fill.solid()
    bar.fill.fore_color.rgb = BLUE
    bar.line.fill.background()
    add_chip(slide, Inches(11.25), Inches(0.28), Inches(1.25), Inches(0.32), slide_id, fill=code_fill, color=code_color, line=code_fill)
    title_size = 22 if len(title) > 34 else 26
    title_height = Inches(0.78) if len(title) > 34 else Inches(0.55)
    add_textbox(slide, Inches(0.55), Inches(0.78), Inches(7.6), title_height, title, size=title_size, bold=True)
    if subtitle:
        add_textbox(slide, Inches(0.55), Inches(1.4), Inches(7.4), Inches(0.62), subtitle, size=13, color=DARK_GRAY)


def add_footer(slide, text):
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.55), Inches(7.08), Inches(12.2), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = LIGHT_GRAY
    line.line.fill.background()
    add_textbox(slide, Inches(0.55), Inches(7.12), Inches(5.0), Inches(0.2), text, size=9, color=DARK_GRAY)


def new_slide(prs, slide_id, label, title, subtitle=None, footer_text=None, code_fill=DEEP_BLUE, code_color=WHITE):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_block(slide, slide_id, label, title, subtitle, code_fill=code_fill, code_color=code_color)
    add_footer(slide, footer_text or label.replace("•", "•"))
    return slide


def add_panel(slide, x, y, w, h, header, bullets=None, accent=DEEP_BLUE, fill=WHITE, level="Опора"):
    add_rect(slide, x, y, w, h, fill=fill, line=LIGHT_GRAY)
    top = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, y, w, Inches(0.04))
    top.fill.solid()
    top.fill.fore_color.rgb = accent
    top.line.fill.background()
    add_chip(slide, x + Inches(0.14), y + Inches(0.11), Inches(0.95), Inches(0.25), level, fill=PALE_BLUE if fill != SOFT_RED else SOFT_RED, color=accent, line=PALE_BLUE if fill != SOFT_RED else SOFT_RED)
    add_textbox(slide, x + Inches(0.14), y + Inches(0.43), w - Inches(0.28), Inches(0.46), header, size=15, bold=True)
    if bullets:
        add_bullets(slide, x + Inches(0.14), y + Inches(0.95), w - Inches(0.28), h - Inches(1.08), bullets, size=12, bullet_color=accent)


def add_placeholder_panel(slide, x, y, w, h, header, placeholder_lines, accent=DEEP_BLUE):
    add_rect(slide, x, y, w, h, fill=SOFT_BLUE, line=LIGHT_GRAY)
    top = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, y, w, Inches(0.04))
    top.fill.solid()
    top.fill.fore_color.rgb = accent
    top.line.fill.background()
    add_chip(slide, x + Inches(0.14), y + Inches(0.11), Inches(1.1), Inches(0.25), "Медиа", fill=WHITE, color=accent, line=LIGHT_GRAY)
    add_textbox(slide, x + Inches(0.14), y + Inches(0.43), w - Inches(0.28), Inches(0.42), header, size=15, bold=True)
    add_bullets(slide, x + Inches(0.14), y + Inches(0.98), w - Inches(0.28), h - Inches(1.12), placeholder_lines, size=12, color=DARK_GRAY, bullet_color=accent)


def add_step_card(slide, x, y, w, h, num, title, text):
    add_rect(slide, x, y, w, h, fill=WHITE, line=LIGHT_GRAY)
    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, y, w, Inches(0.03))
    bar.fill.solid()
    bar.fill.fore_color.rgb = BLUE
    bar.line.fill.background()
    add_textbox(slide, x + Inches(0.12), y + Inches(0.08), Inches(0.5), Inches(0.28), f"{num:02d}", size=14, bold=True, color=DEEP_BLUE)
    add_textbox(slide, x + Inches(0.12), y + Inches(0.35), w - Inches(0.24), Inches(0.36), title, size=12, bold=True)
    add_textbox(slide, x + Inches(0.12), y + Inches(0.73), w - Inches(0.24), h - Inches(0.84), text, size=10, color=DARK_GRAY)


def add_stop_badge(slide, x, y, w, h, text):
    add_rect(slide, x, y, w, h, fill=SOFT_RED, line=SOFT_RED)
    add_textbox(slide, x + Inches(0.1), y + Inches(0.08), w - Inches(0.2), h - Inches(0.16), text, size=13, bold=True, color=WARN_RED, align=PP_ALIGN.CENTER)


def add_table_text(slide, x, y, w, row_h, rows, accent=AMBER):
    col_widths = [Inches(2.0), Inches(2.0), Inches(2.0), Inches(2.7)]
    headers = ["Условие", "Кто контролирует", "Что запрещено", "Когда стоп"]
    cx = x
    for idx, header in enumerate(headers):
        cell = add_rect(slide, cx, y, col_widths[idx], row_h, fill=SOFT_AMBER, line=LIGHT_GRAY)
        cell.line.width = Pt(1.0)
        add_textbox(slide, cx + Inches(0.06), y + Inches(0.04), col_widths[idx] - Inches(0.12), row_h - Inches(0.08), header, size=10, bold=True, color=accent, align=PP_ALIGN.CENTER)
        cx += col_widths[idx]
    for r, row in enumerate(rows, start=1):
        cx = x
        top = y + row_h * r
        for idx, value in enumerate(row):
            add_rect(slide, cx, top, col_widths[idx], row_h, fill=WHITE, line=LIGHT_GRAY)
            add_textbox(slide, cx + Inches(0.06), top + Inches(0.04), col_widths[idx] - Inches(0.12), row_h - Inches(0.08), value, size=9, color=GRAPHITE)
            cx += col_widths[idx]


def build_deck():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    theme_label = "Тема 2.3 • Производство работ"
    footer_text = "ЭВ Стропальщик • Тема 2.3 • preview"

    slide = new_slide(
        prs,
        "S042",
        theme_label,
        "Тема 2.3. Производство работ",
        "Опасные и специальные условия по Игумнову: зона риска, запреты, ППР и ЛЭП.",
        footer_text,
    )
    add_panel(slide, Inches(0.55), Inches(2.0), Inches(5.65), Inches(3.95), "Что входит в блок", [
        "Как читать пространство работ и где проходит опасная зона.",
        "Какие простые нарушения чаще всего приводят к тяжелым последствиям.",
        "Почему специальные условия и ЛЭП требуют отдельного порядка работ.",
    ], level="Вход")
    add_placeholder_panel(slide, Inches(6.45), Inches(2.0), Inches(5.8), Inches(3.95), "Основной формат: статичный ключевой слайд", [
        "(изображение: рабочая площадка с краном и зоной перемещения)",
        "(схема: контур опасной зоны и траектория груза)",
        "(акцент: EV-оформление без перегруза деталями)",
        "(без интерактива: это вводный экран)",
    ])

    slide = new_slide(prs, "S043", theme_label, "Что будет в подтеме", "Маршрут блока по Игумнову и по принятой структуре курса.", footer_text)
    step_texts = [
        "Опасная зона и границы риска на площадке.",
        "Требования безопасности и обязательные запреты.",
        "Специальные условия и организация работ по схеме.",
        "Базовые правила работы стреловыми кранами вблизи ЛЭП.",
        "Усиление 5-6: руководство операцией, документы, сложный подъем и ЛЭП.",
        "Разбор нарушений и промежуточная проверка.",
    ]
    for idx, text in enumerate(step_texts):
        row = idx // 3
        col = idx % 3
        x = Inches(0.55) + col * Inches(3.95)
        y = Inches(2.0) + row * Inches(1.65)
        width = Inches(3.65) if col < 2 else Inches(4.0)
        add_step_card(slide, x, y, width, Inches(1.42), idx + 1, f"Шаг {idx + 1}", text)
    add_placeholder_panel(slide, Inches(8.65), Inches(5.35), Inches(3.6), Inches(1.45), "Формат: схема маршрута с озвучкой", [
        "(слайдшоу с озвучкой: маршрут подтемы по шагам)"
    ])

    slide = new_slide(prs, "S044", theme_label, "Урок 1. Опасная зона", "Игумнов: опасная зона должна быть определена заранее, а не по месту.", footer_text)
    add_panel(slide, Inches(0.55), Inches(2.0), Inches(5.8), Inches(4.15), "Что нужно понять", [
        "Опасная зона есть всегда.",
        "Ее определяют до начала работы.",
        "Под грузом и в зоне возможного падения находиться нельзя.",
        "Границы риска зависят и от груза, и от площадки, и от расположения крана.",
    ], level="Урок 1")
    add_placeholder_panel(slide, Inches(6.55), Inches(2.0), Inches(5.7), Inches(4.15), "Основной формат: слайдшоу с озвучкой", [
        "(слайдшоу: площадка -> кран -> груз -> опасная зона)",
        "(слайд: стройгенплан с выделением зоны риска)",
        "(изображение: ограждения и знаки безопасности)",
    ])
    add_chip(slide, Inches(0.55), Inches(6.35), Inches(2.05), Inches(0.28), "Есть подвал: S044-P01", fill=PALE_BLUE, color=DEEP_BLUE, line=PALE_BLUE)

    slide = new_slide(prs, "S044-P01", "Тема 2.3 • Подвал", "Подвал. Как читать опасную зону", "Расширение по Игумнову: ППРк, стройгенплан, ограждения и знаки.", footer_text)
    add_panel(slide, Inches(0.55), Inches(2.0), Inches(5.65), Inches(4.2), "Что проверяем до начала работ", [
        "Где зона показана в ППРк и на стройгенплане.",
        "Как обозначены ограждения и знаки безопасности.",
        "Где люди, проходы, техника и точки возможного падения груза.",
        "Можно ли считать расположение крана и груза безопасным.",
    ], level="Подвал")
    add_placeholder_panel(slide, Inches(6.45), Inches(2.0), Inches(5.8), Inches(4.2), "Основной формат: слайдшоу с озвучкой", [
        "(слайдшоу: ППРк -> стройгенплан -> ограждения -> проходы)",
        "(схема: разметка опасной зоны с пояснениями)",
        "(изображение: учебная площадка с комментариями)",
    ])

    slide = new_slide(prs, "S045", theme_label, "Урок 1. Как понимать опасную обстановку", "Опасность складывается из траектории, раскачивания, помех и людей вокруг.", footer_text)
    add_panel(slide, Inches(0.55), Inches(2.0), Inches(5.75), Inches(4.1), "На что смотреть", [
        "Смотри на траекторию груза, а не только на точку зацепки.",
        "Учитывай возможное раскачивание и разворот.",
        "Если груз проходит над препятствиями, нужен зазор не менее 500 мм.",
        "Проверяй ограниченное пространство, проемы и помехи.",
        "Думай о людях вокруг и о том, куда они могут попасть.",
    ], level="Контроль")
    add_panel(slide, Inches(6.5), Inches(2.0), Inches(2.65), Inches(4.1), "Короткий чек-лист", [
        "Траектория",
        "Раскачивание",
        "Помехи",
        "Люди",
    ], accent=BLUE, level="Чек-лист")
    add_placeholder_panel(slide, Inches(9.35), Inches(2.0), Inches(2.9), Inches(4.1), "Основной формат: слайдшоу с озвучкой", [
        "(слайдшоу: траектория -> раскачивание -> помехи -> люди)",
        "(схема: опасная ситуация)",
        "(схема: груз проходит над препятствием с зазором не менее 500 мм)",
    ])

    slide = new_slide(prs, "S046", theme_label, "Урок 2. Требования безопасности при производстве работ", "Игумнов собирает запреты и безопасные правила в один обязательный рабочий блок.", footer_text)
    add_panel(slide, Inches(0.55), Inches(2.0), Inches(3.8), Inches(4.3), "Базовые требования", [
        "Работай только в безопасных условиях.",
        "Соблюдай организацию работ и указания ответственных лиц.",
        "При плохой видимости или сложной связи работай через назначенного сигнальщика.",
        "Работы не начинают и прекращают при сильном ветре, грозе, тумане, снегопаде и резком ухудшении видимости.",
        "Работы прекращают, если температура ниже указанной в паспорте или если кран и грузозахватные приспособления неисправны.",
        "Не продолжай операцию, если условия не ясны.",
    ], level="База")
    add_panel(slide, Inches(4.55), Inches(2.0), Inches(3.6), Inches(4.3), "Что подчеркивает Игумнов", [
        "Большинство тяжелых происшествий связано с нарушением простых правил.",
        "Любые ускоряющие приемы обычно и создают риск.",
        "Стропальщик обязан остановить работу при опасности.",
    ], accent=BLUE, level="Смысл")
    add_placeholder_panel(slide, Inches(8.35), Inches(2.0), Inches(3.9), Inches(4.3), "Основной формат: статичный учебный слайд", [
        "(слайд: правила и иконки)",
        "(изображение: безопасная организация работ)",
        "(схема: маршрут и контроль)",
        "(часть изображений можно взять из плакатных схем)",
    ])

    slide = new_slide(prs, "S047", theme_label, "Урок 2. Что запрещено и когда работу нужно остановить", "Запреты, рывки, подтаскивание и любые действия без понятной схемы работ.", footer_text)
    add_panel(slide, Inches(0.55), Inches(2.0), Inches(5.2), Inches(4.2), "Что запрещено", [
        "Подтаскивать груз краном и делать рывки.",
        "Исправлять строповку под висящим грузом.",
        "Работать без понятной схемы, ППР или организации работ.",
        "Оставлять людей в опасной зоне, в кузове машины или в кабине при опускании груза.",
        "Перемещать стрелу над кабиной автомобиля.",
        "Работать без установки всех выносных опор.",
    ], accent=WARN_RED, fill=SOFT_RED, level="Запрет")
    add_stop_badge(slide, Inches(5.95), Inches(2.0), Inches(2.0), Inches(0.7), "При опасности — стоп")
    add_panel(slide, Inches(5.95), Inches(2.9), Inches(2.0), Inches(3.3), "Когда прекращаем работу", [
        "Нет безопасных условий.",
        "Не ясна организация работ.",
        "Есть риск для людей и оборудования.",
        "Погода, видимость или состояние крана не позволяют работать безопасно.",
    ], accent=WARN_RED, fill=WHITE, level="Стоп")
    add_placeholder_panel(slide, Inches(8.15), Inches(2.0), Inches(4.1), Inches(4.2), "Основной формат: предупреждающий слайд", [
        "(слайд: красный экран запретов)",
        "(схема: нельзя / почему опасно)",
        "(часть схем можно взять из плакатов по работе краном)",
    ], accent=WARN_RED)
    add_chip(slide, Inches(0.55), Inches(6.35), Inches(2.05), Inches(0.28), "Есть подвал: S047-P01", fill=SOFT_RED, color=WARN_RED, line=SOFT_RED)

    slide = new_slide(prs, "S047-P01", "Тема 2.3 • Подвал", "Подвал. Запрещенные действия при производстве работ", "Карточка в формате нельзя / почему опасно.", footer_text)
    add_panel(slide, Inches(0.55), Inches(2.0), Inches(5.65), Inches(4.25), "Нельзя", [
        "Подтаскивать груз краном.",
        "Исправлять строповку под висящим грузом.",
        "Продолжать работу при людях в опасной зоне.",
        "Импровизировать без понятной схемы и организации.",
    ], accent=WARN_RED, fill=SOFT_RED, level="Запрет")
    add_panel(slide, Inches(6.45), Inches(2.0), Inches(5.8), Inches(4.25), "Почему это опасно", [
        "Рывки и боковые усилия выводят груз из расчетной схемы.",
        "Под грузом и рядом с ним риск для людей становится мгновенным.",
        "Без схемы и организации стропальщик начинает действовать на глаз.",
        "Любой такой прием должен закончиться остановкой работ.",
    ], accent=DEEP_BLUE, level="Разбор")

    slide = new_slide(prs, "S048", theme_label, "Урок 3. Специальные условия работ", "Проемы, балконы, стесненные места и подача груза в сложной обстановке.", footer_text)
    labels = [
        ("Проемы", "Груз подают только при понятной траектории и безопасной организации."),
        ("Балконы", "Нельзя работать на глаз и без защищенного приема груза."),
        ("Стесненные места", "Нужны точная схема, контроль помех и понятная роль каждого."),
        ("Канава или котлован", "Положение крана и траекторию груза заранее сверяют со схемой площадки."),
    ]
    for idx, (head, body) in enumerate(labels):
        row = idx // 2
        col = idx % 2
        add_step_card(slide, Inches(0.55) + col * Inches(4.1), Inches(2.0) + row * Inches(1.55), Inches(3.8), Inches(1.3), idx + 1, head, body)
    add_placeholder_panel(slide, Inches(8.8), Inches(2.0), Inches(3.45), Inches(3.85), "Основной формат: слайдшоу с озвучкой", [
        "(слайдшоу: проем -> стесненное место -> канава или котлован)",
        "(схема: расположение крана относительно канавы или котлована)",
        "(интерактив: нажать на условие -> увидеть безопасную схему подачи)",
        "(без видео: здесь важнее схема, чем движение)",
    ])
    add_chip(slide, Inches(0.55), Inches(1.72), Inches(1.45), Inches(0.24), "Интерактив 1", fill=PALE_BLUE, color=DEEP_BLUE, line=PALE_BLUE)

    slide = new_slide(prs, "S049", theme_label, "Урок 3. Технологические карты и организация работ", "Игумнов: стропальщик работает не по памяти, а по схеме, ППР и понятному порядку.", footer_text)
    add_panel(slide, Inches(0.55), Inches(2.0), Inches(5.55), Inches(4.25), "Что должно быть понятно", [
        "По какой схеме выполняется операция.",
        "Какая роль у ППР и технологической карты.",
        "Какой маршрут перемещения, кто его контролирует и где обеспечен зазор не менее 500 мм над препятствиями.",
        "Где проходит запрещенная зона у поворотной части крана и обеспечен ли зазор не менее 1 м до строений, штабелей и конструкций.",
        "Где может стоять кран, если рядом канава, котлован или другая опасная граница площадки.",
        "Кто подает сигналы: стропальщик, сигнальщик или старший стропальщик.",
        "Кто отвечает за организацию и безопасность работ.",
    ], level="Организация")
    add_panel(slide, Inches(6.35), Inches(2.0), Inches(2.95), Inches(4.25), "Ключевой вывод", [
        "Не действуй самовольно.",
        "Если стропальщиков несколько, сигнал подает старший стропальщик.",
        "Не заменяй схему опытом.",
        "Не начинай работу при неясном порядке и неясной роли участников.",
    ], accent=BLUE, level="Вывод")
    add_placeholder_panel(slide, Inches(9.55), Inches(2.0), Inches(2.7), Inches(4.25), "Основной формат: слайдшоу с озвучкой", [
        "(слайдшоу: план -> маршрут -> контроль -> опасная граница площадки)",
        "(схема: кран, канава или котлован, безопасная позиция)",
        "(схема: поворотная часть крана и зона не менее 1 м; можно взять из плаката)",
        "(слайд: ППР и техкарта)",
    ])

    slide = new_slide(prs, "S050", theme_label, "Урок 4. Работа вблизи ЛЭП", "Игумнов: даже без касания ЛЭП работа может быть смертельно опасной.", footer_text)
    add_panel(slide, Inches(0.55), Inches(2.0), Inches(5.4), Inches(4.2), "Что нужно знать до начала работ", [
        "ЛЭП — особая опасность.",
        "Нужны специальные меры безопасности и понятные дистанции.",
        "Работу рядом с ЛЭП не ведут по обычной схеме и не оценивают на глаз.",
        "Если безопасные условия не разъяснены, работу не начинают.",
    ], accent=WARN_RED, level="ЛЭП")
    add_panel(slide, Inches(6.15), Inches(2.0), Inches(2.6), Inches(4.2), "Главная мысль", [
        "Нельзя работать на глаз.",
        "Нельзя действовать как обычно.",
        "При сомнении работу останавливают.",
    ], accent=WARN_RED, fill=SOFT_RED, level="Стоп")
    add_placeholder_panel(slide, Inches(8.95), Inches(2.0), Inches(3.3), Inches(4.2), "Основной формат: слайдшоу с озвучкой", [
        "(слайдшоу: кран -> ЛЭП -> зона риска -> решение стоп)",
        "(схема: кран рядом с ЛЭП)",
        "(предупреждающий визуал)",
    ], accent=WARN_RED)
    add_chip(slide, Inches(0.55), Inches(6.35), Inches(1.9), Inches(0.28), "Подвал: S050-P01", fill=PALE_BLUE, color=DEEP_BLUE, line=PALE_BLUE)
    add_chip(slide, Inches(2.6), Inches(6.35), Inches(3.35), Inches(0.28), "Далее: блок 5-6 разряда S051-S057", fill=SOFT_AMBER, color=AMBER, line=SOFT_AMBER)

    slide = new_slide(prs, "S050-P01", "Тема 2.3 • Подвал", "Подвал. ЛЭП: краткая памятка", "Краткий справочный экран по зоне риска, допуску и прекращению работ.", footer_text)
    add_panel(slide, Inches(0.55), Inches(2.0), Inches(5.8), Inches(4.25), "Памятка стропальщику", [
        "Работа рядом с ЛЭП допускается только по специальному порядку.",
        "Дистанции и границы безопасной зоны должны быть понятны всем участникам.",
        "Без разъясненной схемы и специальных мер безопасности работу не начинать.",
        "При угрозе приближения к токоведущим частям — немедленно стоп.",
    ], accent=WARN_RED, level="Памятка")
    add_placeholder_panel(slide, Inches(6.6), Inches(2.0), Inches(5.65), Inches(4.25), "Основной формат: статичная памятка", [
        "(слайд: зона риска у ЛЭП)",
        "(схема: дистанции и запретные позиции)",
        "(изображение: организованный порядок работ)",
    ], accent=WARN_RED)

    slide = new_slide(
        prs,
        "S051",
        "Тема 2.3 • Усиление 5-6",
        "5-6 разряд. Кто руководит сложной операцией",
        "Будущий старший стропальщик должен понимать общую схему руководства, а не только свою точку зацепки.",
        footer_text,
        code_fill=AMBER,
    )
    add_panel(slide, Inches(0.55), Inches(2.0), Inches(5.45), Inches(4.25), "Как устроено руководство операцией", [
        "За организацию и безопасное выполнение отвечает назначенное ответственное лицо.",
        "В сложной операции бригада работает по единой схеме и под единым управлением.",
        "Если стропальщиков несколько, команды крановщику подает один назначенный участник.",
        "Будущий старший стропальщик держит в поле зрения всю операцию: роли, маршрут, ограничения и очередность действий.",
    ], accent=AMBER, fill=SOFT_AMBER, level="5-6")
    add_panel(slide, Inches(6.25), Inches(2.0), Inches(2.75), Inches(4.25), "Кто есть кто", [
        "Ответственный за работы",
        "Крановщик",
        "Стропальщики",
        "Старший стропальщик",
    ], accent=AMBER, level="Роли")
    add_placeholder_panel(slide, Inches(9.2), Inches(2.0), Inches(3.05), Inches(4.25), "Основной формат: схема ролей", [
        "(схема: ответственный -> старший стропальщик -> крановщик)",
        "(слайд: роли в сложной операции)",
        "(без интерактива: не дробим этот экран)",
    ], accent=AMBER)

    slide = new_slide(
        prs,
        "S052",
        "Тема 2.3 • Усиление 5-6",
        "5-6 разряд. Что старший стропальщик делает до начала подъема",
        "До команды на подъем будущий старший стропальщик собирает операцию в рабочую систему.",
        footer_text,
        code_fill=AMBER,
    )
    add_panel(slide, Inches(0.55), Inches(2.0), Inches(5.55), Inches(4.35), "Что он проверяет заранее", [
        "Уточняет массу, схему строповки, центр тяжести и маршрут перемещения.",
        "Проверяет ППР, технологическую карту или другую рабочую схему, если она обязательна для операции.",
        "Сверяет, хватает ли оснастки, подкладок, оттяжек, проходов и места установки.",
        "Проверяет, что место приема груза подготовлено так же, как место подъема.",
        "До начала подъема распределяет роли: кто стропит, кто сопровождает, кто подает команды, кто контролирует прием груза.",
    ], accent=AMBER, fill=SOFT_AMBER, level="Подготовка")
    add_placeholder_panel(slide, Inches(6.35), Inches(2.0), Inches(5.9), Inches(4.35), "Основной формат: слайдшоу с озвучкой", [
        "(слайдшоу: чек-лист до начала сложного подъема)",
        "(схема: маршрут, опоры, место приема груза)",
        "(изображение: бригада на коротком инструктаже)",
    ], accent=AMBER)

    slide = new_slide(
        prs,
        "S053",
        "Тема 2.3 • Усиление 5-6",
        "5-6 разряд. Какие документы и роли обязательны",
        "Сложную операцию нельзя держать только в голове: она должна быть оформлена и разъяснена участникам.",
        footer_text,
        code_fill=AMBER,
    )
    add_panel(slide, Inches(0.55), Inches(2.0), Inches(5.35), Inches(4.35), "Что должно быть на руках и в понимании бригады", [
        "ППР или технологическая карта — если операция относится к сложным, ответственным или специально организуемым работам.",
        "Схема строповки и перемещения груза, меры безопасности и последовательность действий.",
        "Назначенные роли: ответственный за безопасное производство работ, крановщик, стропальщики, при необходимости старший стропальщик.",
        "Для специальных условий дополнительно проверяют, нужен ли наряд-допуск и отдельный инструктаж.",
    ], accent=AMBER, fill=SOFT_AMBER, level="Документы")
    add_panel(slide, Inches(6.1), Inches(2.0), Inches(2.55), Inches(4.35), "Документы не заменяют", [
        "Опыт",
        "Привычка",
        "Устные договоренности",
        "Импровизация",
    ], accent=AMBER, level="Нельзя")
    add_placeholder_panel(slide, Inches(8.9), Inches(2.0), Inches(3.35), Inches(4.35), "Основной формат: статичный документальный слайд", [
        "(слайд: ППР, техкарта, наряд-допуск)",
        "(схема: распределение ролей)",
    ], accent=AMBER)

    slide = new_slide(
        prs,
        "S054",
        "Тема 2.3 • Усиление 5-6",
        "5-6 разряд. Как организуют сложный и ответственный подъем",
        "Игумнов подчеркивает: чем сложнее операция, тем важнее не сила, а согласованность и точная последовательность.",
        footer_text,
        code_fill=AMBER,
    )
    add_step_card(slide, Inches(0.55), Inches(2.0), Inches(2.8), Inches(1.55), 1, "Разъясняют схему", "Все участники понимают способ строповки, маршрут, точки приема и ограничения.")
    add_step_card(slide, Inches(3.55), Inches(2.0), Inches(2.8), Inches(1.55), 2, "Готовят площадку", "Проверяют проходы, место установки, ограждение опасной зоны, подкладки и помехи.")
    add_step_card(slide, Inches(6.55), Inches(2.0), Inches(2.8), Inches(1.55), 3, "Согласуют команды", "До начала работ всем понятны сигналы, очередность действий и единый канал управления.")
    add_step_card(slide, Inches(9.55), Inches(2.0), Inches(2.7), Inches(1.55), 4, "Ведут подъем по схеме", "По ходу операции не меняют схему и не ускоряют подъем рывками.")
    add_panel(slide, Inches(0.55), Inches(3.9), Inches(5.8), Inches(2.4), "Что особенно важно для будущего старшего стропальщика", [
        "Уметь заранее увидеть, где операция может выйти из расчетной схемы: разворот груза, перекос, помеха, несогласованные действия.",
        "Собирать подъем как управляемую последовательность, а не как набор отдельных действий.",
    ], accent=AMBER, fill=SOFT_AMBER, level="Логика")
    add_placeholder_panel(slide, Inches(6.6), Inches(3.9), Inches(5.65), Inches(2.4), "Основной формат: слайдшоу с озвучкой", [
        "(слайдшоу: этап 1 -> этап 2 -> этап 3 -> этап 4)",
        "(интерактив: расставить этапы сложного подъема по порядку)",
        "(схема: сложный подъем по этапам)",
    ], accent=AMBER)
    add_chip(slide, Inches(0.55), Inches(1.72), Inches(1.45), Inches(0.24), "Интерактив 2", fill=SOFT_AMBER, color=AMBER, line=SOFT_AMBER)

    slide = new_slide(
        prs,
        "S055",
        "Тема 2.3 • Усиление 5-6",
        "5-6 разряд. Контрольный подъем и проверка перед перемещением",
        "Перед основным перемещением груз не «сразу в путь»: сначала подтверждают устойчивость и правильность схемы.",
        footer_text,
        code_fill=AMBER,
    )
    add_panel(slide, Inches(0.55), Inches(2.0), Inches(5.55), Inches(4.25), "Что проверяют на контрольном подъеме", [
        "Надежность строповки и равномерность нагрузки.",
        "Не уходит ли груз в перекос, разворот или нежелательное раскачивание.",
        "Нет ли помех по траектории и понятно ли, как груз пойдет дальше.",
        "Готова ли бригада сопровождать груз без входа в опасную зону.",
        "Если груз сразу показывает перекос или разворот, основное перемещение не начинают.",
    ], accent=AMBER, fill=SOFT_AMBER, level="Контроль")
    add_panel(slide, Inches(6.35), Inches(2.0), Inches(2.85), Inches(4.25), "Почему это важно", [
        "Ошибку нужно выявить на контрольном подъеме, а не во время перемещения.",
        "Именно здесь становится видно, верна ли схема в реальных условиях.",
        "Контрольный этап экономит время лучше, чем аварийная остановка.",
    ], accent=AMBER, level="Смысл")
    add_placeholder_panel(slide, Inches(9.45), Inches(2.0), Inches(2.8), Inches(4.25), "Основной формат: видео", [
        "(видео: контрольный подъем груза на малую высоту)",
        "(врезки-подписи: перекос, разворот, помеха, стоп)",
        "(резерв: если нет видео, собрать как псевдовидео из 4 кадров)",
    ], accent=AMBER)

    slide = new_slide(
        prs,
        "S056",
        "Тема 2.3 • Усиление 5-6",
        "5-6 разряд. Работа у ЛЭП: усиленный порядок и контроль",
        "Здесь добавляется не «еще одна осторожность», а отдельный организационный режим работ.",
        footer_text,
        code_fill=AMBER,
    )
    add_panel(slide, Inches(0.55), Inches(2.0), Inches(5.45), Inches(4.35), "Что усиливается по сравнению с базой", [
        "Работа у ЛЭП ведется по наряду-допуску и под непосредственным руководством ответственного лица.",
        "Место установки крана, границы безопасной зоны и порядок действий определяют заранее и доводят до всех участников.",
        "Кран нельзя самовольно устанавливать или переставлять ближе к ЛЭП.",
        "Старший стропальщик в такой операции удерживает дисциплину схемы и не допускает самовольных изменений.",
    ], accent=AMBER, fill=SOFT_AMBER, level="ЛЭП 5-6")
    add_panel(slide, Inches(6.2), Inches(2.0), Inches(2.8), Inches(4.35), "Что контролируют особенно жестко", [
        "Допуск",
        "Место установки",
        "Сигналы и порядок действий",
        "Запрет на работу «на глаз»",
    ], accent=AMBER, level="Контроль")
    add_placeholder_panel(slide, Inches(9.25), Inches(2.0), Inches(3.0), Inches(4.35), "Основной формат: слайдшоу с озвучкой", [
        "(слайдшоу: точка установки -> безопасная зона -> допуск -> стоп)",
        "(интерактив: можно / нельзя по ситуации у ЛЭП)",
        "(схема: кран, ЛЭП, безопасная зона, точка установки)",
    ], accent=AMBER)
    add_chip(slide, Inches(0.55), Inches(1.72), Inches(1.45), Inches(0.24), "Интерактив 3", fill=SOFT_AMBER, color=AMBER, line=SOFT_AMBER)

    slide = new_slide(
        prs,
        "S057",
        "Тема 2.3 • Усиление 5-6",
        "5-6 разряд. Чем будущий старший стропальщик отличается от аккуратного исполнителя",
        "Разница не в том, кто «смелее», а в том, кто умеет держать подъем как целую управляемую операцию.",
        footer_text,
        code_fill=AMBER,
    )
    add_panel(slide, Inches(0.55), Inches(2.0), Inches(5.65), Inches(4.35), "Будущий старший стропальщик", [
        "Читает операцию целиком, а не только свой участок работы.",
        "До начала подъема собирает роли, схему, маршрут и контрольные точки.",
        "Замечает, где ошибка организации может превратиться в ошибку строповки или перемещения.",
        "Умеет вести бригаду через понятную последовательность действий без суеты и самодеятельности.",
    ], accent=AMBER, fill=SOFT_AMBER, level="Профиль")
    add_panel(slide, Inches(6.45), Inches(2.0), Inches(2.8), Inches(4.35), "Аккуратный исполнитель", [
        "Видит в основном свой участок работы.",
        "Хорошо исполняет указания.",
        "Следит прежде всего за своей строповкой.",
        "Реже замечает организационные риски заранее.",
    ], accent=BLUE, level="Сравнение")
    add_placeholder_panel(slide, Inches(9.5), Inches(2.0), Inches(2.75), Inches(4.35), "Основной формат: статичное сравнение", [
        "(слайд: сравнение двух ролей)",
        "(схема: от исполнителя к старшему стропальщику)",
    ], accent=AMBER)

    slide = new_slide(prs, "S058", theme_label, "Урок 4. Когда работу нужно остановить из-за ЛЭП", "При малейшем сомнении работа прекращается до уточнения безопасных условий.", footer_text)
    add_stop_badge(slide, Inches(0.55), Inches(2.0), Inches(2.25), Inches(0.72), "Работа немедленно останавливается")
    add_panel(slide, Inches(0.55), Inches(2.95), Inches(5.65), Inches(3.8), "Стоп-условия", [
        "Нет безопасных условий.",
        "Не ясны дистанции и границы безопасной зоны.",
        "Есть риск опасного приближения к токоведущим частям.",
        "Схема работ не разъяснена участникам.",
    ], accent=WARN_RED, fill=SOFT_RED, level="Стоп")
    add_placeholder_panel(slide, Inches(6.45), Inches(2.0), Inches(5.8), Inches(4.75), "Плейсхолдер под предупреждающий визуал", [
        "(слайд: стоп-экран по ЛЭП)",
        "(схема: опасное сближение и запреты)",
        "(изображение: учебный кейс без реальной опасной съемки)",
    ], accent=WARN_RED)

    slide = new_slide(prs, "S059", theme_label, "Урок 5. Разбор ЧП и нарушений", "Кейс-логика: что произошло, что нарушили, как нужно было действовать.", footer_text)
    add_step_card(slide, Inches(0.55), Inches(2.0), Inches(3.6), Inches(1.5), 1, "Что произошло", "Коротко фиксируем ситуацию, зону риска и исходные условия.")
    add_step_card(slide, Inches(0.55), Inches(3.75), Inches(3.6), Inches(1.5), 2, "Что нарушили", "Показываем правило, которое было проигнорировано.")
    add_step_card(slide, Inches(0.55), Inches(5.5), Inches(3.6), Inches(1.0), 3, "Как правильно", "Формулируем безопасный порядок действий.")
    add_panel(slide, Inches(4.4), Inches(2.0), Inches(3.25), Inches(4.5), "Что важно в разборе", [
        "Не искать виноватого, а разобрать логику риска.",
        "Показывать, где операция должна была быть остановлена.",
        "Связывать кейс с правилом из Игумнова.",
    ], level="Разбор")
    add_placeholder_panel(slide, Inches(7.9), Inches(2.0), Inches(4.35), Inches(4.5), "Плейсхолдер", [
        "(кейс-слайд: что произошло / что нарушено)",
        "(изображение или схема аварийной ситуации)",
    ])

    slide = new_slide(prs, "S060", theme_label, "Промежуточный тест", "Проверка по опасной зоне, запретам, специальным условиям, ЛЭП и усилению 5-6.", footer_text)
    add_panel(slide, Inches(0.55), Inches(2.0), Inches(5.75), Inches(4.2), "Что проверяем", [
        "Где проходит опасная зона до начала работ.",
        "Как понять опасную обстановку.",
        "Когда работу нужно остановить.",
        "Правило: не менее 1 м между поворотной частью крана и строениями, штабелями и конструкциями.",
        "Когда работы запрещены из-за погоды, видимости, температуры и неисправности крана или оснастки.",
        "Какой подход верен рядом с ЛЭП.",
        "Чем для 5-6 разряда дополняются организация и контроль сложного подъема.",
    ], level="Тест")
    add_placeholder_panel(slide, Inches(6.55), Inches(2.0), Inches(5.7), Inches(4.2), "Плейсхолдер под тестовый экран", [
        "(слайд: тестовый интерфейс)",
        "(вопросы S060-Q01...Q05)",
        "(включить вопрос по правилу 1 м у поворотной части крана)",
        "(ошибка -> S060-P01...P05)",
    ])

    review_slides = [
        ("S060-P01", "Ошибка: опасная зона не определена до начала работы", [
            "Опасную зону определяют заранее.",
            "До начала перемещения груза.",
            "Людей в этой зоне быть не должно.",
            "Без этого работу начинать нельзя.",
        ], "(схема: опасная зона на площадке)"),
        ("S060-P02", "Ошибка: опасная обстановка оценивается слишком узко", [
            "Смотри не только на груз.",
            "Оценивай траекторию.",
            "Учитывай раскачивание и помехи.",
            "Думай о людях вокруг.",
        ], "(кейс-слайд: опасная обстановка)"),
        ("S060-P03", "Ошибка: работа продолжается без безопасных условий", [
            "Нет безопасных условий.",
            "Работу нужно прекратить.",
            "Нельзя работать без понятной организации.",
            "Риск сначала устраняют.",
        ], "(стоп-слайд: нет безопасных условий)"),
        ("S060-P04", "Ошибка: работа рядом с ЛЭП ведется на глаз", [
            "ЛЭП — особая опасность.",
            "Нельзя оценивать расстояние на глаз.",
            "Нужны специальные меры.",
            "Нужен организованный порядок работ.",
        ], "(схема: ЛЭП и зона риска)"),
        ("S060-P05", "Ошибка: в сложной операции участники не понимают схему и роли", [
            "Сложный подъем нельзя держать только на устных договоренностях.",
            "Участникам должны быть понятны схема, роли и последовательность действий.",
            "Будущий старший стропальщик удерживает подъем как единую операцию.",
            "Если порядок не разъяснен, работу не начинают.",
        ], "(схема: роли, маршрут и точки контроля)"),
    ]
    for code, title, bullets, placeholder in review_slides:
        slide = new_slide(prs, code, "Тема 2.3 • Разбор ошибки", title, "Короткий разбор неверного ответа с опорой на Игумнова.", footer_text)
        accent = WARN_RED if code in {"S060-P03", "S060-P04"} else DEEP_BLUE
        fill = SOFT_RED if code in {"S060-P03", "S060-P04"} else WHITE
        add_panel(slide, Inches(0.55), Inches(2.0), Inches(5.75), Inches(4.25), "Что нужно запомнить", bullets, accent=accent, fill=fill, level="Ошибка")
        add_placeholder_panel(slide, Inches(6.55), Inches(2.0), Inches(5.7), Inches(4.25), "Плейсхолдер под разбор", [
            placeholder,
            "(кнопка: вернуться к тесту)",
        ], accent=accent)

    slide = new_slide(prs, "S061", theme_label, "Финальный вывод по подтеме", "Главный итог блока по Игумнову: пространство работ нужно читать заранее, а сложные операции — заранее организовывать.", footer_text)
    add_panel(slide, Inches(0.55), Inches(2.0), Inches(5.7), Inches(4.25), "Что должен уметь слушатель", [
        "Видеть опасную зону до начала перемещения груза.",
        "Не работать в неясных и неорганизованных условиях.",
        "Останавливать работу при риске, а не после происшествия.",
        "Понимать, что рядом с ЛЭП действует особый порядок.",
        "Для 5-6 разряда — удерживать сложный подъем как управляемую операцию с понятными ролями и документами.",
    ], accent=DEEP_BLUE, level="Итог")
    add_placeholder_panel(slide, Inches(6.5), Inches(2.0), Inches(5.75), Inches(4.25), "Плейсхолдер под финальный визуал", [
        "(плакат-памятка: опасная зона / организация / ЛЭП / 5-6)",
        "(изображение: безопасная организация работ)",
    ])

    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT_PATH)
    print(OUT_PATH)


if __name__ == "__main__":
    build_deck()
