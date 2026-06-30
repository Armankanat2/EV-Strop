from pathlib import Path

from pptx import Presentation


PATH = Path(
    "presentations/ispring-course/module-01-stropovka-gruzov/live-preview/"
    "S029_theme22_preview_2026-06-29_v09.pptx"
)


def safe_text(shape):
    if hasattr(shape, "text") and shape.text:
        return shape.text.replace("\n", " | ")[:140]
    return ""


prs = Presentation(str(PATH))
print(f"FILE: {PATH.name}")
print(f"SLIDES: {len(prs.slides)}")
print("---")

for idx, slide in enumerate(prs.slides, start=1):
    print(f"SLIDE {idx}")
    for s_idx, shape in enumerate(slide.shapes, start=1):
        try:
            name = shape.name
        except Exception:
            name = "?"
        text = safe_text(shape)
        print(
            f"  {s_idx:02d} | {shape.shape_type} | "
            f"l={shape.left} t={shape.top} w={shape.width} h={shape.height} | "
            f"{name} | {text}"
        )
    print("---")
