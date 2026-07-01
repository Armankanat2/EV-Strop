from pathlib import Path
import shutil

from pptx import Presentation


SRC = Path(
    "presentations/ispring-course/module-01-stropovka-gruzov/live-preview/"
    "S042-S061_theme23_preview_2026-07-01_v07_editorial-tighten.pptx"
)
OUT = Path(
    "presentations/ispring-course/module-01-stropovka-gruzov/live-preview/"
    "S042-S061_theme23_preview_2026-07-01_v08_media-cleanup.pptx"
)


def set_text(shape, text: str) -> None:
    shape.text_frame.text = text


def find_slide(prs: Presentation, code: str):
    for slide in prs.slides:
        texts = [shape.text for shape in slide.shapes if hasattr(shape, "text") and shape.text]
        if code in texts:
            return slide
    raise ValueError(f"Slide {code} not found")


def main() -> None:
    OUT.parent.mkdir(parents=True, exist_ok=True)
    shutil.copy2(SRC, OUT)

    prs = Presentation(OUT)

    slide = find_slide(prs, "S042")
    set_text(slide.shapes[17], "Опорный визуал")
    set_text(slide.shapes[18], "Что должно быть видно")
    set_text(
        slide.shapes[19],
        "• Рабочая площадка, кран и траектория груза.\n"
        "• Контур опасной зоны.\n"
        "• Безопасная организация пространства без перегруза деталями.",
    )

    slide = find_slide(prs, "S043")
    set_text(slide.shapes[41], "Как читать маршрут")
    set_text(slide.shapes[42], "Шесть шагов подтемы")
    set_text(
        slide.shapes[43],
        "• Сначала опасная зона.\n"
        "• Затем запреты и специальные условия.\n"
        "• Далее работа у ЛЭП, усиление 5-6 и проверка понимания.",
    )

    slide = find_slide(prs, "S044")
    set_text(slide.shapes[17], "Опорный визуал")
    set_text(slide.shapes[18], "Что должно быть видно")
    set_text(
        slide.shapes[19],
        "• Площадка, кран и груз.\n"
        "• Границы опасной зоны.\n"
        "• Ограждения и знаки безопасности.",
    )
    set_text(slide.shapes[21], "Есть дополнительное пояснение")

    slide = find_slide(prs, "S045")
    set_text(slide.shapes[23], "Короткий ориентир")
    set_text(slide.shapes[24], "Что помогает увидеть риск")
    set_text(
        slide.shapes[25],
        "• Траектория груза.\n"
        "• Возможное раскачивание и разворот.\n"
        "• Помехи и люди в зоне.\n"
        "• Зазор не менее 500 мм над препятствием.",
    )

    slide = find_slide(prs, "S046")
    set_text(slide.shapes[23], "Ключевой вывод")
    set_text(slide.shapes[24], "Что нужно запомнить")
    set_text(
        slide.shapes[25],
        "• Работают только в безопасных условиях.\n"
        "• Простые нарушения чаще всего и приводят к тяжелым последствиям.\n"
        "• При опасности работу останавливают.",
    )

    slide = find_slide(prs, "S047")
    set_text(slide.shapes[25], "Опорный визуал")
    set_text(slide.shapes[26], "Как читать запреты")
    set_text(
        slide.shapes[27],
        "• Что запрещено.\n"
        "• Почему это опасно.\n"
        "• Когда работу нужно остановить.",
    )
    set_text(slide.shapes[29], "Есть дополнительное пояснение")

    slide = find_slide(prs, "S048")
    set_text(slide.shapes[31], "Для сборки")
    set_text(slide.shapes[32], "Лучше схема с озвучкой")
    set_text(
        slide.shapes[33],
        "• Показать проем, стесненное место и кран у канавы или котлована.\n"
        "• Здесь важнее схема и логика подачи груза, чем видео.",
    )
    set_text(slide.shapes[35], "Можно добавить интерактив: выбрать безопасную схему")

    slide = find_slide(prs, "S049")
    set_text(slide.shapes[23], "Опорная схема")
    set_text(slide.shapes[24], "Что должно быть понятно")
    set_text(
        slide.shapes[25],
        "• По какой схеме идет операция.\n"
        "• Где проходит маршрут и опасные границы.\n"
        "• Кто подает сигналы и кто отвечает за организацию работ.",
    )

    slide = find_slide(prs, "S050")
    set_text(slide.shapes[23], "Опорный визуал")
    set_text(slide.shapes[24], "Что должно быть видно")
    set_text(
        slide.shapes[25],
        "• Кран рядом с ЛЭП.\n"
        "• Зона риска и решение стоп.\n"
        "• Работу рядом с ЛЭП не оценивают на глаз.",
    )
    set_text(slide.shapes[27], "Есть дополнительное пояснение")
    set_text(slide.shapes[29], 'Далее: блок "5-6 разряд"')

    slide = find_slide(prs, "S051")
    set_text(slide.shapes[23], "Опорная схема")
    set_text(slide.shapes[24], "Как читать роли")
    set_text(
        slide.shapes[25],
        "• Ответственный за работы.\n"
        "• Старший стропальщик.\n"
        "• Крановщик.\n"
        "• Стропальщики.",
    )

    slide = find_slide(prs, "S052")
    set_text(slide.shapes[17], "Опорный чек-лист")
    set_text(slide.shapes[18], "Что собрать до начала")
    set_text(
        slide.shapes[19],
        "• Масса, схема, центр тяжести и маршрут.\n"
        "• Оснастка, подкладки, оттяжки и место приема груза.\n"
        "• Распределение ролей.",
    )

    slide = find_slide(prs, "S053")
    set_text(slide.shapes[23], "Опорный блок")
    set_text(slide.shapes[24], "Что должно быть оформлено")
    set_text(
        slide.shapes[25],
        "• ППР или техкарта.\n"
        "• Схема строповки и перемещения.\n"
        "• Роли участников.\n"
        "• При необходимости наряд-допуск и инструктаж.",
    )

    slide = find_slide(prs, "S054")
    set_text(slide.shapes[37], "Опорная схема")
    set_text(slide.shapes[38], "Как читать сложный подъем")
    set_text(
        slide.shapes[39],
        "• Разъясняют схему.\n"
        "• Готовят площадку.\n"
        "• Согласуют команды.\n"
        "• Ведут подъем без самовольных изменений.",
    )
    set_text(slide.shapes[41], "Можно добавить интерактив: расставить этапы по порядку")

    slide = find_slide(prs, "S055")
    set_text(slide.shapes[23], "Для сборки")
    set_text(slide.shapes[24], "Лучше короткое видео")
    set_text(
        slide.shapes[25],
        "• Контрольный подъем на малую высоту.\n"
        "• Отметить перекос, разворот, помеху и решение стоп.\n"
        "• Если видео нет, собрать псевдовидео из 4 кадров.",
    )

    slide = find_slide(prs, "S056")
    set_text(slide.shapes[23], "Для сборки")
    set_text(slide.shapes[24], "Лучше схема с выбором ситуации")
    set_text(
        slide.shapes[25],
        "• Показать точку установки, безопасную зону и порядок допуска.\n"
        "• Подходит интерактив: можно / нельзя по ситуациям у ЛЭП.",
    )
    set_text(slide.shapes[27], "Интерактив: можно / нельзя")

    slide = find_slide(prs, "S057")
    set_text(slide.shapes[23], "Опорное сравнение")
    set_text(slide.shapes[24], "Что отличает будущего старшего стропальщика")
    set_text(
        slide.shapes[25],
        "• Видит операцию целиком.\n"
        "• Заранее собирает роли, маршрут и контрольные точки.\n"
        "• Замечает организационный риск до ошибки в строповке.",
    )

    slide = find_slide(prs, "S058")
    set_text(slide.shapes[19], "Опорный визуал")
    set_text(slide.shapes[20], "Что должно быть видно")
    set_text(
        slide.shapes[21],
        "• Стоп-условия у ЛЭП.\n"
        "• Опасное сближение и запреты.\n"
        "• Учебный кейс без реальной опасной съемки.",
    )

    slide = find_slide(prs, "S059")
    set_text(slide.shapes[32], "Для разбора")
    set_text(slide.shapes[33], "Как собрать кейс")
    set_text(
        slide.shapes[34],
        "• Что произошло.\n"
        "• Какое правило нарушили.\n"
        "• Где работу нужно было остановить.",
    )

    slide = find_slide(prs, "S060")
    set_text(slide.shapes[17], "Для сборки")
    set_text(slide.shapes[18], "Тестовый экран")
    set_text(
        slide.shapes[19],
        "• Вопросы по опасной зоне, правилу 1 м, погодным стоп-условиям, ЛЭП и усилению 5-6.\n"
        "• Ошибки должны вести в разбор причин, а не только к правильному ответу.",
    )

    slide = find_slide(prs, "S061")
    set_text(slide.shapes[17], "Главный ориентир")
    set_text(slide.shapes[18], "Что нужно унести из блока")
    set_text(
        slide.shapes[19],
        "• Опасную зону читают заранее.\n"
        "• Не работают в неясных условиях.\n"
        "• У ЛЭП действует особый порядок.\n"
        "• Для 5-6 разряда сложный подъем заранее организуют.",
    )

    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    main()
