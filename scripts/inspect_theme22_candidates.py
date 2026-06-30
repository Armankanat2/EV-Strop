from pathlib import Path

from pptx import Presentation


FILES = [
    Path("presentations/ispring-course/module-01-stropovka-gruzov/live-preview/S001-S040_live_preview_working_2026-06-29_v53.pptx"),
    Path("presentations/ispring-course/module-01-stropovka-gruzov/live-preview/S001-S040_live_preview_working_2026-06-29_v54.pptx"),
    Path("presentations/ispring-course/module-01-stropovka-gruzov/live-preview/S001-S040_live_preview_working_2026-06-29_v55.pptx"),
    Path("presentations/ispring-course/module-01-stropovka-gruzov/live-preview/S029_theme22_preview_2026-06-29_v09.pptx"),
]

TARGETS = [
    "S031-P01",
    "S034-P01",
    "S035-P01",
    "S035-PP01",
    "S037-P01",
    "S037-PP01",
    "S039-P01",
    "S039-PP01",
]


def collect_slide_text(slide) -> str:
    parts = []
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text:
            parts.append(shape.text.replace("\n", " | "))
    return " || ".join(parts)


for path in FILES:
    prs = Presentation(str(path))
    print(f"FILE: {path.name}")
    print(f"SLIDES: {len(prs.slides)}")
    found = []
    for idx, slide in enumerate(prs.slides, start=1):
        text = collect_slide_text(slide)
        hits = [target for target in TARGETS if target in text]
        if hits:
            found.extend(hit for hit in hits if hit not in found)
            print(f"  slide {idx}: {', '.join(hits)}")
            print(f"    {text[:280]}")
    print("FOUND:", ", ".join(found) if found else "none")
    print("---")
