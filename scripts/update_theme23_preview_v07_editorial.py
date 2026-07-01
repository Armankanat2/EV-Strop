from pathlib import Path
import shutil

from pptx import Presentation


SRC = Path(
    "presentations/ispring-course/module-01-stropovka-gruzov/live-preview/"
    "S042-S061_theme23_preview_2026-07-01_v06_ev-text_5-6-linear_media-plan_rules.pptx"
)
OUT = Path(
    "presentations/ispring-course/module-01-stropovka-gruzov/live-preview/"
    "S042-S061_theme23_preview_2026-07-01_v07_editorial-tighten.pptx"
)


def set_text(shape, text: str) -> None:
    shape.text_frame.text = text


def find_slide(prs: Presentation, code: str):
    for slide in prs.slides:
        texts = [shape.text for shape in slide.shapes if hasattr(shape, "text") and shape.text]
        if code in texts:
            return slide
    raise ValueError(f"Slide {code} not found")


def main() -> None:
    OUT.parent.mkdir(parents=True, exist_ok=True)
    shutil.copy2(SRC, OUT)

    prs = Presentation(OUT)

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text == "ЭВ Стропальщик • Тема 2.3 • preview":
                set_text(shape, "ЭВ Стропальщик • Тема 2.3")
            if hasattr(shape, "text") and shape.text == "ЭВ Стропальщик • Тема 2.3 • Усиление 5-6":
                set_text(shape, "ЭВ Стропальщик • Тема 2.3 • 5-6 разряд")

    slide = find_slide(prs, "S042")
    set_text(slide.shapes[5], "Опасная зона, запреты, специальные условия, ЛЭП и усиление 5-6.")

    slide = find_slide(prs, "S043")
    set_text(
        slide.shapes[5],
        "Маршрут блока: от опасной зоны и запретов к ЛЭП, сложным подъемам и разбору нарушений.",
    )
    set_text(slide.shapes[32], "Усиление 5-6: руководство операцией, документы, сложный подъем и работа у ЛЭП.")

    slide = find_slide(prs, "S046")
    set_text(
        slide.shapes[5],
        "Базовые правила и запреты образуют один обязательный рабочий блок.",
    )

    slide = find_slide(prs, "S047")
    set_text(
        slide.shapes[5],
        "Без понятной схемы, без безопасных условий и с опасными приемами работу не ведут.",
    )

    slide = find_slide(prs, "S049")
    set_text(
        slide.shapes[19],
        "• Не действуй самовольно.\n"
        "• Если стропальщиков несколько, команды крановщику подает один назначенный участник.\n"
        "• Не заменяй схему опытом.\n"
        "• Не начинай работу при неясном порядке и неясной роли участников.",
    )

    slide = find_slide(prs, "S055")
    set_text(
        slide.shapes[19],
        "• Ошибку нужно выявить на контрольном подъеме, а не во время перемещения.\n"
        "• Именно здесь становится видно, верна ли схема в реальных условиях.\n"
        "• Контрольный подъем позволяет увидеть риск до основного перемещения.",
    )

    slide = find_slide(prs, "S056")
    set_text(
        slide.shapes[5],
        "У ЛЭП действует отдельный организационный режим работ.",
    )

    slide = find_slide(prs, "S057")
    set_text(
        slide.shapes[5],
        "Разница не в смелости, а в умении держать подъем как управляемую операцию.",
    )

    slide = find_slide(prs, "S059")
    set_text(
        slide.shapes[5],
        "Разбор строится по трем вопросам: что произошло, что нарушили, как нужно было действовать.",
    )

    slide = find_slide(prs, "S061")
    set_text(
        slide.shapes[5],
        "Главный итог блока: пространство работ читают заранее, а сложные операции заранее организуют.",
    )

    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    main()
