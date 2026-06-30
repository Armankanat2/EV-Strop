from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


SRC_PATH = Path(
    "presentations/ispring-course/module-01-stropovka-gruzov/live-preview/"
    "S029_theme22_preview_2026-06-29_v09.pptx"
)
OUT_PATH = Path(
    "presentations/ispring-course/module-01-stropovka-gruzov/live-preview/"
    "S029_theme22_preview_2026-06-30_v13_ev-group-techstyle-premium-subslides.pptx"
)

PNG_OUT_DIR = Path(
    "presentations/ispring-course/module-01-stropovka-gruzov/live-preview/"
    "S029_theme22_preview_2026-06-30_v13_ev-group-techstyle-premium-subslides_png"
)


BLUE = RGBColor(0x00, 0x57, 0xFF)
DEEP_BLUE = RGBColor(0x00, 0x38, 0xB8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAPHITE = RGBColor(0x11, 0x11, 0x11)
DARK_GRAY = RGBColor(0x4B, 0x55, 0x63)
LIGHT_GRAY = RGBColor(0xE5, 0xE7, 0xEB)
PALE_BLUE = RGBColor(0xEA, 0xF0, 0xFF)
RED = RGBColor(0xC5, 0x1F, 0x1F)
SOFT_RED = RGBColor(0xFE, 0xF2, 0xF2)


ASSETS = {
    "S029": Path("assets/course-media/module-01-stropovka-gruzov/images/S004_avatar-1_signal-podnyat-gruz_white-bg.png"),
    "S031": Path("assets/course-media/module-02-tekhnologiya-raboty-stropalshchika/diagrams/error-analysis/S031-P01_podgotovka-k-rabote_checklist.png"),
    "S032": Path("assets/course-media/module-02-tekhnologiya-raboty-stropalshchika/diagrams/error-analysis/S031-P02_plohaia-osveshchennost_warning.png"),
    "S035": Path("assets/course-media/module-01-stropovka-gruzov/images/S004_avatar-1_signal-podnyat-gruz_white-bg.png"),
    "S036": Path("assets/course-media/module-02-tekhnologiya-raboty-stropalshchika/diagrams/error-analysis/S031-P05_stop-prioritet_poster.png"),
    "S037": Path("assets/course-media/module-02-tekhnologiya-raboty-stropalshchika/diagrams/error-analysis/S031-P08_opasnaia-zona_movement-diagram.png"),
    "S038": Path("assets/course-media/module-02-tekhnologiya-raboty-stropalshchika/diagrams/error-analysis/S031-P11_raskachivanie-liudi_case-diagram.png"),
    "S039": Path("assets/course-media/module-02-tekhnologiya-raboty-stropalshchika/diagrams/error-analysis/S031-P09_skladovanie_podkladki-diagram.png"),
    "S031-P01": Path("assets/course-media/module-02-tekhnologiya-raboty-stropalshchika/diagrams/error-analysis/S031-P01_podgotovka-k-rabote_checklist.png"),
    "S034-P01": Path("assets/course-media/module-02-tekhnologiya-raboty-stropalshchika/diagrams/error-analysis/S031-P03_obviazka-pered-podem_steps.png"),
    "S035-P01": Path("assets/course-media/module-01-stropovka-gruzov/images/S004_avatar-1_signal-palm-down-source.png"),
    "S035-PP01": Path("assets/course-media/module-01-stropovka-gruzov/images/S004_avatar-1_signal-podnyat-gruz_white-bg.png"),
    "S037-P01": Path("assets/course-media/module-02-tekhnologiya-raboty-stropalshchika/diagrams/error-analysis/S031-P08_opasnaia-zona_movement-diagram.png"),
    "S037-PP01": Path("assets/course-media/module-02-tekhnologiya-raboty-stropalshchika/diagrams/error-analysis/S031-P07_ottiagka_safe-guiding-diagram.png"),
    "S039-P01": Path("assets/course-media/module-02-tekhnologiya-raboty-stropalshchika/diagrams/error-analysis/S031-P10_rasstropovka_rano-pravilno.png"),
    "S039-PP01": Path("assets/course-media/module-02-tekhnologiya-raboty-stropalshchika/diagrams/error-analysis/S031-P09_skladovanie_podkladki-diagram.png"),
}


CAPTIONS = {
    "S029": "Стропальщик подает команду перед началом грузоподъемной операции",
    "S031": "Памятка по подготовке до начала работ",
    "S032": "Пример условия, при котором работу начинать нельзя",
    "S035": "Ключевой ручной сигнал для взаимодействия с крановщиком",
    "S036": "Команда «Стоп» имеет безусловный приоритет",
    "S037": "Опасная зона и безопасная позиция при перемещении груза",
    "S038": "Нештатная ситуация: перемещение нужно остановить",
    "S039": "Подкладки и устойчивое положение перед расстроповкой",
    "S031-P01": "Подготовка: задание, условия, оснастка и рабочая зона",
    "S034-P01": "Раскрытие маршрута операции по этапам",
    "S035-P01": "Способы подачи сигналов и пример ручной команды",
    "S035-PP01": "Галерея жестов должна читаться как единая система",
    "S037-P01": "Схема опасной зоны для отдельного разбора",
    "S037-PP01": "Работа с оттяжкой только на безопасной дистанции",
    "S039-P01": "Расстроповка допустима только после устойчивой установки",
    "S039-PP01": "Варианты безопасного складирования",
}


STEP_LABELS = {
    "S029": "Шаг 00",
    "S030": "Шаг 00",
    "S031": "Шаг 01",
    "S031-P01": "Шаг 01",
    "S032": "Шаг 01",
    "S033": "Шаг 02",
    "S034": "Шаг 02",
    "S034-P01": "Шаг 02",
    "S035": "Шаг 03",
    "S035-P01": "Шаг 03",
    "S035-PP01": "Шаг 03",
    "S036": "Шаг 03",
    "S037": "Шаг 04",
    "S037-P01": "Шаг 04",
    "S037-PP01": "Шаг 04",
    "S038": "Шаг 04",
    "S039": "Шаг 05",
    "S039-P01": "Шаг 05",
    "S039-PP01": "Шаг 05",
    "S040": "Шаг 05",
    "S041": "Шаг 05",
}


LEVEL_LABELS = {
    "S029": "Опора",
    "S031": "Контроль",
    "S031-P01": "Опора",
    "S032": "Критично",
    "S033": "Роль",
    "S034": "Маршрут",
    "S034-P01": "Опора",
    "S035": "Сигнал",
    "S035-P01": "Сигнал",
    "S035-PP01": "Сигнал",
    "S036": "Критично",
    "S037": "Контроль",
    "S037-P01": "Критично",
    "S037-PP01": "Опора",
    "S038": "Критично",
    "S039": "Контроль",
    "S039-P01": "Критично",
    "S039-PP01": "Опора",
    "S040": "Проверка",
    "S041": "Итог",
}


VISUAL_TITLES = {
    "S029": "Ключевой визуальный вход",
    "S031": "Памятка и схема",
    "S031-P01": "Опорная памятка",
    "S032": "Разбор опасного условия",
    "S035": "Базовый жест",
    "S035-P01": "Сигнал рукой",
    "S035-PP01": "Система жестов",
    "S036": "Приоритет остановки",
    "S037": "Опасная зона",
    "S037-P01": "Контур опасной зоны",
    "S037-PP01": "Сопровождение груза",
    "S038": "Нештатная ситуация",
    "S039": "Устойчивое основание",
    "S039-P01": "Финальная проверка",
    "S039-PP01": "Безопасное складирование",
}


PREMIUM_VISUAL_CODES = {
    "S035",
    "S035-P01",
    "S035-PP01",
    "S036",
    "S037",
    "S037-P01",
    "S037-PP01",
    "S038",
    "S039",
    "S039-P01",
    "S039-PP01",
}


def set_fill(shape, rgb, transparency=None):
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb
    if transparency is not None:
        shape.fill.transparency = transparency


def set_line(shape, rgb, width_pt=1.0):
    shape.line.color.rgb = rgb
    shape.line.width = Pt(width_pt)


def style_text_shape(shape, size, color, bold=False, align=PP_ALIGN.LEFT, wrap=True):
    tf = shape.text_frame
    tf.word_wrap = wrap
    tf.margin_left = Pt(8)
    tf.margin_right = Pt(8)
    tf.margin_top = Pt(6)
    tf.margin_bottom = Pt(6)
    for p in tf.paragraphs:
        p.alignment = align
        for run in p.runs:
            run.font.name = "Arial"
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = color


def contain_image(slide, image_path, left, top, width, height):
    with Image.open(image_path) as img:
        img_w, img_h = img.size

    box_ratio = width / height
    img_ratio = img_w / img_h
    if img_ratio > box_ratio:
        pic_w = width
        pic_h = width / img_ratio
    else:
        pic_h = height
        pic_w = height * img_ratio

    pic_left = left + (width - pic_w) / 2
    pic_top = top + (height - pic_h) / 2
    return slide.shapes.add_picture(str(image_path), pic_left, pic_top, width=pic_w, height=pic_h)


def add_small_accent(slide):
    accent = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0.56),
        Inches(0.74),
        Inches(1.12),
        Inches(0.05),
    )
    set_fill(accent, BLUE)
    accent.line.fill.background()


def add_strip(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, width, height)
    set_fill(shape, color)
    shape.line.fill.background()
    return shape


def add_chip(slide, left, top, width, height, text, fill, color):
    chip = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, width, height)
    set_fill(chip, fill)
    chip.line.fill.background()
    tf = chip.text_frame
    tf.clear()
    tf.word_wrap = False
    tf.margin_left = Pt(4)
    tf.margin_right = Pt(4)
    tf.margin_top = Pt(1)
    tf.margin_bottom = Pt(1)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name = "Arial"
    run.font.size = Pt(8.5)
    run.font.bold = True
    run.font.color.rgb = color
    return chip


def add_panel_meta(slide, left, top, width, code, title):
    step = STEP_LABELS.get(code, "Шаг")
    level = LEVEL_LABELS.get(code, "Опора")
    level_fill = PALE_BLUE if level not in {"Критично"} else SOFT_RED
    level_color = DEEP_BLUE if level not in {"Критично"} else RED

    add_strip(slide, left, top, width, Inches(0.04), DEEP_BLUE)

    chip_top = top + Inches(0.10)
    add_chip(slide, left + Inches(0.12), chip_top, Inches(0.95 if len(code) <= 7 else 1.12), Inches(0.22), code, PALE_BLUE, DEEP_BLUE)
    add_chip(slide, left + Inches(1.18 if len(code) <= 7 else 1.35), chip_top, Inches(0.85), Inches(0.22), step, WHITE, DARK_GRAY)
    add_chip(slide, left + width - Inches(1.12), chip_top, Inches(1.0), Inches(0.22), level, level_fill, level_color)

    label = slide.shapes.add_textbox(left + Inches(0.12), top + Inches(0.38), width - Inches(0.24), Inches(0.24))
    tf = label.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.name = "Arial"
    run.font.size = Pt(13)
    run.font.bold = True
    run.font.color.rgb = GRAPHITE
    return label


def restyle_slide(slide):
    code = slide.shapes[2].text.strip()
    left_panel = slide.shapes[4]
    left_header = slide.shapes[5]
    left_body = slide.shapes[6]
    right_panel = slide.shapes[7]
    right_body = slide.shapes[8]
    back_btn = slide.shapes[9]
    next_btn = slide.shapes[10]
    has_image = code in ASSETS and ASSETS[code].exists()
    premium_visual = code in PREMIUM_VISUAL_CODES

    # Background band
    set_fill(slide.shapes[0], WHITE)
    slide.shapes[0].line.fill.background()
    add_small_accent(slide)

    # Titles
    slide.shapes[1].text = "Тема 2.2\nТехнология работы"
    style_text_shape(slide.shapes[1], 10, GRAPHITE, bold=True, wrap=False)
    style_text_shape(slide.shapes[3], 27 if len(slide.shapes[3].text) < 34 else 24, GRAPHITE, bold=True)
    slide.shapes[1].top = Inches(0.30)
    slide.shapes[1].left = Inches(0.64)
    slide.shapes[1].width = Inches(1.45)
    slide.shapes[1].height = Inches(0.42)
    slide.shapes[3].top = Inches(0.93)
    slide.shapes[3].left = Inches(0.64)
    slide.shapes[3].width = Inches(6.2)
    slide.shapes[3].height = Inches(0.8)

    # Code badge
    set_fill(slide.shapes[2], DEEP_BLUE)
    slide.shapes[2].line.fill.background()
    badge_width = 1.15
    badge_font = 13
    if len(code) >= 8:
        badge_width = 1.32
        badge_font = 12
    if len(code) >= 9:
        badge_width = 1.48
        badge_font = 11
    style_text_shape(slide.shapes[2], badge_font, WHITE, bold=True, align=PP_ALIGN.CENTER, wrap=False)
    slide.shapes[2].left = Inches(12.20) - Inches(badge_width)
    slide.shapes[2].top = Inches(0.34)
    slide.shapes[2].width = Inches(badge_width)
    slide.shapes[2].height = Inches(0.34)

    # Panels
    for panel in [left_panel, right_panel]:
        set_fill(panel, WHITE)
        set_line(panel, LIGHT_GRAY, width_pt=1.25)

    set_fill(left_header, WHITE)
    left_header.line.fill.background()
    style_text_shape(left_header, 13, GRAPHITE, bold=True)

    # Body text
    left_size = 16
    right_size = 14
    if "-P01" in code or "-PP01" in code:
        left_size = 13
        right_size = 13
    if code in {"S031-P01", "S034-P01", "S035-P01", "S035-PP01", "S039-P01", "S039-PP01"}:
        left_size = 12
        right_size = 12

    style_text_shape(left_body, left_size, GRAPHITE)
    style_text_shape(right_body, right_size, DARK_GRAY)

    # Strict geometry with larger visual containers
    left_panel.top = Inches(2.0)
    left_panel.left = Inches(0.74)
    left_panel.width = Inches(4.45 if has_image else 5.1)
    left_panel.height = Inches(4.28)
    left_header.top = Inches(2.34 if premium_visual else 2.26)
    left_header.left = Inches(0.74)
    left_header.width = left_panel.width
    left_header.height = Inches(0.34 if premium_visual else 0.40)
    left_body.top = Inches(2.92 if premium_visual else 2.84)
    left_body.left = Inches(0.98)
    left_body.width = left_panel.width - Inches(0.42)
    left_body.height = Inches(2.96 if premium_visual else 3.02)

    right_panel.top = Inches(2.0)
    right_panel.left = Inches(5.18 if premium_visual else (5.45 if has_image else 6.08))
    right_panel.width = Inches(7.03 if premium_visual else (6.76 if has_image else 5.96))
    right_panel.height = Inches(4.28)

    add_strip(slide, left_panel.left, left_panel.top, left_panel.width, Inches(0.04), DEEP_BLUE)
    add_chip(slide, left_panel.left + Inches(0.12), left_panel.top + Inches(0.10), Inches(0.85), Inches(0.22), STEP_LABELS.get(code, "Шаг"), PALE_BLUE, DEEP_BLUE)
    add_chip(slide, left_panel.left + left_panel.width - Inches(1.12), left_panel.top + Inches(0.10), Inches(1.0), Inches(0.22), LEVEL_LABELS.get(code, "Опора"), SOFT_RED if LEVEL_LABELS.get(code) == "Критично" else WHITE, RED if LEVEL_LABELS.get(code) == "Критично" else DARK_GRAY)

    if premium_visual:
        add_strip(slide, left_panel.left + Inches(0.12), left_panel.top + Inches(0.58), left_panel.width - Inches(0.24), Inches(0.018), LIGHT_GRAY)
        style_text_shape(left_body, left_size, GRAPHITE)
        left_body.text_frame.margin_left = Pt(10)
        left_body.text_frame.margin_right = Pt(8)

    # Right panel: text-only vs image-driven
    image_path = ASSETS.get(code)
    if image_path and image_path.exists():
        if premium_visual:
            set_fill(right_panel, WHITE)
            add_strip(slide, right_panel.left + Inches(0.16), right_panel.top + Inches(0.74), right_panel.width - Inches(0.32), Inches(0.012), LIGHT_GRAY)
        right_body.top = right_panel.top + Inches(3.90 if premium_visual else 3.88)
        right_body.left = right_panel.left + Inches(0.20)
        right_body.width = right_panel.width - Inches(0.40)
        right_body.height = Inches(0.24)
        style_text_shape(right_body, 9 if premium_visual else 10, DARK_GRAY)
        right_body.text = CAPTIONS.get(code, "")
        add_panel_meta(slide, right_panel.left, right_panel.top, right_panel.width, code, VISUAL_TITLES.get(code, "Технический визуал"))
        contain_image(slide, image_path, right_panel.left + Inches(0.18), right_panel.top + Inches(0.88), right_panel.width - Inches(0.36), Inches(2.98 if premium_visual else 2.92))
        add_strip(slide, right_panel.left + Inches(0.18), right_panel.top + Inches(3.80), right_panel.width - Inches(0.36), Inches(0.02), LIGHT_GRAY)
    else:
        add_panel_meta(slide, right_panel.left, right_panel.top, right_panel.width, code, "Технический комментарий")
        right_body.top = right_panel.top + Inches(0.94)
        right_body.left = right_panel.left + Inches(0.20)
        right_body.width = right_panel.width - Inches(0.40)
        right_body.height = Inches(3.05)

    # Navigation
    set_fill(back_btn, WHITE)
    set_line(back_btn, LIGHT_GRAY, width_pt=1.0)
    style_text_shape(back_btn, 12, DARK_GRAY, bold=True, align=PP_ALIGN.CENTER)
    back_btn.left = Inches(0.74)
    back_btn.top = Inches(6.55)
    back_btn.width = Inches(1.85)
    back_btn.height = Inches(0.38)

    set_fill(next_btn, BLUE)
    next_btn.line.fill.background()
    style_text_shape(next_btn, 12, WHITE, bold=True, align=PP_ALIGN.CENTER)
    next_btn.left = Inches(10.15)
    next_btn.top = Inches(6.55)
    next_btn.width = Inches(2.08)
    next_btn.height = Inches(0.38)


def main():
    prs = Presentation(str(SRC_PATH))
    for slide in prs.slides:
        restyle_slide(slide)
    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT_PATH)
    print(OUT_PATH)


if __name__ == "__main__":
    main()
