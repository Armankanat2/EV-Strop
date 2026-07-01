from pathlib import Path
import shutil

from pptx import Presentation
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.util import Inches


SRC = Path(
    "presentations/ispring-course/module-01-stropovka-gruzov/live-preview/"
    "S042-S061_theme23_preview_2026-07-01_v09_overflow-fix.pptx"
)
OUT = Path(
    "presentations/ispring-course/module-01-stropovka-gruzov/live-preview/"
    "S042-S061_theme23_preview_2026-07-01_v10_targeted-layout-fix.pptx"
)


def set_text(shape, text: str) -> None:
    shape.text_frame.text = text


def tighten(shape) -> None:
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.margin_left = Inches(0.04)
    tf.margin_right = Inches(0.04)
    tf.margin_top = Inches(0.03)
    tf.margin_bottom = Inches(0.03)


def find_slide(prs: Presentation, code: str):
    for slide in prs.slides:
        texts = [shape.text for shape in slide.shapes if hasattr(shape, "text") and shape.text]
        if code in texts:
            return slide
    raise ValueError(f"Slide {code} not found")


def main() -> None:
    OUT.parent.mkdir(parents=True, exist_ok=True)
    shutil.copy2(SRC, OUT)

    prs = Presentation(OUT)

    slide = find_slide(prs, "S043")
    set_text(
        slide.shapes[43],
        "• Опасная зона.\n"
        "• Запреты и специальные условия.\n"
        "• ЛЭП, блок 5-6 и проверка понимания.",
    )

    slide = find_slide(prs, "S045")
    set_text(
        slide.shapes[13],
        "• Смотри на траекторию груза.\n"
        "• Учитывай раскачивание и разворот.\n"
        "• Над препятствием нужен зазор не менее 500 мм.\n"
        "• Проверяй помехи, проемы и людей в зоне.",
    )
    set_text(
        slide.shapes[25],
        "• Траектория.\n"
        "• Раскачивание.\n"
        "• Помехи и люди.\n"
        "• Зазор 500 мм.",
    )

    slide = find_slide(prs, "S046")
    set_text(
        slide.shapes[13],
        "• Работай только в безопасных условиях.\n"
        "• Соблюдай порядок работ и указания ответственных лиц.\n"
        "• При плохой видимости и сложной связи работай через сигнальщика.\n"
        "• При сильном ветре, грозе, тумане, снегопаде и резком ухудшении видимости работу прекращают.\n"
        "• При неисправности крана или оснастки работу прекращают.",
    )

    slide = find_slide(prs, "S047")
    set_text(
        slide.shapes[13],
        "• Не подтаскивай груз краном и не делай рывков.\n"
        "• Не исправляй строповку под висящим грузом.\n"
        "• Не работай без понятной схемы и организации работ.\n"
        "• Не опускай груз, если люди находятся в кузове или в кабине.\n"
        "• Не перемещай стрелу над кабиной автомобиля.\n"
        "• Не работай без всех выносных опор.",
    )
    set_text(
        slide.shapes[21],
        "• Нет безопасных условий.\n"
        "• Не ясна организация работ.\n"
        "• Есть риск для людей и оборудования.\n"
        "• Погода, видимость или состояние крана не позволяют работать.",
    )

    slide = find_slide(prs, "S048")
    set_text(slide.shapes[12], "Груз подают только по понятной траектории.")

    slide = find_slide(prs, "S054")
    set_text(slide.shapes[12], "Все понимают схему, маршрут и ограничения.")
    set_text(slide.shapes[33], "• Заранее видеть перекос, разворот, помеху и несогласованные действия.\n• Собирать подъем как управляемую последовательность.")
    set_text(
        slide.shapes[39],
        "• Разъясняют схему.\n"
        "• Готовят площадку.\n"
        "• Согласуют команды.\n"
        "• Ведут подъем без самовольных изменений.",
    )
    set_text(slide.shapes[41], "Интерактив: расставить этапы по порядку")

    slide = find_slide(prs, "S057")
    set_text(
        slide.shapes[19],
        "• Видит в основном свой участок.\n"
        "• Хорошо исполняет указания.\n"
        "• Следит прежде всего за своей строповкой.\n"
        "• Реже замечает организационный риск заранее.",
    )
    set_text(
        slide.shapes[25],
        "• Видит операцию целиком.\n"
        "• Заранее собирает роли, маршрут и контроль.\n"
        "• Замечает риск до ошибки в строповке.",
    )

    slide = find_slide(prs, "S061")
    set_text(
        slide.shapes[19],
        "• Опасную зону читают заранее.\n"
        "• Не работают в неясных условиях.\n"
        "• У ЛЭП действует особый порядок.\n"
        "• Для 5-6 разряда сложный подъем заранее организуют.",
    )

    for code, indices in {
        "S043": [43],
        "S045": [13, 25],
        "S046": [13],
        "S047": [13, 21],
        "S048": [12],
        "S054": [12, 33, 39, 41],
        "S057": [19, 25],
        "S061": [19],
    }.items():
        slide = find_slide(prs, code)
        for idx in indices:
            tighten(slide.shapes[idx])

    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    main()
