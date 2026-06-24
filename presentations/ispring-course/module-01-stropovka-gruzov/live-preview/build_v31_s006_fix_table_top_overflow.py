from pathlib import Path

from pptx import Presentation
from pptx.util import Inches


BASE_DIR = Path(__file__).resolve().parent
SRC = BASE_DIR / "S001-S007_live_preview_working_2026-06-24_v30.pptx"
OUT = BASE_DIR / "S001-S007_live_preview_working_2026-06-24_v31.pptx"

S006_INDEX = 16
TARGET_TOP = Inches(2.17)


def main() -> None:
    prs = Presentation(SRC)
    slide = prs.slides[S006_INDEX]

    table_shape = None
    for shape in slide.shapes:
        if shape.has_table:
            table_shape = shape
            break

    if table_shape is None:
        raise RuntimeError("Could not find S006 table")

    table_shape.top = TARGET_TOP

    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    main()
