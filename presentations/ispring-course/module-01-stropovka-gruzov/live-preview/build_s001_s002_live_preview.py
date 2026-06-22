from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


OUT_DIR = Path(__file__).resolve().parent
OUT_FILE = OUT_DIR / "S001-S007_live_preview_working_2026-06-22_v4.pptx"


def u(text):
    return text.encode("ascii").decode("unicode_escape")


NAVY = RGBColor(18, 43, 66)
BLUE = RGBColor(38, 90, 137)
GREEN = RGBColor(66, 133, 95)
STEEL = RGBColor(83, 119, 149)
PALE = RGBColor(240, 244, 247)
ORANGE = RGBColor(222, 116, 54)
WHITE = RGBColor(255, 255, 255)
TEXT = RGBColor(33, 43, 54)
GRAY = RGBColor(110, 122, 134)
LINE = RGBColor(210, 218, 226)


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def add_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, x, y, w, h, text, size=18, bold=False, color=TEXT, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return box


def add_band(slide, title, code, subtitle=""):
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.72))
    band.fill.solid()
    band.fill.fore_color.rgb = NAVY
    band.line.fill.background()

    add_text(slide, Inches(0.45), Inches(0.14), Inches(8.6), Inches(0.3), title, 24, True, WHITE)

    tag = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(11.45), Inches(0.14), Inches(1.35), Inches(0.35)
    )
    tag.fill.solid()
    tag.fill.fore_color.rgb = ORANGE
    tag.line.fill.background()
    tf = tag.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = code
    r.font.size = Pt(16)
    r.font.bold = True
    r.font.color.rgb = WHITE

    if subtitle:
        add_text(slide, Inches(0.5), Inches(0.82), Inches(6.5), Inches(0.25), subtitle, 11, False, GRAY)


def add_button(slide, x, y, w, h, text, fill, font_size=18):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = fill
    tf = shape.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.size = Pt(font_size)
    r.font.bold = True
    r.font.color.rgb = WHITE
    return shape


def add_linear_nav(slide, prev_text=None, next_text=None):
    prev_btn = None
    next_btn = None
    if prev_text:
        prev_btn = add_button(slide, Inches(0.7), Inches(6.55), Inches(2.35), Inches(0.5), prev_text, STEEL)
    if next_text:
        next_btn = add_button(slide, Inches(10.25), Inches(6.55), Inches(2.35), Inches(0.5), next_text, ORANGE)
    return prev_btn, next_btn


def add_panel(slide, x, y, w, h, title, lines, accent):
    panel = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
    panel.fill.solid()
    panel.fill.fore_color.rgb = WHITE
    panel.line.color.rgb = accent
    panel.line.width = Pt(1.5)

    head = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, Inches(0.72))
    head.fill.solid()
    head.fill.fore_color.rgb = accent
    head.line.color.rgb = accent
    tf = head.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = title
    r.font.size = Pt(22)
    r.font.bold = True
    r.font.color.rgb = WHITE

    body = slide.shapes.add_textbox(x + Inches(0.28), y + Inches(1.0), w - Inches(0.56), h - Inches(1.2))
    tf = body.text_frame
    tf.clear()
    tf.word_wrap = True
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = line
        p.level = 0
        p.bullet = True
        p.space_after = Pt(8)
        if p.runs:
            p.runs[0].font.size = Pt(17)
            p.runs[0].font.color.rgb = TEXT
    return panel


def add_header_cell(slide, x, y, w, h, text, fill, font_size=15):
    cell = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, y, w, h)
    cell.fill.solid()
    cell.fill.fore_color.rgb = fill
    cell.line.color.rgb = fill
    tf = cell.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.size = Pt(font_size)
    r.font.bold = True
    r.font.color.rgb = WHITE
    return cell


def add_body_cell(slide, x, y, w, h, text, align=PP_ALIGN.LEFT, fill=WHITE, size=14, bold=False):
    cell = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, y, w, h)
    cell.fill.solid()
    cell.fill.fore_color.rgb = fill
    cell.line.color.rgb = LINE
    tf = cell.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = TEXT
    return cell


def add_info_panel(slide, x, y, w, h, title, lines, accent=GREEN):
    panel = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
    panel.fill.solid()
    panel.fill.fore_color.rgb = WHITE
    panel.line.color.rgb = accent
    panel.line.width = Pt(1.5)

    head = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, Inches(0.62))
    head.fill.solid()
    head.fill.fore_color.rgb = accent
    head.line.color.rgb = accent
    tf = head.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = title
    r.font.size = Pt(19)
    r.font.bold = True
    r.font.color.rgb = WHITE

    body = slide.shapes.add_textbox(x + Inches(0.25), y + Inches(0.86), w - Inches(0.5), h - Inches(1.06))
    tf = body.text_frame
    tf.clear()
    tf.word_wrap = True
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = line
        p.bullet = False
        p.space_after = Pt(8)
        if p.runs:
            p.runs[0].font.size = Pt(15)
            p.runs[0].font.color.rgb = TEXT
    return panel, head, body


def add_glossary_column(slide, x, y, w, h, items, term_size=12, body_size=11):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for idx, (term, definition) in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(3)
        term_run = p.add_run()
        term_run.text = term
        term_run.font.size = Pt(term_size)
        term_run.font.bold = True
        term_run.font.color.rgb = BLUE
        def_run = p.add_run()
        def_run.text = u(r" - ") + definition
        def_run.font.size = Pt(body_size)
        def_run.font.color.rgb = TEXT
    return box
    return panel


def add_detail_table(slide, title, code, doc_name, about, importance, accent):
    add_bg(slide, WHITE)
    add_band(slide, title, code, u(r"\u041f\u043e\u0434\u0432\u0430\u043b \u043e\u0442 S003"))

    add_text(
        slide,
        Inches(0.7),
        Inches(1.1),
        Inches(5.6),
        Inches(0.42),
        u(r"\u041a\u043e\u0440\u043e\u0442\u043a\u0430\u044f \u0441\u043f\u0440\u0430\u0432\u043a\u0430 \u043f\u043e \u0434\u043e\u043a\u0443\u043c\u0435\u043d\u0442\u0443"),
        22,
        True,
        NAVY,
    )

    table_x = Inches(0.7)
    table_y = Inches(1.75)
    left_w = Inches(2.6)
    right_w = Inches(7.35)

    add_header_cell(slide, table_x, table_y, left_w, Inches(0.55), u(r"\u041f\u043e\u043b\u0435"), accent)
    add_header_cell(slide, table_x + left_w, table_y, right_w, Inches(0.55), u(r"\u0421\u043e\u0434\u0435\u0440\u0436\u0430\u043d\u0438\u0435"), accent)

    labels = [
        u(r"\u0414\u043e\u043a\u0443\u043c\u0435\u043d\u0442"),
        u(r"\u041e \u0447\u0435\u043c \u0434\u043e\u043a\u0443\u043c\u0435\u043d\u0442"),
        u(r"\u0427\u0435\u043c \u0432\u0430\u0436\u043d\u043e \u0434\u043b\u044f \u0441\u0442\u0440\u043e\u043f\u0430\u043b\u044c\u0449\u0438\u043a\u0430"),
    ]
    values = [doc_name, about, importance]
    heights = [Inches(0.8), Inches(1.0), Inches(1.75)]

    current_y = table_y + Inches(0.55)
    for label, value, height in zip(labels, values, heights):
        add_body_cell(slide, table_x, current_y, left_w, height, label, PP_ALIGN.LEFT, PALE, 14, True)
        add_body_cell(slide, table_x + left_w, current_y, right_w, height, value, PP_ALIGN.LEFT, WHITE, 14, False)
        current_y += height



# S001
s1 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s1, PALE)
add_band(s1, u(r"\u0422\u0438\u0442\u0443\u043b\u044c\u043d\u044b\u0439 \u044d\u043a\u0440\u0430\u043d \u043a\u0443\u0440\u0441\u0430"), "S001")
add_text(s1, Inches(0.65), Inches(1.15), Inches(6.5), Inches(0.55), u(r"\u042d\u0412 \u0421\u0442\u0440\u043e\u043f\u0430\u043b\u044c\u0449\u0438\u043a"), 30, True)
add_text(
    s1,
    Inches(0.67),
    Inches(1.95),
    Inches(4.5),
    Inches(0.3),
    u(r"\u0421\u0442\u0440\u043e\u043f\u0430\u043b\u044c\u0449\u0438\u043a 2-4 \u0440\u0430\u0437\u0440\u044f\u0434\u0430"),
    18,
    True,
    BLUE,
)
add_text(
    s1,
    Inches(0.67),
    Inches(2.28),
    Inches(4.5),
    Inches(0.3),
    u(r"\u0421\u0442\u0440\u043e\u043f\u0430\u043b\u044c\u0449\u0438\u043a 5-6 \u0440\u0430\u0437\u0440\u044f\u0434\u0430"),
    18,
    True,
    ORANGE,
)
add_text(
    s1,
    Inches(0.67),
    Inches(2.85),
    Inches(6.2),
    Inches(0.42),
    u(r"\u041f\u0440\u0430\u043a\u0442\u0438\u0447\u0435\u0441\u043a\u0438\u0439 \u043a\u0443\u0440\u0441 \u043f\u043e \u0431\u0435\u0437\u043e\u043f\u0430\u0441\u043d\u043e\u0439 \u0440\u0430\u0431\u043e\u0442\u0435 \u0441 \u0433\u0440\u0443\u0437\u0430\u043c\u0438"),
    20,
    True,
)
note = s1.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.67), Inches(3.52), Inches(6.2), Inches(1.08))
note.fill.solid()
note.fill.fore_color.rgb = WHITE
note.line.color.rgb = RGBColor(210, 218, 226)
tf = note.text_frame
tf.clear()
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]
r = p.add_run()
r.text = u(r"\u0411\u0430\u0437\u043e\u0432\u0430\u044f \u0447\u0430\u0441\u0442\u044c \u043a\u0443\u0440\u0441\u0430 - \u0434\u043b\u044f 2-4 \u0440\u0430\u0437\u0440\u044f\u0434\u0430. \u0414\u043b\u044f 5-6 \u0440\u0430\u0437\u0440\u044f\u0434\u0430 \u043e\u0442\u043a\u0440\u044b\u0432\u0430\u0439\u0442\u0435 \u0434\u043e\u043f\u043e\u043b\u043d\u0438\u0442\u0435\u043b\u044c\u043d\u044b\u0435 \u0431\u043b\u043e\u043a\u0438.")
r.font.size = Pt(17)
r.font.color.rgb = TEXT

right = s1.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(7.9), Inches(0.72), Inches(5.43), Inches(6.78))
right.fill.solid()
right.fill.fore_color.rgb = STEEL
right.line.fill.background()
add_text(s1, Inches(8.35), Inches(1.2), Inches(4.3), Inches(0.3), u(r"\u0416\u0438\u0432\u0430\u044f demo-\u0441\u0431\u043e\u0440\u043a\u0430"), 18, True, WHITE, PP_ALIGN.CENTER)
add_text(
    s1,
    Inches(8.3),
    Inches(1.8),
    Inches(4.45),
    Inches(2.7),
    u(
        r"\u0421\u0435\u0439\u0447\u0430\u0441 \u043c\u043e\u0436\u043d\u043e \u043f\u043e\u0441\u043c\u043e\u0442\u0440\u0435\u0442\u044c:\n\n\u2022 S001 \u0442\u0438\u0442\u0443\u043b\n\u2022 S002 \u044d\u043a\u0440\u0430\u043d \u0432\u044b\u0431\u043e\u0440\u0430\n\u2022 S002-P01 \u0432\u0438\u0434\u0435\u043e\u0438\u043d\u0441\u0442\u0440\u0443\u043a\u0446\u0438\u044f\n\u2022 S002-P02 \u0441\u043b\u043e\u0432\u0430\u0440\u044c \u0442\u0435\u0440\u043c\u0438\u043d\u043e\u0432"
    ),
    18,
    False,
    WHITE,
)
start_btn = add_button(s1, Inches(0.67), Inches(5.15), Inches(2.6), Inches(0.62), u(r"\u041d\u0410\u0427\u0410\u0422\u042c \u041a\u0423\u0420\u0421"), ORANGE)


# S002
s2 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s2, PALE)
add_band(
    s2,
    u(r"\u0412\u0438\u0434\u0435\u043e\u0438\u043d\u0441\u0442\u0440\u0443\u043a\u0446\u0438\u044f \u043f\u043e \u043a\u0443\u0440\u0441\u0443 \u0438 \u0433\u043b\u043e\u0441\u0430\u0440\u0438\u0439"),
    "S002",
    u(r"\u0422\u0435\u043c\u0430 1 \u2022 \u0432\u0432\u043e\u0434\u043d\u044b\u0439 \u044d\u043a\u0440\u0430\u043d"),
)
add_text(s2, Inches(0.65), Inches(1.0), Inches(4.5), Inches(0.3), u(r"\u0412\u044b\u0431\u0435\u0440\u0438\u0442\u0435, \u0447\u0442\u043e \u043e\u0442\u043a\u0440\u044b\u0442\u044c:"), 20, True)
card_video = add_panel(
    s2,
    Inches(0.7),
    Inches(1.65),
    Inches(5.7),
    Inches(3.9),
    u(r"\u0412\u0438\u0434\u0435\u043e\u0438\u043d\u0441\u0442\u0440\u0443\u043a\u0446\u0438\u044f"),
    [
        u(r"\u042d\u0442\u043e \u043d\u0430\u0432\u0438\u0433\u0430\u0446\u0438\u044f \u043f\u043e \u043a\u0443\u0440\u0441\u0443."),
        u(r"\u041e\u0442\u043a\u0440\u043e\u0439\u0442\u0435 \u0438 \u043e\u0437\u043d\u0430\u043a\u043e\u043c\u044c\u0442\u0435\u0441\u044c."),
        u(r"\u041f\u0435\u0440\u0435\u0445\u043e\u0434 \u0432 \u043f\u043e\u0434\u0432\u0430\u043b S002-P01."),
    ],
    BLUE,
)
card_gloss = add_panel(
    s2,
    Inches(6.9),
    Inches(1.65),
    Inches(5.7),
    Inches(3.9),
    u(r"\u0421\u043b\u043e\u0432\u0430\u0440\u044c \u0442\u0435\u0440\u043c\u0438\u043d\u043e\u0432"),
    [
        u(r"\u041e\u0442\u043a\u0440\u043e\u0439\u0442\u0435 \u0438 \u043f\u0440\u043e\u0447\u0442\u0438\u0442\u0435."),
        u(r"\u0422\u0435\u0440\u043c\u0438\u043d\u044b \u0432\u044b\u043d\u0435\u0441\u0435\u043d\u044b \u0432 \u043e\u0442\u0434\u0435\u043b\u044c\u043d\u044b\u0439 \u043f\u043e\u0434\u0432\u0430\u043b."),
        u(r"\u041f\u0435\u0440\u0435\u0445\u043e\u0434 \u0432 \u043f\u043e\u0434\u0432\u0430\u043b S002-P02."),
    ],
    GREEN,
)
btn_video = add_button(s2, Inches(1.75), Inches(5.0), Inches(3.55), Inches(0.62), u(r"\u041e\u0422\u041a\u0420\u042b\u0422\u042c S002-P01"), BLUE)
btn_gloss = add_button(s2, Inches(7.95), Inches(5.0), Inches(3.55), Inches(0.62), u(r"\u041e\u0422\u041a\u0420\u042b\u0422\u042c S002-P02"), GREEN)
back_s2, next_s2 = add_linear_nav(s2, u(r"\u041d\u0410\u0417\u0410\u0414"), u(r"\u0414\u0410\u041b\u0415\u0415"))


# S002-P01
s3 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s3, WHITE)
add_band(
    s3,
    u(r"\u0412\u0438\u0434\u0435\u043e\u0438\u043d\u0441\u0442\u0440\u0443\u043a\u0446\u0438\u044f \u043f\u043e \u043a\u0443\u0440\u0441\u0443"),
    "S002-P01",
    u(r"\u041f\u043e\u0434\u0432\u0430\u043b \u043e\u0442 S002"),
)
preview = s3.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(1.25), Inches(4.8), Inches(3.85))
preview.fill.solid()
preview.fill.fore_color.rgb = RGBColor(224, 232, 239)
preview.line.color.rgb = STEEL
add_text(s3, Inches(1.05), Inches(1.58), Inches(4.1), Inches(0.45), u(r"\u0417\u0434\u0435\u0441\u044c \u0431\u0443\u0434\u0435\u0442 \u0432\u0438\u0434\u0435\u043e\u0440\u043e\u043b\u0438\u043a"), 24, True, NAVY, PP_ALIGN.CENTER)
add_text(
    s3,
    Inches(0.95),
    Inches(2.35),
    Inches(4.25),
    Inches(1.1),
    u(r"\u0420\u0430\u0431\u043e\u0447\u0430\u044f \u043e\u0437\u0432\u0443\u0447\u043a\u0430 \u0443\u0436\u0435 \u043f\u043e\u0434\u0433\u043e\u0442\u043e\u0432\u043b\u0435\u043d\u0430:\nS002-P01_ozvuchka_55_sec_stress_2026-06-20.mp3"),
    17,
    False,
    TEXT,
    PP_ALIGN.CENTER,
)
open_video_btn = add_button(
    s3,
    Inches(1.05),
    Inches(4.48),
    Inches(4.1),
    Inches(0.58),
    u(r"\u041e\u0422\u041a\u0420\u042b\u0422\u042c \u0412\u0418\u0414\u0415\u041e\u0418\u041d\u0421\u0422\u0420\u0423\u041a\u0426\u0418\u042e"),
    ORANGE,
)
add_text(s3, Inches(5.95), Inches(1.3), Inches(5.9), Inches(0.3), u(r"\u0427\u0442\u043e \u0434\u043e\u043b\u0436\u0435\u043d \u043f\u043e\u043d\u044f\u0442\u044c \u0441\u043b\u0443\u0448\u0430\u0442\u0435\u043b\u044c"), 20, True)
add_panel(
    s3,
    Inches(5.9),
    Inches(1.8),
    Inches(6.0),
    Inches(3.55),
    u(r"\u041a\u043e\u0440\u043e\u0442\u043a\u043e \u043e \u043b\u043e\u0433\u0438\u043a\u0435 \u043a\u0443\u0440\u0441\u0430"),
    [
        u(r"\u041a\u0443\u0440\u0441 \u0438\u043d\u0442\u0435\u0440\u0430\u043a\u0442\u0438\u0432\u043d\u044b\u0439."),
        u(r"\u0422\u0435\u043c\u044b \u0438\u0434\u0443\u0442 \u043f\u043e \u043f\u043e\u0440\u044f\u0434\u043a\u0443."),
        u(r"\u041f\u043e\u0441\u043b\u0435 \u0442\u0435\u043c \u0438 \u0431\u043b\u043e\u043a\u043e\u0432 \u0431\u0443\u0434\u0443\u0442 \u0442\u0435\u0441\u0442\u044b."),
        u(r"\u041f\u043e \u043a\u043d\u043e\u043f\u043a\u0430\u043c \u043e\u0442\u043a\u0440\u044b\u0432\u0430\u0439\u0442\u0435 \u043f\u043e\u044f\u0441\u043d\u0435\u043d\u0438\u044f."),
        u(r"\u041f\u043e\u0441\u043b\u0435 \u043f\u0440\u043e\u0441\u043c\u043e\u0442\u0440\u0430 \u0432\u0435\u0440\u043d\u0438\u0442\u0435\u0441\u044c \u0432 \u043a\u0443\u0440\u0441."),
    ],
    BLUE,
)
back1 = add_button(s3, Inches(9.15), Inches(6.55), Inches(3.05), Inches(0.5), u(r"\u041d\u0410\u0417\u0410\u0414 \u041a \u041a\u0423\u0420\u0421\u0423"), NAVY)


# S002-P02
s4 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s4, WHITE)
add_band(
    s4,
    u(r"\u0421\u043b\u043e\u0432\u0430\u0440\u044c \u0442\u0435\u0440\u043c\u0438\u043d\u043e\u0432"),
    "S002-P02",
    u(r"\u041f\u043e\u0434\u0432\u0430\u043b \u043e\u0442 S002"),
)
glossary_panel = s4.shapes.add_shape(
    MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.62), Inches(1.25), Inches(12.08), Inches(5.05)
)
glossary_panel.fill.solid()
glossary_panel.fill.fore_color.rgb = WHITE
glossary_panel.line.color.rgb = GREEN
glossary_panel.line.width = Pt(1.5)

add_text(
    s4,
    Inches(0.95),
    Inches(1.48),
    Inches(5.2),
    Inches(0.24),
    u(r"\u0422\u0435\u0440\u043c\u0438\u043d\u044b \u0440\u0430\u0441\u043f\u043e\u043b\u043e\u0436\u0435\u043d\u044b \u043f\u043e \u0430\u043b\u0444\u0430\u0432\u0438\u0442\u0443"),
    11,
    False,
    GRAY,
)

divider = s4.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(6.64), Inches(1.78), Inches(0.02), Inches(4.15))
divider.fill.solid()
divider.fill.fore_color.rgb = LINE
divider.line.fill.background()

glossary_items = [
    (
        u(r"\u0413\u0430\u043a"),
        u(r"\u043a\u0440\u044e\u043a \u0433\u0440\u0443\u0437\u043e\u043f\u043e\u0434\u044a\u0451\u043c\u043d\u043e\u0433\u043e \u043c\u0435\u0445\u0430\u043d\u0438\u0437\u043c\u0430."),
    ),
    (
        u(r"\u0413\u0440\u0443\u0437"),
        u(r"\u043f\u0440\u0435\u0434\u043c\u0435\u0442 \u0438\u043b\u0438 \u043a\u043e\u043d\u0441\u0442\u0440\u0443\u043a\u0446\u0438\u044f \u0434\u043b\u044f \u043f\u043e\u0434\u044a\u0451\u043c\u0430."),
    ),
    (
        u(r"\u0413\u0440\u0443\u0437\u043e\u0437\u0430\u0445\u0432\u0430\u0442\u043d\u043e\u0435 \u043f\u0440\u0438\u0441\u043f\u043e\u0441\u043e\u0431\u043b\u0435\u043d\u0438\u0435"),
        u(r"\u0443\u0441\u0442\u0440\u043e\u0439\u0441\u0442\u0432\u043e \u0434\u043b\u044f \u0437\u0430\u0445\u0432\u0430\u0442\u0430 \u0433\u0440\u0443\u0437\u0430."),
    ),
    (
        u(r"\u0417\u0430\u0436\u0438\u043c\u044b"),
        u(r"\u0443\u0441\u0442\u0440\u043e\u0439\u0441\u0442\u0432\u0430 \u0434\u043b\u044f \u0437\u0430\u043a\u0440\u0435\u043f\u043b\u0435\u043d\u0438\u044f \u043a\u043e\u043d\u0446\u043e\u0432 \u043a\u0430\u043d\u0430\u0442\u0430."),
    ),
    (
        u(r"\u041a\u043e\u0443\u0448"),
        u(r"\u043c\u0435\u0442\u0430\u043b\u043b\u0438\u0447\u0435\u0441\u043a\u0430\u044f \u0432\u0441\u0442\u0430\u0432\u043a\u0430, \u0437\u0430\u0449\u0438\u0449\u0430\u044e\u0449\u0430\u044f \u043f\u0435\u0442\u043b\u044e \u043a\u0430\u043d\u0430\u0442\u0430 \u043e\u0442 \u0437\u0430\u043b\u043e\u043c\u0430 \u0438 \u0438\u0437\u043d\u043e\u0441\u0430."),
    ),
    (
        u(r"\u041c\u0430\u0433\u043d\u0438\u0442\u043d\u044b\u0435 \u0437\u0430\u0436\u0438\u043c\u044b"),
        u(r"\u0443\u0441\u0442\u0440\u043e\u0439\u0441\u0442\u0432\u0430, \u0443\u0434\u0435\u0440\u0436\u0438\u0432\u0430\u044e\u0449\u0438\u0435 \u043c\u0435\u0442\u0430\u043b\u043b\u0438\u0447\u0435\u0441\u043a\u0438\u0439 \u0433\u0440\u0443\u0437 \u0441 \u043f\u043e\u043c\u043e\u0449\u044c\u044e \u043c\u0430\u0433\u043d\u0438\u0442\u0430."),
    ),
    (
        u(r"\u041c\u0451\u0440\u0442\u0432\u044b\u0439 \u0433\u0440\u0443\u0437"),
        u(r"\u0433\u0440\u0443\u0437, \u043c\u0430\u0441\u0441\u0443 \u043a\u043e\u0442\u043e\u0440\u043e\u0433\u043e \u043e\u043f\u0440\u0435\u0434\u0435\u043b\u0438\u0442\u044c \u043d\u0435\u0432\u043e\u0437\u043c\u043e\u0436\u043d\u043e."),
    ),
    (
        u(r"\u041e\u0442\u0442\u044f\u0436\u043a\u0430"),
        u(r"\u043a\u0430\u043d\u0430\u0442 \u0438\u043b\u0438 \u0432\u0435\u0440\u0451\u0432\u043a\u0430 \u0434\u043b\u044f \u0443\u043f\u0440\u0430\u0432\u043b\u0435\u043d\u0438\u044f \u043f\u043e\u043b\u043e\u0436\u0435\u043d\u0438\u0435\u043c \u0433\u0440\u0443\u0437\u0430 \u0438 \u043f\u0440\u0435\u0434\u043e\u0442\u0432\u0440\u0430\u0449\u0435\u043d\u0438\u044f \u0435\u0433\u043e \u0440\u0430\u0441\u043a\u0430\u0447\u0438\u0432\u0430\u043d\u0438\u044f."),
    ),
    (
        u(r"\u0420\u044b\u043c"),
        u(r"\u043a\u043e\u043b\u044c\u0446\u043e \u0438\u043b\u0438 \u043f\u0440\u043e\u0443\u0448\u0438\u043d\u0430 \u0434\u043b\u044f \u043a\u0440\u0435\u043f\u043b\u0435\u043d\u0438\u044f \u0441\u0442\u0440\u043e\u043f\u0430 \u043a \u0433\u0440\u0443\u0437\u0443."),
    ),
    (
        u(r"\u0421\u0438\u0433\u043d\u0430\u043b\u044c\u0449\u0438\u043a"),
        u(r"\u0440\u0430\u0431\u043e\u0442\u043d\u0438\u043a, \u043a\u043e\u0442\u043e\u0440\u044b\u0439 \u043f\u043e\u0434\u0430\u0451\u0442 \u043a\u043e\u043c\u0430\u043d\u0434\u044b \u043a\u0440\u0430\u043d\u043e\u0432\u0449\u0438\u043a\u0443 \u043f\u0440\u0438 \u043f\u0435\u0440\u0435\u043c\u0435\u0449\u0435\u043d\u0438\u0438 \u0433\u0440\u0443\u0437\u0430, \u0435\u0441\u043b\u0438 \u043c\u0435\u0436\u0434\u0443 \u0441\u0442\u0440\u043e\u043f\u0430\u043b\u044c\u0449\u0438\u043a\u043e\u043c \u0438 \u043a\u0440\u0430\u043d\u043e\u0432\u0449\u0438\u043a\u043e\u043c \u043a\u043e\u043d\u0442\u0430\u043a\u0442 \u043e\u0442\u0441\u0443\u0442\u0441\u0442\u0432\u0443\u0435\u0442."),
    ),
    (
        u(r"\u0421\u0418\u0417"),
        u(r"\u0441\u0440\u0435\u0434\u0441\u0442\u0432\u0430 \u0438\u043d\u0434\u0438\u0432\u0438\u0434\u0443\u0430\u043b\u044c\u043d\u043e\u0439 \u0437\u0430\u0449\u0438\u0442\u044b."),
    ),
    (
        u(r"\u0421\u0442\u0440\u043e\u043f"),
        u(r"\u0433\u0440\u0443\u0437\u043e\u0437\u0430\u0445\u0432\u0430\u0442\u043d\u043e\u0435 \u0441\u0440\u0435\u0434\u0441\u0442\u0432\u043e."),
    ),
    (
        u(r"\u0421\u0442\u0440\u043e\u043f\u0430\u043b\u044c\u0449\u0438\u043a"),
        u(r"\u0440\u0430\u0431\u043e\u0447\u0438\u0439, \u043a\u043e\u0442\u043e\u0440\u044b\u0439 \u0432\u044b\u043f\u043e\u043b\u043d\u044f\u0435\u0442 \u0441\u0442\u0440\u043e\u043f\u043e\u0432\u043a\u0443 \u0433\u0440\u0443\u0437\u0430."),
    ),
    (
        u(r"\u0421\u0442\u0440\u043e\u043f\u043e\u0432\u043a\u0430"),
        u(r"\u0441\u043f\u043e\u0441\u043e\u0431 \u043e\u0431\u0432\u044f\u0437\u043a\u0438 \u0438 \u0437\u0430\u0446\u0435\u043f\u043a\u0438."),
    ),
    (
        u(r"\u0422\u0430\u0440\u0430"),
        u(r"\u0451\u043c\u043a\u043e\u0441\u0442\u044c \u0438\u043b\u0438 \u043f\u0440\u0438\u0441\u043f\u043e\u0441\u043e\u0431\u043b\u0435\u043d\u0438\u0435 \u0434\u043b\u044f \u0440\u0430\u0437\u043c\u0435\u0449\u0435\u043d\u0438\u044f \u0438 \u043f\u0435\u0440\u0435\u043c\u0435\u0449\u0435\u043d\u0438\u044f \u0433\u0440\u0443\u0437\u0430."),
    ),
    (
        u(r"\u0422\u0440\u0430\u0432\u0435\u0440\u0441\u0430"),
        u(r"\u0433\u0440\u0443\u0437\u043e\u0437\u0430\u0445\u0432\u0430\u0442\u043d\u0430\u044f \u0431\u0430\u043b\u043a\u0430 \u0434\u043b\u044f \u043f\u043e\u0434\u044a\u0451\u043c\u0430 \u0434\u043b\u0438\u043d\u043d\u044b\u0445 \u0438\u043b\u0438 \u043a\u0440\u0443\u043f\u043d\u043e\u0433\u0430\u0431\u0430\u0440\u0438\u0442\u043d\u044b\u0445 \u0433\u0440\u0443\u0437\u043e\u0432."),
    ),
    (
        u(r"\u0426\u0430\u043f\u0444\u0430"),
        u(r"\u0432\u044b\u0441\u0442\u0443\u043f\u0430\u044e\u0449\u0430\u044f \u0447\u0430\u0441\u0442\u044c \u0433\u0440\u0443\u0437\u0430 \u0438\u043b\u0438 \u043e\u0431\u043e\u0440\u0443\u0434\u043e\u0432\u0430\u043d\u0438\u044f, \u0438\u0441\u043f\u043e\u043b\u044c\u0437\u0443\u0435\u043c\u0430\u044f \u043a\u0430\u043a \u043e\u043f\u043e\u0440\u0430 \u0438\u043b\u0438 \u043c\u0435\u0441\u0442\u043e \u0441\u0442\u0440\u043e\u043f\u043e\u0432\u043a\u0438."),
    ),
]

left_items = glossary_items[:9]
right_items = glossary_items[9:]

add_glossary_column(s4, Inches(0.95), Inches(1.88), Inches(5.35), Inches(4.1), left_items, 12, 11)
add_glossary_column(s4, Inches(6.95), Inches(1.88), Inches(5.35), Inches(4.1), right_items, 12, 11)
back2 = add_button(s4, Inches(9.15), Inches(6.55), Inches(3.05), Inches(0.5), u(r"\u041d\u0410\u0417\u0410\u0414 \u041a \u041a\u0423\u0420\u0421\u0423"), NAVY)


# S003
s5 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s5, PALE)
add_band(
    s5,
    u(r"\u041a\u0430\u043a\u0438\u0435 \u0434\u043e\u043a\u0443\u043c\u0435\u043d\u0442\u044b \u0440\u0435\u0433\u0443\u043b\u0438\u0440\u0443\u044e\u0442 \u0440\u0430\u0431\u043e\u0442\u0443 \u0441\u0442\u0440\u043e\u043f\u0430\u043b\u044c\u0449\u0438\u043a\u0430"),
    "S003",
    u(r"\u0422\u0435\u043c\u0430 1 \u2022 \u043e\u0431\u0449\u0430\u044f \u0447\u0430\u0441\u0442\u044c"),
)
add_text(
    s5,
    Inches(0.7),
    Inches(1.02),
    Inches(8.6),
    Inches(0.45),
    u(r"\u042d\u0442\u043e \u043e\u0441\u043d\u043e\u0432\u043d\u044b\u0435 \u043f\u0440\u0430\u0432\u043e\u0432\u044b\u0435 \u0430\u043a\u0442\u044b, \u0440\u0435\u0433\u0443\u043b\u0438\u0440\u0443\u044e\u0449\u0438\u0435 \u0440\u0430\u0431\u043e\u0442\u0443 \u0441\u0442\u0440\u043e\u043f\u0430\u043b\u044c\u0449\u0438\u043a\u0430"),
    20,
    True,
    NAVY,
)

table_x = Inches(0.7)
table_y = Inches(1.85)
law_w = Inches(6.15)
date_w = Inches(1.65)
issuer_w = Inches(3.35)
row_h = Inches(0.82)

add_header_cell(s5, table_x, table_y, law_w, Inches(0.56), u(r"\u0417\u0430\u043a\u043e\u043d"), BLUE)
add_header_cell(s5, table_x + law_w, table_y, date_w, Inches(0.56), u(r"\u0414\u0430\u0442\u0430 \u043f\u0440\u0438\u043d\u044f\u0442\u0438\u044f"), BLUE, 13)
add_header_cell(s5, table_x + law_w + date_w, table_y, issuer_w, Inches(0.56), u(r"\u041a\u0442\u043e \u0438\u0437\u0434\u0430\u043b"), BLUE)

rows = [
    (
        u(r"\u0422\u0440\u0443\u0434\u043e\u0432\u043e\u0439 \u043a\u043e\u0434\u0435\u043a\u0441 \u0420\u0424 \u2116 197-\u0424\u0417"),
        "30.12.2001",
        u(r"\u0424\u0435\u0434\u0435\u0440\u0430\u043b\u044c\u043d\u043e\u0435 \u0421\u043e\u0431\u0440\u0430\u043d\u0438\u0435 \u0420\u0424,\n\u041f\u0440\u0435\u0437\u0438\u0434\u0435\u043d\u0442 \u0420\u0424"),
    ),
    (
        u(r"\u041f\u043e\u0441\u0442\u0430\u043d\u043e\u0432\u043b\u0435\u043d\u0438\u0435 \u041f\u0440\u0430\u0432\u0438\u0442\u0435\u043b\u044c\u0441\u0442\u0432\u0430 \u0420\u0424 \u2116 2464"),
        "24.12.2021",
        u(r"\u041f\u0440\u0430\u0432\u0438\u0442\u0435\u043b\u044c\u0441\u0442\u0432\u043e \u0420\u0424"),
    ),
    (
        u(r"\u0424\u0435\u0434\u0435\u0440\u0430\u043b\u044c\u043d\u044b\u0439 \u0437\u0430\u043a\u043e\u043d \u2116 116-\u0424\u0417"),
        "21.07.1997",
        u(r"\u0424\u0435\u0434\u0435\u0440\u0430\u043b\u044c\u043d\u043e\u0435 \u0421\u043e\u0431\u0440\u0430\u043d\u0438\u0435 \u0420\u0424,\n\u041f\u0440\u0435\u0437\u0438\u0434\u0435\u043d\u0442 \u0420\u0424"),
    ),
    (
        u(r"\u041f\u0440\u0438\u043a\u0430\u0437 \u0420\u043e\u0441\u0442\u0435\u0445\u043d\u0430\u0434\u0437\u043e\u0440\u0430 \u2116 461"),
        "26.11.2020",
        u(r"\u0420\u043e\u0441\u0442\u0435\u0445\u043d\u0430\u0434\u0437\u043e\u0440"),
    ),
]

main_click_shapes = []
current_y = table_y + Inches(0.56)
for idx, row in enumerate(rows):
    fill = WHITE if idx % 2 == 0 else RGBColor(248, 250, 252)
    law_cell = add_body_cell(s5, table_x, current_y, law_w, row_h, row[0], PP_ALIGN.LEFT, fill, 15, True)
    add_body_cell(s5, table_x + law_w, current_y, date_w, row_h, row[1], PP_ALIGN.CENTER, fill, 14, False)
    add_body_cell(s5, table_x + law_w + date_w, current_y, issuer_w, row_h, row[2], PP_ALIGN.LEFT, fill, 13, False)
    main_click_shapes.append(law_cell)
    current_y += row_h

info_panel_shapes = add_info_panel(
    s5,
    Inches(10.45),
    Inches(1.85),
    Inches(2.15),
    Inches(2.15),
    u(r"\u0418\u043d\u0442\u0435\u0440\u0430\u043a\u0442\u0438\u0432"),
    [
        u(r"\u041d\u0430\u0436\u043c\u0438\u0442\u0435 \u043d\u0430 \u043d\u0430\u0437\u0432\u0430\u043d\u0438\u0435 \u0434\u043e\u043a\u0443\u043c\u0435\u043d\u0442\u0430."),
    ],
    GREEN,
)
back_s3, next_s3 = add_linear_nav(s5, u(r"\u041d\u0410\u0417\u0410\u0414"), u(r"\u0414\u0410\u041b\u0415\u0415"))


# S003-P01
s5_p01 = prs.slides.add_slide(prs.slide_layouts[6])
add_detail_table(
    s5_p01,
    u(r"\u0422\u0440\u0443\u0434\u043e\u0432\u043e\u0439 \u043a\u043e\u0434\u0435\u043a\u0441 \u0420\u0424 \u2116 197-\u0424\u0417"),
    "S003-P01",
    u(r"\u0422\u0440\u0443\u0434\u043e\u0432\u043e\u0439 \u043a\u043e\u0434\u0435\u043a\u0441 \u0420\u0424 \u2116 197-\u0424\u0417"),
    u(r"\u041e\u0431\u0449\u0430\u044f \u0440\u0430\u043c\u043a\u0430 \u043f\u043e \u043e\u0445\u0440\u0430\u043d\u0435 \u0442\u0440\u0443\u0434\u0430 \u0438 \u0434\u043e\u043f\u0443\u0441\u043a\u0443 \u043a \u0440\u0430\u0431\u043e\u0442\u0435"),
    u(r"\u0420\u0430\u0437\u0434\u0435\u043b X \"\u041e\u0445\u0440\u0430\u043d\u0430 \u0442\u0440\u0443\u0434\u0430\": \u043e\u0431\u044f\u0437\u0430\u043d\u043d\u043e\u0441\u0442\u0438 \u0440\u0430\u0431\u043e\u0442\u043e\u0434\u0430\u0442\u0435\u043b\u044f \u0438 \u0440\u0430\u0431\u043e\u0442\u043d\u0438\u043a\u0430, \u0421\u0418\u0417, \u0431\u0435\u0437\u043e\u043f\u0430\u0441\u043d\u044b\u0435 \u0443\u0441\u043b\u043e\u0432\u0438\u044f \u0442\u0440\u0443\u0434\u0430, \u0440\u0430\u0441\u0441\u043b\u0435\u0434\u043e\u0432\u0430\u043d\u0438\u0435 \u043d\u0435\u0441\u0447\u0430\u0441\u0442\u043d\u044b\u0445 \u0441\u043b\u0443\u0447\u0430\u0435\u0432"),
    BLUE,
)
p01_back_main = add_button(s5_p01, Inches(0.7), Inches(6.55), Inches(2.35), Inches(0.5), u(r"\u041d\u0410\u0417\u0410\u0414"), STEEL, 18)


# S003-P02
s5_p02 = prs.slides.add_slide(prs.slide_layouts[6])
add_detail_table(
    s5_p02,
    u(r"\u041f\u043e\u0441\u0442\u0430\u043d\u043e\u0432\u043b\u0435\u043d\u0438\u0435 \u041f\u0440\u0430\u0432\u0438\u0442\u0435\u043b\u044c\u0441\u0442\u0432\u0430 \u0420\u0424 \u2116 2464"),
    "S003-P02",
    u(r"\u041f\u043e\u0441\u0442\u0430\u043d\u043e\u0432\u043b\u0435\u043d\u0438\u0435 \u041f\u0440\u0430\u0432\u0438\u0442\u0435\u043b\u044c\u0441\u0442\u0432\u0430 \u0420\u0424 \u2116 2464"),
    u(r"\u041e\u0431\u0443\u0447\u0435\u043d\u0438\u0435 \u043f\u043e \u043e\u0445\u0440\u0430\u043d\u0435 \u0442\u0440\u0443\u0434\u0430 \u0438 \u043f\u0440\u043e\u0432\u0435\u0440\u043a\u0430 \u0437\u043d\u0430\u043d\u0438\u0439"),
    u(r"\u0418\u043d\u0441\u0442\u0440\u0443\u043a\u0442\u0430\u0436\u0438, \u0441\u0442\u0430\u0436\u0438\u0440\u043e\u0432\u043a\u0430, \u043e\u0431\u0443\u0447\u0435\u043d\u0438\u0435, \u043f\u0435\u0440\u0432\u0430\u044f \u043f\u043e\u043c\u043e\u0449\u044c, \u043f\u0440\u043e\u0432\u0435\u0440\u043a\u0430 \u0437\u043d\u0430\u043d\u0438\u0439"),
    GREEN,
)
p02_back_main = add_button(s5_p02, Inches(0.7), Inches(6.55), Inches(2.35), Inches(0.5), u(r"\u041d\u0410\u0417\u0410\u0414"), STEEL, 18)


# S003-P03
s5_p03 = prs.slides.add_slide(prs.slide_layouts[6])
add_detail_table(
    s5_p03,
    u(r"\u0424\u0435\u0434\u0435\u0440\u0430\u043b\u044c\u043d\u044b\u0439 \u0437\u0430\u043a\u043e\u043d \u2116 116-\u0424\u0417"),
    "S003-P03",
    u(r"\u0424\u0435\u0434\u0435\u0440\u0430\u043b\u044c\u043d\u044b\u0439 \u0437\u0430\u043a\u043e\u043d \u2116 116-\u0424\u0417"),
    u(r"\u0423\u0441\u0442\u0430\u043d\u0430\u0432\u043b\u0438\u0432\u0430\u0435\u0442 \u0442\u0440\u0435\u0431\u043e\u0432\u0430\u043d\u0438\u044f \u0431\u0435\u0437\u043e\u043f\u0430\u0441\u043d\u043e\u0439 \u0440\u0430\u0431\u043e\u0442\u044b \u043d\u0430 \u043e\u043f\u0430\u0441\u043d\u044b\u0445 \u043f\u0440\u043e\u0438\u0437\u0432\u043e\u0434\u0441\u0442\u0432\u0435\u043d\u043d\u044b\u0445 \u043e\u0431\u044a\u0435\u043a\u0442\u0430\u0445, \u043f\u043e\u0434\u0433\u043e\u0442\u043e\u0432\u043a\u0438 \u0440\u0430\u0431\u043e\u0442\u043d\u0438\u043a\u043e\u0432, \u043f\u0440\u0435\u0434\u0443\u043f\u0440\u0435\u0436\u0434\u0435\u043d\u0438\u044f \u0430\u0432\u0430\u0440\u0438\u0439 \u0438 \u043f\u0440\u043e\u0438\u0437\u0432\u043e\u0434\u0441\u0442\u0432\u0435\u043d\u043d\u043e\u0433\u043e \u043a\u043e\u043d\u0442\u0440\u043e\u043b\u044f"),
    u(r"\u041d\u0443\u0436\u0435\u043d \u043f\u0440\u0438 \u0440\u0430\u0431\u043e\u0442\u0435 \u043d\u0430 \u041e\u041f\u041e: \u043f\u043e\u0434\u0433\u043e\u0442\u043e\u0432\u043a\u0430 \u0440\u0430\u0431\u043e\u0442\u043d\u0438\u043a\u043e\u0432, \u0442\u0440\u0435\u0431\u043e\u0432\u0430\u043d\u0438\u044f \u043f\u0440\u043e\u043c\u0431\u0435\u0437\u043e\u043f\u0430\u0441\u043d\u043e\u0441\u0442\u0438, \u043f\u0440\u043e\u0438\u0437\u0432\u043e\u0434\u0441\u0442\u0432\u0435\u043d\u043d\u044b\u0439 \u043a\u043e\u043d\u0442\u0440\u043e\u043b\u044c. \u0414\u043b\u044f \u0441\u0442\u0440\u043e\u043f\u0430\u043b\u044c\u0449\u0438\u043a\u0430 \u0432\u0430\u0436\u0435\u043d \u0442\u0435\u043c, \u0447\u0442\u043e \u0442\u0440\u0435\u0431\u0443\u0435\u0442 \u043f\u0440\u043e\u0445\u043e\u0434\u0438\u0442\u044c \u043e\u0431\u0443\u0447\u0435\u043d\u0438\u0435, \u0441\u043e\u0431\u043b\u044e\u0434\u0430\u0442\u044c \u043f\u0440\u0430\u0432\u0438\u043b\u0430 \u0431\u0435\u0437\u043e\u043f\u0430\u0441\u043d\u043e\u0441\u0442\u0438, \u0432\u044b\u043f\u043e\u043b\u043d\u044f\u0442\u044c \u0438\u043d\u0441\u0442\u0440\u0443\u043a\u0446\u0438\u0438 \u0438 \u0441\u043e\u043e\u0431\u0449\u0430\u0442\u044c \u043e \u043d\u0435\u0438\u0441\u043f\u0440\u0430\u0432\u043d\u043e\u0441\u0442\u044f\u0445 \u0438 \u043e\u043f\u0430\u0441\u043d\u044b\u0445 \u0441\u0438\u0442\u0443\u0430\u0446\u0438\u044f\u0445"),
    STEEL,
)
p03_back_main = add_button(s5_p03, Inches(0.7), Inches(6.55), Inches(2.35), Inches(0.5), u(r"\u041d\u0410\u0417\u0410\u0414"), STEEL, 18)


# S003-P04
s5_p04 = prs.slides.add_slide(prs.slide_layouts[6])
add_detail_table(
    s5_p04,
    u(r"\u041f\u0440\u0438\u043a\u0430\u0437 \u0420\u043e\u0441\u0442\u0435\u0445\u043d\u0430\u0434\u0437\u043e\u0440\u0430 \u2116 461"),
    "S003-P04",
    u(r"\u041f\u0440\u0438\u043a\u0430\u0437 \u0420\u043e\u0441\u0442\u0435\u0445\u043d\u0430\u0434\u0437\u043e\u0440\u0430 \u2116 461"),
    u(r"\u0423\u0441\u0442\u0430\u043d\u0430\u0432\u043b\u0438\u0432\u0430\u0435\u0442 \u043f\u0440\u0430\u0432\u0438\u043b\u0430 \u0431\u0435\u0437\u043e\u043f\u0430\u0441\u043d\u043e\u0439 \u044d\u043a\u0441\u043f\u043b\u0443\u0430\u0442\u0430\u0446\u0438\u0438 \u043a\u0440\u0430\u043d\u043e\u0432 \u0438 \u0434\u0440\u0443\u0433\u0438\u0445 \u043f\u043e\u0434\u044a\u0435\u043c\u043d\u044b\u0445 \u0441\u043e\u043e\u0440\u0443\u0436\u0435\u043d\u0438\u0439 \u043d\u0430 \u043e\u043f\u0430\u0441\u043d\u044b\u0445 \u043f\u0440\u043e\u0438\u0437\u0432\u043e\u0434\u0441\u0442\u0432\u0435\u043d\u043d\u044b\u0445 \u043e\u0431\u044a\u0435\u043a\u0442\u0430\u0445"),
    u(r"\u043f. 19, \u043f. 98-99, \u043f. 155-163: \u043f\u0435\u0440\u0441\u043e\u043d\u0430\u043b, \u0441\u0438\u0433\u043d\u0430\u043b\u044b, \u041f\u041f\u0420/\u0422\u041a, \u0441\u0442\u0440\u043e\u043f\u043e\u0432\u043a\u0430, \u041f\u0420\u0420, \u0441\u043a\u043b\u0430\u0434\u0438\u0440\u043e\u0432\u0430\u043d\u0438\u0435. \u0414\u043b\u044f \u0441\u0442\u0440\u043e\u043f\u0430\u043b\u044c\u0449\u0438\u043a\u0430 \u0432\u0430\u0436\u0435\u043d \u0442\u0435\u043c, \u0447\u0442\u043e \u043e\u043f\u0440\u0435\u0434\u0435\u043b\u044f\u0435\u0442 \u043f\u0440\u0430\u0432\u0438\u043b\u0430 \u0441\u0442\u0440\u043e\u043f\u043e\u0432\u043a\u0438 \u0438 \u043f\u0435\u0440\u0435\u043c\u0435\u0449\u0435\u043d\u0438\u044f \u0433\u0440\u0443\u0437\u043e\u0432, \u043f\u0440\u043e\u0432\u0435\u0440\u043a\u0438 \u0441\u0442\u0440\u043e\u043f\u043e\u0432 \u0438 \u0442\u0430\u0440\u044b, \u043f\u043e\u0434\u0430\u0447\u0438 \u0441\u0438\u0433\u043d\u0430\u043b\u043e\u0432 \u0438 \u043f\u0440\u0435\u043a\u0440\u0430\u0449\u0435\u043d\u0438\u044f \u0440\u0430\u0431\u043e\u0442\u044b \u043f\u0440\u0438 \u043e\u043f\u0430\u0441\u043d\u043e\u0441\u0442\u0438"),
    ORANGE,
)
p04_back_main = add_button(s5_p04, Inches(0.7), Inches(6.55), Inches(2.35), Inches(0.5), u(r"\u041d\u0410\u0417\u0410\u0414"), STEEL, 18)


# S004
s6 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s6, PALE)
add_band(
    s6,
    u(r"\u041a\u0442\u043e \u0442\u0430\u043a\u043e\u0439 \u0441\u0442\u0440\u043e\u043f\u0430\u043b\u044c\u0449\u0438\u043a"),
    "S004",
    u(r"\u0422\u0435\u043c\u0430 1 \u2022 \u043e\u0431\u0449\u0430\u044f \u0447\u0430\u0441\u0442\u044c"),
)
add_panel(
    s6,
    Inches(0.7),
    Inches(1.45),
    Inches(5.8),
    Inches(4.7),
    u(r"\u041e\u0441\u043d\u043e\u0432\u043d\u044b\u0435 \u0434\u0435\u0439\u0441\u0442\u0432\u0438\u044f"),
    [
        u(r"\u041f\u043e\u0434\u0431\u0438\u0440\u0430\u0435\u0442 \u0438 \u043f\u0440\u043e\u0432\u0435\u0440\u044f\u0435\u0442 \u0441\u0442\u0440\u043e\u043f\u044b."),
        u(r"\u0421\u0442\u0440\u043e\u043f\u0438\u0442 \u0433\u0440\u0443\u0437 \u043f\u043e \u0441\u0445\u0435\u043c\u0435."),
        u(r"\u041f\u043e\u0434\u0430\u0435\u0442 \u0441\u0438\u0433\u043d\u0430\u043b\u044b \u043a\u0440\u0430\u043d\u043e\u0432\u0449\u0438\u043a\u0443."),
        u(r"\u0421\u043b\u0435\u0434\u0438\u0442 \u0437\u0430 \u0431\u0435\u0437\u043e\u043f\u0430\u0441\u043d\u044b\u043c \u043f\u0435\u0440\u0435\u043c\u0435\u0449\u0435\u043d\u0438\u0435\u043c \u0433\u0440\u0443\u0437\u0430."),
    ],
    BLUE,
)
add_panel(
    s6,
    Inches(6.9),
    Inches(1.45),
    Inches(5.7),
    Inches(4.7),
    u(r"\u041e\u0442\u0432\u0435\u0442\u0441\u0442\u0432\u0435\u043d\u043d\u043e\u0441\u0442\u044c"),
    [
        u(r"\u041e\u043d \u043e\u0442\u0432\u0435\u0447\u0430\u0435\u0442 \u043d\u0435 \u0442\u043e\u043b\u044c\u043a\u043e \u0437\u0430 \u0441\u0442\u0440\u043e\u043f, \u043d\u043e \u0438 \u0437\u0430 \u0431\u0435\u0437\u043e\u043f\u0430\u0441\u043d\u043e\u0441\u0442\u044c."),
        u(r"\u041e\u0448\u0438\u0431\u043a\u0430 \u043f\u0440\u0438 \u0441\u0442\u0440\u043e\u043f\u043e\u0432\u043a\u0435 \u043c\u043e\u0436\u0435\u0442 \u043f\u0440\u0438\u0432\u0435\u0441\u0442\u0438 \u043a \u0430\u0432\u0430\u0440\u0438\u0438."),
        u(r"\u041f\u043e\u044d\u0442\u043e\u043c\u0443 \u0432\u0441\u0435 \u0434\u0435\u0439\u0441\u0442\u0432\u0438\u044f \u0434\u043e\u043b\u0436\u043d\u044b \u0431\u044b\u0442\u044c \u0442\u043e\u0447\u043d\u044b\u043c\u0438 \u0438 \u043f\u043e \u043f\u0440\u0430\u0432\u0438\u043b\u0430\u043c."),
    ],
    GREEN,
)
back_s4, next_s4 = add_linear_nav(s6, u(r"\u041d\u0410\u0417\u0410\u0414"), u(r"\u0414\u0410\u041b\u0415\u0415"))


# S005
s7 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s7, PALE)
add_band(
    s7,
    u(r"\u041a\u0442\u043e \u043c\u043e\u0436\u0435\u0442 \u0440\u0430\u0431\u043e\u0442\u0430\u0442\u044c \u0441\u0442\u0440\u043e\u043f\u0430\u043b\u044c\u0449\u0438\u043a\u043e\u043c"),
    "S005",
    u(r"\u0422\u0435\u043c\u0430 1 \u2022 \u043e\u0431\u0449\u0430\u044f \u0447\u0430\u0441\u0442\u044c"),
)
add_panel(
    s7,
    Inches(0.7),
    Inches(1.45),
    Inches(5.8),
    Inches(4.7),
    u(r"\u0414\u043e\u043f\u0443\u0441\u043a \u043a \u0440\u0430\u0431\u043e\u0442\u0435"),
    [
        u(r"\u041f\u0440\u043e\u0448\u0435\u043b \u043e\u0431\u0443\u0447\u0435\u043d\u0438\u0435."),
        u(r"\u041f\u0440\u043e\u0448\u0435\u043b \u0438\u043d\u0441\u0442\u0440\u0443\u043a\u0442\u0430\u0436."),
        u(r"\u0421\u0434\u0430\u043b \u043f\u0440\u043e\u0432\u0435\u0440\u043a\u0443 \u0437\u043d\u0430\u043d\u0438\u0439."),
        u(r"\u041f\u043e\u043b\u0443\u0447\u0438\u043b \u0434\u043e\u043f\u0443\u0441\u043a \u0440\u0430\u0431\u043e\u0442\u043e\u0434\u0430\u0442\u0435\u043b\u044f."),
        u(r"\u0418\u0441\u043f\u043e\u043b\u044c\u0437\u0443\u0435\u0442 \u0421\u0418\u0417."),
    ],
    BLUE,
)
add_panel(
    s7,
    Inches(6.9),
    Inches(1.45),
    Inches(5.7),
    Inches(4.7),
    u(r"\u0412\u0430\u0436\u043d\u044b\u0439 \u0432\u044b\u0432\u043e\u0434"),
    [
        u(r"\u0414\u043e\u043f\u0443\u0441\u043a \u043a \u0440\u0430\u0431\u043e\u0442\u0435 - \u044d\u0442\u043e \u043d\u0435 \u0444\u043e\u0440\u043c\u0430\u043b\u044c\u043d\u043e\u0441\u0442\u044c."),
        u(r"\u0411\u0435\u0437 \u043f\u043e\u0434\u0433\u043e\u0442\u043e\u0432\u043a\u0438 \u0438 \u043f\u0440\u043e\u0432\u0435\u0440\u043a\u0438 \u0437\u043d\u0430\u043d\u0438\u0439 \u0440\u0430\u0431\u043e\u0442\u0430\u0442\u044c \u043d\u0435\u043b\u044c\u0437\u044f."),
        u(r"\u0411\u0435\u0437\u043e\u043f\u0430\u0441\u043d\u0430\u044f \u0440\u0430\u0431\u043e\u0442\u0430 \u043d\u0430\u0447\u0438\u043d\u0430\u0435\u0442\u0441\u044f \u0441 \u0434\u043e\u043f\u0443\u0441\u043a\u0430."),
    ],
    GREEN,
)
back_s5, next_s5 = add_linear_nav(s7, u(r"\u041d\u0410\u0417\u0410\u0414"), u(r"\u0414\u0410\u041b\u0415\u0415"))


# S006
s8 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s8, PALE)
add_band(
    s8,
    u(r"\u0420\u0430\u0437\u0440\u044f\u0434\u044b 2-4"),
    "S006",
    u(r"\u0422\u0435\u043c\u0430 1 \u2022 \u043e\u0431\u0449\u0430\u044f \u0447\u0430\u0441\u0442\u044c"),
)
add_panel(
    s8,
    Inches(0.7),
    Inches(1.45),
    Inches(5.8),
    Inches(4.7),
    u(r"\u0427\u0442\u043e \u0432\u0445\u043e\u0434\u0438\u0442 \u0432 \u043a\u0443\u0440\u0441"),
    [
        u(r"2 \u0440\u0430\u0437\u0440\u044f\u0434 - \u043f\u0440\u043e\u0441\u0442\u044b\u0435 \u0433\u0440\u0443\u0437\u044b \u0434\u043e 5 \u0442."),
        u(r"3 \u0440\u0430\u0437\u0440\u044f\u0434 - \u043f\u0440\u043e\u0441\u0442\u044b\u0435 \u0433\u0440\u0443\u0437\u044b \u0441\u0432\u044b\u0448\u0435 5 \u0442 \u0434\u043e 25 \u0442."),
        u(r"4 \u0440\u0430\u0437\u0440\u044f\u0434 - \u0441\u043b\u043e\u0436\u043d\u044b\u0435 \u0433\u0440\u0443\u0437\u044b \u0434\u043e 5 \u0442."),
    ],
    BLUE,
)
add_panel(
    s8,
    Inches(6.9),
    Inches(1.45),
    Inches(5.7),
    Inches(4.7),
    u(r"\u041a\u0430\u043a \u0447\u0438\u0442\u0430\u0442\u044c \u043c\u0430\u0440\u0448\u0440\u0443\u0442"),
    [
        u(r"\u0411\u0430\u0437\u043e\u0432\u0430\u044f \u0447\u0430\u0441\u0442\u044c \u043a\u0443\u0440\u0441\u0430 \u0437\u0430\u043a\u0440\u044b\u0432\u0430\u0435\u0442 2-4 \u0440\u0430\u0437\u0440\u044f\u0434."),
        u(r"\u0414\u043b\u044f 5-6 \u0440\u0430\u0437\u0440\u044f\u0434\u0430 \u0431\u0443\u0434\u0443\u0442 \u043e\u0442\u0434\u0435\u043b\u044c\u043d\u044b\u0435 \u0443\u0433\u043b\u0443\u0431\u043b\u0435\u043d\u043d\u044b\u0435 \u0431\u043b\u043e\u043a\u0438."),
        u(r"\u0414\u0430\u043b\u044c\u0448\u0435 \u043d\u0430\u0447\u0438\u043d\u0430\u0435\u0442\u0441\u044f \u043f\u0440\u0430\u043a\u0442\u0438\u0447\u0435\u0441\u043a\u0430\u044f \u0442\u0435\u043c\u0430 2."),
    ],
    GREEN,
)
back_s6, next_s6 = add_linear_nav(s8, u(r"\u041d\u0410\u0417\u0410\u0414"), u(r"\u0414\u0410\u041b\u0415\u0415"))


# S007
s9 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s9, PALE)
add_band(
    s9,
    u(r"\u0422\u0435\u043c\u0430 2. \u0421\u0442\u0440\u043e\u043f\u0430\u043b\u044c\u043d\u044b\u0435 \u0440\u0430\u0431\u043e\u0442\u044b"),
    "S007",
    u(r"\u041f\u0435\u0440\u0435\u0445\u043e\u0434 \u043a \u043e\u0441\u043d\u043e\u0432\u043d\u043e\u043c\u0443 \u043f\u0440\u0430\u043a\u0442\u0438\u0447\u0435\u0441\u043a\u043e\u043c\u0443 \u0431\u043b\u043e\u043a\u0443"),
)
add_panel(
    s9,
    Inches(0.9),
    Inches(1.6),
    Inches(5.55),
    Inches(4.45),
    u(r"\u0427\u0442\u043e \u0434\u0430\u043b\u044c\u0448\u0435"),
    [
        u(r"\u0412 \u0422\u0435\u043c\u0435 2 \u043d\u0430\u0447\u0438\u043d\u0430\u0435\u0442\u0441\u044f \u043f\u0440\u0438\u043a\u043b\u0430\u0434\u043d\u0430\u044f \u0447\u0430\u0441\u0442\u044c \u043a\u0443\u0440\u0441\u0430."),
        u(r"\u0417\u0434\u0435\u0441\u044c \u0431\u0443\u0434\u0443\u0442 \u0433\u0440\u0443\u0437\u044b, \u0441\u0445\u0435\u043c\u044b \u0441\u0442\u0440\u043e\u043f\u043e\u0432\u043a\u0438, \u0441\u0438\u0433\u043d\u0430\u043b\u044b \u0438 \u043e\u043f\u0430\u0441\u043d\u044b\u0435 \u0437\u043e\u043d\u044b."),
        u(r"\u042d\u0442\u043e \u0442\u043e\u0447\u043a\u0430 \u043f\u0435\u0440\u0435\u0445\u043e\u0434\u0430 \u0438\u0437 \u0432\u0432\u043e\u0434\u043d\u043e\u0439 \u0447\u0430\u0441\u0442\u0438 \u0432 \u043e\u0441\u043d\u043e\u0432\u043d\u0443\u044e \u0440\u0430\u0431\u043e\u0442\u0443."),
    ],
    BLUE,
)
route_box = s9.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.8), Inches(5.0), Inches(2.8))
route_box.fill.solid()
route_box.fill.fore_color.rgb = WHITE
route_box.line.color.rgb = GREEN
add_text(
    s9,
    Inches(7.2),
    Inches(2.15),
    Inches(4.2),
    Inches(1.55),
    u(r"\u0422\u0435\u043c\u0430 2.1 \u0413\u0440\u0443\u0437\u044b\n\u2192 \u0422\u0435\u043c\u0430 2.2 \u0422\u0435\u0445\u043d\u043e\u043b\u043e\u0433\u0438\u044f\n\u2192 \u0422\u0435\u043c\u0430 2.3 \u041f\u0440\u043e\u0438\u0437\u0432\u043e\u0434\u0441\u0442\u0432\u043e \u0440\u0430\u0431\u043e\u0442"),
    19,
    True,
    TEXT,
    PP_ALIGN.CENTER,
)
back_s7, next_s7 = add_linear_nav(s9, u(r"\u041d\u0410\u0417\u0410\u0414"), None)


start_btn.click_action.target_slide = s2
card_video.click_action.target_slide = s3
card_gloss.click_action.target_slide = s4
btn_video.click_action.target_slide = s3
btn_gloss.click_action.target_slide = s4
back_s2.click_action.target_slide = s1
next_s2.click_action.target_slide = s5
back1.click_action.target_slide = s2
back2.click_action.target_slide = s2
back_s3.click_action.target_slide = s2
main_click_shapes[0].click_action.target_slide = s5_p01
main_click_shapes[1].click_action.target_slide = s5_p02
main_click_shapes[2].click_action.target_slide = s5_p03
main_click_shapes[3].click_action.target_slide = s5_p04
next_s3.click_action.target_slide = s6
p01_back_main.click_action.target_slide = s5
p02_back_main.click_action.target_slide = s5
p03_back_main.click_action.target_slide = s5
p04_back_main.click_action.target_slide = s5
back_s4.click_action.target_slide = s5
next_s4.click_action.target_slide = s7
back_s5.click_action.target_slide = s6
next_s5.click_action.target_slide = s8
back_s6.click_action.target_slide = s7
next_s6.click_action.target_slide = s9
back_s7.click_action.target_slide = s8


prs.save(OUT_FILE)
print(OUT_FILE)
