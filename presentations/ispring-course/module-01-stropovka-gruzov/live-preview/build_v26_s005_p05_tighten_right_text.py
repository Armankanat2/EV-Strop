from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt


SRC = Path(__file__).resolve().parent / "S001-S007_live_preview_working_2026-06-23_v25.pptx"
OUT = Path(__file__).resolve().parent / "S001-S007_live_preview_working_2026-06-23_v26.pptx"


def main() -> None:
    prs = Presentation(SRC)
    slide = prs.slides[15]

    text_box = slide.shapes[8]
    tf = text_box.text_frame

    # Shift the content block slightly upward and keep the same bottom edge.
    shift_up = Inches(0.12)
    text_box.top -= shift_up
    text_box.height += shift_up

    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.margin_left = Inches(0.04)
    tf.margin_right = Inches(0.04)

    for paragraph in tf.paragraphs:
        paragraph.space_after = Pt(3)
        for run in paragraph.runs:
            if run.font.size is not None:
                if run.font.bold:
                    run.font.size = Pt(11.8)
                else:
                    run.font.size = Pt(11.4)

    if tf.paragraphs:
        tf.paragraphs[-1].space_after = Pt(0)

    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    main()
