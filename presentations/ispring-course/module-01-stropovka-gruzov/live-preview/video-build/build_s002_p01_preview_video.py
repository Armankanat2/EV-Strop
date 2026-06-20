from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path
from typing import Iterable

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Inches


ROOT = Path(__file__).resolve().parents[5]
LIVE_DIR = Path(__file__).resolve().parent.parent
BASE_DIR = Path(__file__).resolve().parent / "base-slides"
OUT_DIR = Path(__file__).resolve().parent / "s002-p01-preview-2026-06-20"
FRAMES_DIR = OUT_DIR / "frames"
PPTX_FILE = OUT_DIR / "S002-P01_video_preview_61_14sec_2026-06-20.pptx"
PLAN_FILE = OUT_DIR / "S002-P01_video_preview_plan_2026-06-20.txt"

SLIDE_01 = BASE_DIR / "slide_01.png"
SLIDE_02 = BASE_DIR / "slide_02.png"
SLIDE_03 = BASE_DIR / "slide_03.png"
SLIDE_04 = BASE_DIR / "slide_04.png"

AUDIO_FILE = ROOT / "assets" / "course-media" / "module-01-stropovka-gruzov" / "audio" / "S002-P01_ozvuchka_55_sec_stress_2026-06-20.mp3"

FONT_UI = Path(r"C:\Windows\Fonts\segoeui.ttf")
FONT_BOLD = Path(r"C:\Windows\Fonts\arialbd.ttf")
FONT_REG = Path(r"C:\Windows\Fonts\arial.ttf")

WIDTH = 1280
HEIGHT = 720

INTRO_SECONDS = 1.0
BODY_SECONDS = 59.14
OUTRO_SECONDS = 1.0
FPS = 10
FRAME_SECONDS = 1.0 / FPS

SUB_BG = (14, 27, 41, 220)
SUB_TEXT = (255, 255, 255, 255)
GREEN = (57, 214, 103, 255)
GREEN_SOFT = (57, 214, 103, 70)
BLACK = (0, 0, 0, 255)
WHITE = (255, 255, 255, 255)
RED = (221, 74, 58, 255)
ORANGE = (234, 124, 53, 255)
BLUE = (44, 104, 161, 255)
NAVY = (18, 43, 66, 255)
TEXT = (42, 52, 64, 255)
STEEL = (83, 119, 149, 255)


@dataclass(frozen=True)
class Point:
    x: float
    y: float


@dataclass(frozen=True)
class Rect:
    x: int
    y: int
    w: int
    h: int

    @property
    def center(self) -> Point:
        return Point(self.x + self.w / 2, self.y + self.h / 2)

    def inflate(self, px: int) -> "Rect":
        return Rect(self.x - px, self.y - px, self.w + px * 2, self.h + px * 2)


@dataclass(frozen=True)
class Scene:
    name: str
    start: float
    end: float
    slide: str
    subtitle: str
    cursor_anchor: Point | None = None
    highlight: Rect | None = None
    click_window: tuple[float, float] | None = None
    overlay: str | None = None


START_BTN = Rect(64, 494, 250, 60)
VIDEO_BTN = Rect(168, 480, 341, 60)
GLOSSARY_BTN = Rect(763, 480, 341, 60)
VIDEO_CARD = Rect(67, 158, 547, 374)
GLOSSARY_CARD = Rect(662, 158, 547, 374)
BACK_BTN = Rect(878, 629, 293, 48)
ROW_24 = Rect(65, 188, 350, 40)
ROW_56 = Rect(66, 221, 350, 40)

SCENES: list[Scene] = [
    Scene(
        name="intro-title",
        start=0.00,
        end=3.80,
        slide="slide_01",
        subtitle="Перед вами видеоинструкция по курсу «Стропальщик».",
        cursor_anchor=Point(START_BTN.x + START_BTN.w - 72, START_BTN.y + 10),
        highlight=START_BTN,
        click_window=(0.45, 3.65),
        overlay="title-note",
    ),
    Scene(
        name="interactive-route",
        start=3.80,
        end=9.60,
        slide="slide_02",
        subtitle="Курс интерактивный. Вы будете проходить темы по порядку, открывать пояснения и отвечать на вопросы.",
        overlay="route",
    ),
    Scene(
        name="self-check",
        start=9.60,
        end=14.40,
        slide="slide_02",
        subtitle="После тем и отдельных блоков будут тесты для самопроверки.",
        overlay="test-card",
    ),
    Scene(
        name="open-video",
        start=14.40,
        end=18.80,
        slide="slide_02",
        subtitle="Если на слайде есть активная кнопка, нажмите на нее.",
        cursor_anchor=Point(VIDEO_BTN.x + VIDEO_BTN.w - 112, VIDEO_BTN.y + 12),
        highlight=VIDEO_BTN,
        click_window=(14.85, 18.35),
    ),
    Scene(
        name="video-footer",
        start=18.80,
        end=24.20,
        slide="slide_03",
        subtitle="Так вы откроете термин, схему, картинку или короткое пояснение.",
        highlight=Rect(567, 173, 575, 342),
        overlay="video-points",
    ),
    Scene(
        name="return-from-video",
        start=24.20,
        end=28.80,
        slide="slide_03",
        subtitle="После просмотра нажмите «НАЗАД К КУРСУ» и продолжайте обучение.",
        cursor_anchor=Point(BACK_BTN.x + BACK_BTN.w - 112, BACK_BTN.y - 8),
        highlight=BACK_BTN,
        click_window=(24.70, 28.25),
    ),
    Scene(
        name="open-glossary",
        start=28.80,
        end=33.50,
        slide="slide_02",
        subtitle="Если в промежуточном тесте ответ неверный, откройте пояснение, посмотрите материал и вернитесь к вопросу.",
        cursor_anchor=Point(GLOSSARY_BTN.x + GLOSSARY_BTN.w - 112, GLOSSARY_BTN.y + 12),
        highlight=GLOSSARY_BTN,
        click_window=(29.20, 32.85),
        overlay="hint-chip",
    ),
    Scene(
        name="glossary-screen",
        start=33.50,
        end=37.20,
        slide="slide_04",
        subtitle="Если в промежуточном тесте ответ неверный, откройте пояснение, посмотрите материал и вернитесь к вопросу.",
        overlay="glossary-focus",
    ),
    Scene(
        name="final-test",
        start=37.20,
        end=44.40,
        slide="slide_02",
        subtitle="В конце курса будет итоговый тест: 20 вопросов и не больше 3 ошибок.",
        overlay="final-test-card",
    ),
    Scene(
        name="base-course",
        start=44.40,
        end=51.80,
        slide="slide_01",
        subtitle="Базовая часть курса предназначена для стро́пальщика 2-4 разряда.",
        highlight=ROW_24,
        overlay="base-chip",
    ),
    Scene(
        name="extra-course",
        start=51.80,
        end=59.14,
        slide="slide_01",
        subtitle="Если вы проходите подготовку на 5-6 разряд, открывайте дополнительные блоки.",
        highlight=ROW_56,
        overlay="extra-chip",
    ),
]


def get_font(size: int, bold: bool = False) -> ImageFont.FreeTypeFont:
    path = FONT_BOLD if bold and FONT_BOLD.exists() else FONT_UI if FONT_UI.exists() else FONT_REG
    return ImageFont.truetype(str(path), size)


def wrap_text(draw: ImageDraw.ImageDraw, text: str, font: ImageFont.FreeTypeFont, max_width: int) -> list[str]:
    words = text.split()
    lines: list[str] = []
    current = ""
    for word in words:
        probe = word if not current else f"{current} {word}"
        box = draw.textbbox((0, 0), probe, font=font)
        if box[2] - box[0] <= max_width:
            current = probe
        else:
            if current:
                lines.append(current)
            current = word
    if current:
        lines.append(current)
    return lines


def draw_multiline_centered(
    draw: ImageDraw.ImageDraw,
    box: Rect,
    text: str,
    font: ImageFont.FreeTypeFont,
    fill: tuple[int, int, int, int],
    spacing: int = 6,
) -> None:
    lines = wrap_text(draw, text, font, box.w - 24)
    heights = []
    widths = []
    for line in lines:
        left, top, right, bottom = draw.textbbox((0, 0), line, font=font)
        widths.append(right - left)
        heights.append(bottom - top)
    total_h = sum(heights) + spacing * max(0, len(lines) - 1)
    y = box.y + (box.h - total_h) / 2
    for line, width, height in zip(lines, widths, heights):
        x = box.x + (box.w - width) / 2
        draw.text((x, y), line, font=font, fill=fill)
        y += height + spacing


def draw_subtitle(draw: ImageDraw.ImageDraw, text: str) -> None:
    box = Rect(55, 610, 1170, 82)
    draw.rounded_rectangle((box.x, box.y, box.x + box.w, box.y + box.h), radius=22, fill=SUB_BG)
    font = get_font(28, bold=True)
    draw_multiline_centered(draw, box, text, font, SUB_TEXT, spacing=4)


def draw_cursor(base: Image.Image, pos: Point) -> None:
    draw = ImageDraw.Draw(base, "RGBA")
    x = int(pos.x)
    y = int(pos.y)
    body = [
        (x, y),
        (x + 16, y + 40),
        (x + 22, y + 28),
        (x + 38, y + 46),
        (x + 46, y + 40),
        (x + 28, y + 22),
        (x + 40, y + 16),
    ]
    shadow = [(px + 3, py + 3) for px, py in body]
    draw.polygon(shadow, fill=(0, 0, 0, 70))
    draw.polygon(body, fill=WHITE, outline=BLACK)


def draw_highlight(draw: ImageDraw.ImageDraw, rect: Rect, pulse: bool) -> None:
    if not pulse:
        return
    outer = rect.inflate(10)
    draw.rounded_rectangle((outer.x, outer.y, outer.x + outer.w, outer.y + outer.h), radius=20, outline=GREEN, width=6)
    fill = Image.new("RGBA", (WIDTH, HEIGHT), (0, 0, 0, 0))
    fill_draw = ImageDraw.Draw(fill, "RGBA")
    fill_draw.rounded_rectangle(
        (rect.x, rect.y, rect.x + rect.w, rect.y + rect.h),
        radius=16,
        fill=GREEN_SOFT,
    )
    draw._image.alpha_composite(fill)


def draw_chip(draw: ImageDraw.ImageDraw, rect: Rect, text: str, fill: tuple[int, int, int, int]) -> None:
    draw.rounded_rectangle((rect.x, rect.y, rect.x + rect.w, rect.y + rect.h), radius=18, fill=fill)
    font = get_font(22, bold=True)
    draw_multiline_centered(draw, rect, text, font, WHITE, spacing=2)


def overlay_title_note(draw: ImageDraw.ImageDraw) -> None:
    box = Rect(80, 355, 560, 92)
    draw.rounded_rectangle((box.x, box.y, box.x + box.w, box.y + box.h), radius=20, fill=(255, 255, 255, 235), outline=(210, 218, 226, 255), width=2)
    draw_multiline_centered(draw, box, "Базовая часть курса для 2-4 разряда.\nДля 5-6 разряда есть дополнительные блоки.", get_font(23, bold=False), TEXT, spacing=4)


def overlay_route(draw: ImageDraw.ImageDraw, local_t: float) -> None:
    band = Rect(110, 548, 1040, 46)
    draw.rounded_rectangle((band.x, band.y, band.x + band.w, band.y + band.h), radius=22, fill=(255, 255, 255, 235))
    labels = ["Тема 1", "Тема 2", "Тема 3", "Тема 4", "Тесты"]
    dots = [160, 380, 600, 820, 1040]
    active = min(int(local_t * 5), 4)
    font = get_font(18, bold=True)
    for idx, (label, x) in enumerate(zip(labels, dots)):
        color = ORANGE if idx <= active else STEEL
        draw.line((x + 40, 571, x + 140, 571), fill=(170, 182, 194, 255), width=4)
        draw.ellipse((x, 553, x + 36, 589), fill=color, outline=color)
        draw.text((x + 48, 553), label, font=font, fill=TEXT)


def overlay_test_card(draw: ImageDraw.ImageDraw, local_t: float) -> None:
    box = Rect(388, 210, 504, 220)
    draw.rounded_rectangle((box.x, box.y, box.x + box.w, box.y + box.h), radius=26, fill=(255, 255, 255, 245), outline=(44, 104, 161, 255), width=3)
    title_font = get_font(24, bold=True)
    body_font = get_font(20, bold=False)
    draw.text((440, 240), "Самопроверка", font=title_font, fill=TEXT)
    draw.text((430, 285), "Ответ: неверно", font=body_font, fill=RED)
    draw.text((430, 320), "Откройте пояснение", font=body_font, fill=TEXT)
    btn = Rect(460, 365, 360, 48)
    pulse = int(local_t * 8) % 2 == 0
    fill = GREEN if pulse else BLUE
    draw_chip(draw, btn, "ПОКАЗАТЬ ПОЯСНЕНИЕ", fill)


def overlay_hint_chip(draw: ImageDraw.ImageDraw) -> None:
    draw_chip(draw, Rect(485, 112, 310, 44), "Открыть пояснение", ORANGE)


def overlay_video_points(draw: ImageDraw.ImageDraw, local_t: float) -> None:
    lines = [
        "Темы идут по порядку.",
        "После тем будут тесты.",
        "По кнопкам открывайте пояснения.",
    ]
    box = Rect(675, 185, 365, 225)
    draw.rounded_rectangle((box.x, box.y, box.x + box.w, box.y + box.h), radius=18, fill=(255, 255, 255, 220))
    font = get_font(22, bold=False)
    for idx, line in enumerate(lines):
        y = 215 + idx * 56
        color = ORANGE if idx == min(int(local_t * 3), 2) else TEXT
        draw.text((705, y), f"• {line}", font=font, fill=color)


def overlay_glossary_focus(draw: ImageDraw.ImageDraw, local_t: float) -> None:
    focus_left = local_t < 0.5
    rect = Rect(82, 152, 520, 116) if focus_left else Rect(675, 152, 500, 116)
    draw_highlight(draw, rect, True)


def overlay_final_test_card(draw: ImageDraw.ImageDraw, local_t: float) -> None:
    box = Rect(365, 182, 550, 250)
    draw.rounded_rectangle((box.x, box.y, box.x + box.w, box.y + box.h), radius=28, fill=(255, 255, 255, 246), outline=(18, 43, 66, 255), width=3)
    draw_multiline_centered(
        draw,
        Rect(410, 200, 460, 70),
        "Итоговый тест",
        get_font(30, bold=True),
        NAVY,
        spacing=0,
    )
    draw_multiline_centered(
        draw,
        Rect(410, 278, 460, 58),
        "20 вопросов",
        get_font(28, bold=True),
        ORANGE,
        spacing=0,
    )
    draw_multiline_centered(
        draw,
        Rect(390, 337, 500, 64),
        "Не больше 3 ошибок",
        get_font(25, bold=False),
        TEXT,
        spacing=0,
    )
    pill = Rect(515, 386, 250, 36)
    fill = GREEN if int(local_t * 8) % 2 == 0 else BLUE
    draw_chip(draw, pill, "ДОПУСК К ИТОГУ", fill)


def overlay_base_chip(draw: ImageDraw.ImageDraw) -> None:
    draw_chip(draw, Rect(85, 256, 300, 44), "Базовая часть курса", BLUE)


def overlay_extra_chip(draw: ImageDraw.ImageDraw) -> None:
    draw_chip(draw, Rect(85, 256, 338, 44), "Дополнительные блоки 5-6", ORANGE)


def base_image(name: str) -> Image.Image:
    mapping = {
        "slide_01": SLIDE_01,
        "slide_02": SLIDE_02,
        "slide_03": SLIDE_03,
        "slide_04": SLIDE_04,
    }
    return Image.open(mapping[name]).convert("RGBA")


def apply_overlay(img: Image.Image, scene: Scene, local_t: float) -> None:
    draw = ImageDraw.Draw(img, "RGBA")
    if scene.overlay == "title-note":
        overlay_title_note(draw)
    elif scene.overlay == "route":
        overlay_route(draw, local_t)
    elif scene.overlay == "test-card":
        overlay_test_card(draw, local_t)
    elif scene.overlay == "video-points":
        overlay_video_points(draw, local_t)
    elif scene.overlay == "hint-chip":
        overlay_hint_chip(draw)
    elif scene.overlay == "glossary-focus":
        overlay_glossary_focus(draw, local_t)
    elif scene.overlay == "final-test-card":
        overlay_final_test_card(draw, local_t)
    elif scene.overlay == "base-chip":
        overlay_base_chip(draw)
    elif scene.overlay == "extra-chip":
        overlay_extra_chip(draw)


def scene_at(t: float) -> Scene:
    for scene in SCENES:
        if scene.start <= t < scene.end:
            return scene
    return SCENES[-1]


def render_body_frames() -> tuple[list[Path], list[float]]:
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    FRAMES_DIR.mkdir(parents=True, exist_ok=True)

    frames: list[Path] = []
    durations: list[float] = []

    frame_index = 0
    current = 0.0
    while current < BODY_SECONDS - 1e-9:
        duration = min(FRAME_SECONDS, BODY_SECONDS - current)
        t_mid = current + duration / 2
        scene = scene_at(t_mid)
        img = base_image(scene.slide)
        local_t = (t_mid - scene.start) / max(scene.end - scene.start, 0.001)
        local_t = max(0.0, min(1.0, local_t))
        apply_overlay(img, scene, local_t)
        draw = ImageDraw.Draw(img, "RGBA")
        if scene.highlight and scene.click_window:
            pulse = scene.click_window[0] <= t_mid <= scene.click_window[1]
            draw_highlight(draw, scene.highlight, pulse)
        elif scene.highlight and scene.click_window is None:
            pulse = int(local_t * 6) % 2 == 0
            draw_highlight(draw, scene.highlight, pulse)
        if scene.cursor_anchor and scene.click_window and scene.click_window[0] <= t_mid <= scene.click_window[1]:
            draw_cursor(img, scene.cursor_anchor)
        draw_subtitle(draw, scene.subtitle)

        frame_path = FRAMES_DIR / f"frame_{frame_index:04d}.png"
        img.save(frame_path)
        frames.append(frame_path)
        durations.append(duration)
        frame_index += 1
        current += duration

    return frames, durations


def make_intro_outro() -> tuple[Path, Path]:
    intro = OUT_DIR / "intro_frame.png"
    outro = OUT_DIR / "outro_frame.png"
    base_image("slide_01").save(intro)

    outro_img = base_image("slide_01")
    draw = ImageDraw.Draw(outro_img, "RGBA")
    overlay_extra_chip(draw)
    draw_highlight(draw, ROW_56, True)
    outro_img.save(outro)
    return intro, outro


def build_pptx(frame_paths: Iterable[Path], frame_durations: list[float], intro: Path, outro: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    all_images = [intro, *frame_paths, outro]
    for image_path in all_images:
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(str(image_path), 0, 0, width=prs.slide_width, height=prs.slide_height)

    prs.save(PPTX_FILE)

    lines = [
        f"audio={AUDIO_FILE}",
        f"pptx={PPTX_FILE}",
        f"intro_seconds={INTRO_SECONDS}",
        f"body_seconds={BODY_SECONDS}",
        f"outro_seconds={OUTRO_SECONDS}",
        f"fps={FPS}",
        f"frame_count={len(frame_durations)}",
        "",
        "slide_timings_seconds:",
        f"slide_0001={INTRO_SECONDS:.2f}",
    ]
    for idx, duration in enumerate(frame_durations, start=2):
        lines.append(f"slide_{idx:04d}={duration:.2f}")
    lines.append(f"slide_{len(frame_durations) + 2:04d}={OUTRO_SECONDS:.2f}")
    PLAN_FILE.write_text("\n".join(lines) + "\n", encoding="utf-8")


def main() -> None:
    frames, durations = render_body_frames()
    intro, outro = make_intro_outro()
    build_pptx(frames, durations, intro, outro)
    print(f"Built {len(frames)} body frames")
    print(PPTX_FILE)
    print(PLAN_FILE)


if __name__ == "__main__":
    main()
