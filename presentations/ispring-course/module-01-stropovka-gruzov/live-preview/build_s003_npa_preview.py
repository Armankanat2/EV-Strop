from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


OUT_DIR = Path(__file__).resolve().parent
OUT_FILE = OUT_DIR / "S003_npa_preview_2026-06-22_v4.pptx"


def u(text):
    return text.encode("ascii").decode("unicode_escape")


NAVY = RGBColor(18, 43, 66)
BLUE = RGBColor(38, 90, 137)
GREEN = RGBColor(66, 133, 95)
STEEL = RGBColor(83, 119, 149)
PALE = RGBColor(240, 244, 247)
ORANGE = RGBColor(222, 116, 54)
WHITE = RGBColor(255, 255, 255)
TEXT = RGBColor(33, 43, 54)
GRAY = RGBColor(110, 122, 134)
LINE = RGBColor(210, 218, 226)


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def add_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, x, y, w, h, text, size=18, bold=False, color=TEXT, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return box


def add_band(slide, title, code, subtitle=""):
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.72))
    band.fill.solid()
    band.fill.fore_color.rgb = NAVY
    band.line.fill.background()

    add_text(slide, Inches(0.45), Inches(0.14), Inches(8.8), Inches(0.3), title, 24, True, WHITE)

    tag = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(11.35), Inches(0.14), Inches(1.45), Inches(0.35)
    )
    tag.fill.solid()
    tag.fill.fore_color.rgb = ORANGE
    tag.line.fill.background()
    tf = tag.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = code
    r.font.size = Pt(16)
    r.font.bold = True
    r.font.color.rgb = WHITE

    if subtitle:
        add_text(slide, Inches(0.5), Inches(0.82), Inches(6.8), Inches(0.25), subtitle, 11, False, GRAY)


def add_button(slide, x, y, w, h, text, fill, font_size=18):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = fill
    tf = shape.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.size = Pt(font_size)
    r.font.bold = True
    r.font.color.rgb = WHITE
    return shape


def add_header_cell(slide, x, y, w, h, text, fill, font_size=15):
    cell = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, y, w, h)
    cell.fill.solid()
    cell.fill.fore_color.rgb = fill
    cell.line.color.rgb = fill
    tf = cell.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.size = Pt(font_size)
    r.font.bold = True
    r.font.color.rgb = WHITE
    return cell


def add_body_cell(slide, x, y, w, h, text, align=PP_ALIGN.LEFT, fill=WHITE, size=14, bold=False):
    cell = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, y, w, h)
    cell.fill.solid()
    cell.fill.fore_color.rgb = fill
    cell.line.color.rgb = LINE
    tf = cell.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = TEXT
    return cell


def add_info_panel(slide, x, y, w, h, title, lines, accent=GREEN):
    panel = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
    panel.fill.solid()
    panel.fill.fore_color.rgb = WHITE
    panel.line.color.rgb = accent
    panel.line.width = Pt(1.5)

    head = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, Inches(0.62))
    head.fill.solid()
    head.fill.fore_color.rgb = accent
    head.line.color.rgb = accent
    tf = head.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = title
    r.font.size = Pt(19)
    r.font.bold = True
    r.font.color.rgb = WHITE

    body = slide.shapes.add_textbox(x + Inches(0.25), y + Inches(0.86), w - Inches(0.5), h - Inches(1.06))
    tf = body.text_frame
    tf.clear()
    tf.word_wrap = True
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = line
        p.bullet = True
        p.space_after = Pt(8)
        if p.runs:
            p.runs[0].font.size = Pt(15)
            p.runs[0].font.color.rgb = TEXT
    return panel


def add_detail_table(slide, title, code, doc_name, about, importance, accent):
    add_bg(slide, WHITE)
    add_band(slide, title, code, u(r"\u041f\u043e\u0434\u0432\u0430\u043b \u043e\u0442 S003"))

    add_text(
        slide,
        Inches(0.7),
        Inches(1.1),
        Inches(5.6),
        Inches(0.42),
        u(r"\u041a\u043e\u0440\u043e\u0442\u043a\u0430\u044f \u0441\u043f\u0440\u0430\u0432\u043a\u0430 \u043f\u043e \u0434\u043e\u043a\u0443\u043c\u0435\u043d\u0442\u0443"),
        22,
        True,
        NAVY,
    )

    table_x = Inches(0.7)
    table_y = Inches(1.75)
    left_w = Inches(2.6)
    right_w = Inches(7.35)
    row_h = Inches(1.1)

    add_header_cell(slide, table_x, table_y, left_w, Inches(0.55), u(r"\u041f\u043e\u043b\u0435"), accent)
    add_header_cell(slide, table_x + left_w, table_y, right_w, Inches(0.55), u(r"\u0421\u043e\u0434\u0435\u0440\u0436\u0430\u043d\u0438\u0435"), accent)

    labels = [
        u(r"\u0414\u043e\u043a\u0443\u043c\u0435\u043d\u0442"),
        u(r"\u041e \u0447\u0435\u043c \u0434\u043e\u043a\u0443\u043c\u0435\u043d\u0442"),
        u(r"\u0427\u0435\u043c \u0432\u0430\u0436\u043d\u043e \u0434\u043b\u044f \u0441\u0442\u0440\u043e\u043f\u0430\u043b\u044c\u0449\u0438\u043a\u0430"),
    ]
    values = [doc_name, about, importance]

    current_y = table_y + Inches(0.55)
    heights = [Inches(0.8), Inches(1.0), Inches(1.75)]
    for label, value, height in zip(labels, values, heights):
        add_body_cell(slide, table_x, current_y, left_w, height, label, PP_ALIGN.LEFT, PALE, 14, True)
        add_body_cell(slide, table_x + left_w, current_y, right_w, height, value, PP_ALIGN.LEFT, WHITE, 14, False)
        current_y += height



# S003 main slide
s1 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s1, PALE)
add_band(
    s1,
    u(r"\u041a\u0430\u043a\u0438\u0435 \u0434\u043e\u043a\u0443\u043c\u0435\u043d\u0442\u044b \u0440\u0435\u0433\u0443\u043b\u0438\u0440\u0443\u044e\u0442 \u0440\u0430\u0431\u043e\u0442\u0443 \u0441\u0442\u0440\u043e\u043f\u0430\u043b\u044c\u0449\u0438\u043a\u0430"),
    "S003",
    u(r"\u0422\u0435\u043c\u0430 1 \u2022 \u043e\u0431\u0449\u0430\u044f \u0447\u0430\u0441\u0442\u044c"),
)
add_text(
    s1,
    Inches(0.7),
    Inches(1.0),
    Inches(7.8),
    Inches(0.45),
    u(r"\u042d\u0442\u043e \u043e\u0441\u043d\u043e\u0432\u043d\u044b\u0435 \u043f\u0440\u0430\u0432\u043e\u0432\u044b\u0435 \u0430\u043a\u0442\u044b, \u0440\u0435\u0433\u0443\u043b\u0438\u0440\u0443\u044e\u0449\u0438\u0435 \u0440\u0430\u0431\u043e\u0442\u0443 \u0441\u0442\u0440\u043e\u043f\u0430\u043b\u044c\u0449\u0438\u043a\u0430"),
    22,
    True,
    NAVY,
)

table_x = Inches(0.7)
table_y = Inches(1.85)
law_w = Inches(6.15)
date_w = Inches(1.65)
issuer_w = Inches(3.35)
row_h = Inches(0.82)

add_header_cell(s1, table_x, table_y, law_w, Inches(0.56), u(r"\u0417\u0430\u043a\u043e\u043d"), BLUE)
add_header_cell(s1, table_x + law_w, table_y, date_w, Inches(0.56), u(r"\u0414\u0430\u0442\u0430 \u043f\u0440\u0438\u043d\u044f\u0442\u0438\u044f"), BLUE, 13)
add_header_cell(s1, table_x + law_w + date_w, table_y, issuer_w, Inches(0.56), u(r"\u041a\u0442\u043e \u0438\u0437\u0434\u0430\u043b"), BLUE)

rows = [
    (
        u(r"\u0422\u0440\u0443\u0434\u043e\u0432\u043e\u0439 \u043a\u043e\u0434\u0435\u043a\u0441 \u0420\u0424 \u2116 197-\u0424\u0417"),
        "30.12.2001",
        u(r"\u0424\u0435\u0434\u0435\u0440\u0430\u043b\u044c\u043d\u043e\u0435 \u0421\u043e\u0431\u0440\u0430\u043d\u0438\u0435 \u0420\u0424,\n\u041f\u0440\u0435\u0437\u0438\u0434\u0435\u043d\u0442 \u0420\u0424"),
    ),
    (
        u(r"\u041f\u043e\u0441\u0442\u0430\u043d\u043e\u0432\u043b\u0435\u043d\u0438\u0435 \u041f\u0440\u0430\u0432\u0438\u0442\u0435\u043b\u044c\u0441\u0442\u0432\u0430 \u0420\u0424 \u2116 2464"),
        "24.12.2021",
        u(r"\u041f\u0440\u0430\u0432\u0438\u0442\u0435\u043b\u044c\u0441\u0442\u0432\u043e \u0420\u0424"),
    ),
    (
        u(r"\u0424\u0435\u0434\u0435\u0440\u0430\u043b\u044c\u043d\u044b\u0439 \u0437\u0430\u043a\u043e\u043d \u2116 116-\u0424\u0417"),
        "21.07.1997",
        u(r"\u0424\u0435\u0434\u0435\u0440\u0430\u043b\u044c\u043d\u043e\u0435 \u0421\u043e\u0431\u0440\u0430\u043d\u0438\u0435 \u0420\u0424,\n\u041f\u0440\u0435\u0437\u0438\u0434\u0435\u043d\u0442 \u0420\u0424"),
    ),
    (
        u(r"\u041f\u0440\u0438\u043a\u0430\u0437 \u0420\u043e\u0441\u0442\u0435\u0445\u043d\u0430\u0434\u0437\u043e\u0440\u0430 \u2116 461"),
        "26.11.2020",
        u(r"\u0420\u043e\u0441\u0442\u0435\u0445\u043d\u0430\u0434\u0437\u043e\u0440"),
    ),
]

main_click_shapes = []
current_y = table_y + Inches(0.56)
for idx, row in enumerate(rows):
    fill = WHITE if idx % 2 == 0 else RGBColor(248, 250, 252)
    law_cell = add_body_cell(s1, table_x, current_y, law_w, row_h, row[0], PP_ALIGN.LEFT, fill, 15, True)
    add_body_cell(s1, table_x + law_w, current_y, date_w, row_h, row[1], PP_ALIGN.CENTER, fill, 14, False)
    add_body_cell(s1, table_x + law_w + date_w, current_y, issuer_w, row_h, row[2], PP_ALIGN.LEFT, fill, 13, False)
    main_click_shapes.append(law_cell)
    current_y += row_h

add_info_panel(
    s1,
    Inches(10.45),
    Inches(1.85),
    Inches(2.15),
    Inches(3.45),
    u(r"\u0418\u043d\u0442\u0435\u0440\u0430\u043a\u0442\u0438\u0432"),
    [
        u(r"\u041d\u0430\u0436\u043c\u0438\u0442\u0435 \u043d\u0430 \u043d\u0430\u0437\u0432\u0430\u043d\u0438\u0435 \u0434\u043e\u043a\u0443\u043c\u0435\u043d\u0442\u0430."),
        u(r"\u041e\u0442\u043a\u0440\u043e\u0435\u0442\u0441\u044f \u043a\u043e\u0440\u043e\u0442\u043a\u0438\u0439 \u043f\u043e\u0434\u0432\u0430\u043b-\u0441\u043f\u0440\u0430\u0432\u043a\u0430."),
        u(r"\u0418\u0437 \u043f\u043e\u0434\u0432\u0430\u043b\u0430 \u0435\u0441\u0442\u044c \u0432\u043e\u0437\u0432\u0440\u0430\u0442 \u0432 S003."),
    ],
    GREEN,
)
back_btn = add_button(s1, Inches(0.7), Inches(6.55), Inches(2.35), Inches(0.5), u(r"\u041d\u0410\u0417\u0410\u0414"), STEEL)
next_btn = add_button(s1, Inches(10.25), Inches(6.55), Inches(2.35), Inches(0.5), u(r"\u0414\u0410\u041b\u0415\u0415"), ORANGE)


# Detail slides
s2 = prs.slides.add_slide(prs.slide_layouts[6])
add_detail_table(
    s2,
    u(r"\u0422\u0440\u0443\u0434\u043e\u0432\u043e\u0439 \u043a\u043e\u0434\u0435\u043a\u0441 \u0420\u0424 \u2116 197-\u0424\u0417"),
    "S003-P01",
    u(r"\u0422\u0440\u0443\u0434\u043e\u0432\u043e\u0439 \u043a\u043e\u0434\u0435\u043a\u0441 \u0420\u0424 \u2116 197-\u0424\u0417"),
    u(r"\u041e\u0431\u0449\u0430\u044f \u0440\u0430\u043c\u043a\u0430 \u043f\u043e \u043e\u0445\u0440\u0430\u043d\u0435 \u0442\u0440\u0443\u0434\u0430 \u0438 \u0434\u043e\u043f\u0443\u0441\u043a\u0443 \u043a \u0440\u0430\u0431\u043e\u0442\u0435"),
    u(r"\u0420\u0430\u0437\u0434\u0435\u043b X \"\u041e\u0445\u0440\u0430\u043d\u0430 \u0442\u0440\u0443\u0434\u0430\": \u043e\u0431\u044f\u0437\u0430\u043d\u043d\u043e\u0441\u0442\u0438 \u0440\u0430\u0431\u043e\u0442\u043e\u0434\u0430\u0442\u0435\u043b\u044f \u0438 \u0440\u0430\u0431\u043e\u0442\u043d\u0438\u043a\u0430, \u0421\u0418\u0417, \u0431\u0435\u0437\u043e\u043f\u0430\u0441\u043d\u044b\u0435 \u0443\u0441\u043b\u043e\u0432\u0438\u044f \u0442\u0440\u0443\u0434\u0430, \u0440\u0430\u0441\u0441\u043b\u0435\u0434\u043e\u0432\u0430\u043d\u0438\u0435 \u043d\u0435\u0441\u0447\u0430\u0441\u0442\u043d\u044b\u0445 \u0441\u043b\u0443\u0447\u0430\u0435\u0432"),
    BLUE,
)

s3 = prs.slides.add_slide(prs.slide_layouts[6])
add_detail_table(
    s3,
    u(r"\u041f\u043e\u0441\u0442\u0430\u043d\u043e\u0432\u043b\u0435\u043d\u0438\u0435 \u041f\u0440\u0430\u0432\u0438\u0442\u0435\u043b\u044c\u0441\u0442\u0432\u0430 \u0420\u0424 \u2116 2464"),
    "S003-P02",
    u(r"\u041f\u043e\u0441\u0442\u0430\u043d\u043e\u0432\u043b\u0435\u043d\u0438\u0435 \u041f\u0440\u0430\u0432\u0438\u0442\u0435\u043b\u044c\u0441\u0442\u0432\u0430 \u0420\u0424 \u2116 2464"),
    u(r"\u041e\u0431\u0443\u0447\u0435\u043d\u0438\u0435 \u043f\u043e \u043e\u0445\u0440\u0430\u043d\u0435 \u0442\u0440\u0443\u0434\u0430 \u0438 \u043f\u0440\u043e\u0432\u0435\u0440\u043a\u0430 \u0437\u043d\u0430\u043d\u0438\u0439"),
    u(r"\u0418\u043d\u0441\u0442\u0440\u0443\u043a\u0442\u0430\u0436\u0438, \u0441\u0442\u0430\u0436\u0438\u0440\u043e\u0432\u043a\u0430, \u043e\u0431\u0443\u0447\u0435\u043d\u0438\u0435, \u043f\u0435\u0440\u0432\u0430\u044f \u043f\u043e\u043c\u043e\u0449\u044c, \u043f\u0440\u043e\u0432\u0435\u0440\u043a\u0430 \u0437\u043d\u0430\u043d\u0438\u0439"),
    GREEN,
)

s4 = prs.slides.add_slide(prs.slide_layouts[6])
add_detail_table(
    s4,
    u(r"\u0424\u0435\u0434\u0435\u0440\u0430\u043b\u044c\u043d\u044b\u0439 \u0437\u0430\u043a\u043e\u043d \u2116 116-\u0424\u0417"),
    "S003-P03",
    u(r"\u0424\u0435\u0434\u0435\u0440\u0430\u043b\u044c\u043d\u044b\u0439 \u0437\u0430\u043a\u043e\u043d \u2116 116-\u0424\u0417"),
    u(r"\u0423\u0441\u0442\u0430\u043d\u0430\u0432\u043b\u0438\u0432\u0430\u0435\u0442 \u0442\u0440\u0435\u0431\u043e\u0432\u0430\u043d\u0438\u044f \u0431\u0435\u0437\u043e\u043f\u0430\u0441\u043d\u043e\u0439 \u0440\u0430\u0431\u043e\u0442\u044b \u043d\u0430 \u043e\u043f\u0430\u0441\u043d\u044b\u0445 \u043f\u0440\u043e\u0438\u0437\u0432\u043e\u0434\u0441\u0442\u0432\u0435\u043d\u043d\u044b\u0445 \u043e\u0431\u044a\u0435\u043a\u0442\u0430\u0445, \u043f\u043e\u0434\u0433\u043e\u0442\u043e\u0432\u043a\u0438 \u0440\u0430\u0431\u043e\u0442\u043d\u0438\u043a\u043e\u0432, \u043f\u0440\u0435\u0434\u0443\u043f\u0440\u0435\u0436\u0434\u0435\u043d\u0438\u044f \u0430\u0432\u0430\u0440\u0438\u0439 \u0438 \u043f\u0440\u043e\u0438\u0437\u0432\u043e\u0434\u0441\u0442\u0432\u0435\u043d\u043d\u043e\u0433\u043e \u043a\u043e\u043d\u0442\u0440\u043e\u043b\u044f"),
    u(r"\u041d\u0443\u0436\u0435\u043d \u043f\u0440\u0438 \u0440\u0430\u0431\u043e\u0442\u0435 \u043d\u0430 \u041e\u041f\u041e: \u043f\u043e\u0434\u0433\u043e\u0442\u043e\u0432\u043a\u0430 \u0440\u0430\u0431\u043e\u0442\u043d\u0438\u043a\u043e\u0432, \u0442\u0440\u0435\u0431\u043e\u0432\u0430\u043d\u0438\u044f \u043f\u0440\u043e\u043c\u0431\u0435\u0437\u043e\u043f\u0430\u0441\u043d\u043e\u0441\u0442\u0438, \u043f\u0440\u043e\u0438\u0437\u0432\u043e\u0434\u0441\u0442\u0432\u0435\u043d\u043d\u044b\u0439 \u043a\u043e\u043d\u0442\u0440\u043e\u043b\u044c. \u0414\u043b\u044f \u0441\u0442\u0440\u043e\u043f\u0430\u043b\u044c\u0449\u0438\u043a\u0430 \u0432\u0430\u0436\u0435\u043d \u0442\u0435\u043c, \u0447\u0442\u043e \u0442\u0440\u0435\u0431\u0443\u0435\u0442 \u043f\u0440\u043e\u0445\u043e\u0434\u0438\u0442\u044c \u043e\u0431\u0443\u0447\u0435\u043d\u0438\u0435, \u0441\u043e\u0431\u043b\u044e\u0434\u0430\u0442\u044c \u043f\u0440\u0430\u0432\u0438\u043b\u0430 \u0431\u0435\u0437\u043e\u043f\u0430\u0441\u043d\u043e\u0441\u0442\u0438, \u0432\u044b\u043f\u043e\u043b\u043d\u044f\u0442\u044c \u0438\u043d\u0441\u0442\u0440\u0443\u043a\u0446\u0438\u0438 \u0438 \u0441\u043e\u043e\u0431\u0449\u0430\u0442\u044c \u043e \u043d\u0435\u0438\u0441\u043f\u0440\u0430\u0432\u043d\u043e\u0441\u0442\u044f\u0445 \u0438 \u043e\u043f\u0430\u0441\u043d\u044b\u0445 \u0441\u0438\u0442\u0443\u0430\u0446\u0438\u044f\u0445"),
    STEEL,
)

s5 = prs.slides.add_slide(prs.slide_layouts[6])
add_detail_table(
    s5,
    u(r"\u041f\u0440\u0438\u043a\u0430\u0437 \u0420\u043e\u0441\u0442\u0435\u0445\u043d\u0430\u0434\u0437\u043e\u0440\u0430 \u2116 461"),
    "S003-P04",
    u(r"\u041f\u0440\u0438\u043a\u0430\u0437 \u0420\u043e\u0441\u0442\u0435\u0445\u043d\u0430\u0434\u0437\u043e\u0440\u0430 \u2116 461"),
    u(r"\u0423\u0441\u0442\u0430\u043d\u0430\u0432\u043b\u0438\u0432\u0430\u0435\u0442 \u043f\u0440\u0430\u0432\u0438\u043b\u0430 \u0431\u0435\u0437\u043e\u043f\u0430\u0441\u043d\u043e\u0439 \u044d\u043a\u0441\u043f\u043b\u0443\u0430\u0442\u0430\u0446\u0438\u0438 \u043a\u0440\u0430\u043d\u043e\u0432 \u0438 \u0434\u0440\u0443\u0433\u0438\u0445 \u043f\u043e\u0434\u044a\u0435\u043c\u043d\u044b\u0445 \u0441\u043e\u043e\u0440\u0443\u0436\u0435\u043d\u0438\u0439 \u043d\u0430 \u043e\u043f\u0430\u0441\u043d\u044b\u0445 \u043f\u0440\u043e\u0438\u0437\u0432\u043e\u0434\u0441\u0442\u0432\u0435\u043d\u043d\u044b\u0445 \u043e\u0431\u044a\u0435\u043a\u0442\u0430\u0445"),
    u(r"\u043f. 19, \u043f. 98-99, \u043f. 155-163: \u043f\u0435\u0440\u0441\u043e\u043d\u0430\u043b, \u0441\u0438\u0433\u043d\u0430\u043b\u044b, \u041f\u041f\u0420/\u0422\u041a, \u0441\u0442\u0440\u043e\u043f\u043e\u0432\u043a\u0430, \u041f\u0420\u0420, \u0441\u043a\u043b\u0430\u0434\u0438\u0440\u043e\u0432\u0430\u043d\u0438\u0435. \u0414\u043b\u044f \u0441\u0442\u0440\u043e\u043f\u0430\u043b\u044c\u0449\u0438\u043a\u0430 \u0432\u0430\u0436\u0435\u043d \u0442\u0435\u043c, \u0447\u0442\u043e \u043e\u043f\u0440\u0435\u0434\u0435\u043b\u044f\u0435\u0442 \u043f\u0440\u0430\u0432\u0438\u043b\u0430 \u0441\u0442\u0440\u043e\u043f\u043e\u0432\u043a\u0438 \u0438 \u043f\u0435\u0440\u0435\u043c\u0435\u0449\u0435\u043d\u0438\u044f \u0433\u0440\u0443\u0437\u043e\u0432, \u043f\u0440\u043e\u0432\u0435\u0440\u043a\u0438 \u0441\u0442\u0440\u043e\u043f\u043e\u0432 \u0438 \u0442\u0430\u0440\u044b, \u043f\u043e\u0434\u0430\u0447\u0438 \u0441\u0438\u0433\u043d\u0430\u043b\u043e\u0432 \u0438 \u043f\u0440\u0435\u043a\u0440\u0430\u0449\u0435\u043d\u0438\u044f \u0440\u0430\u0431\u043e\u0442\u044b \u043f\u0440\u0438 \u043e\u043f\u0430\u0441\u043d\u043e\u0441\u0442\u0438"),
    ORANGE,
)


# Navigation
main_click_shapes[0].click_action.target_slide = s2
main_click_shapes[1].click_action.target_slide = s3
main_click_shapes[2].click_action.target_slide = s4
main_click_shapes[3].click_action.target_slide = s5
back_btn.click_action.target_slide = s1

next_btn.click_action.target_slide = s2

back_to_s003 = [
    add_button(s2, Inches(0.7), Inches(6.55), Inches(2.35), Inches(0.5), u(r"\u041d\u0410\u0417\u0410\u0414"), STEEL, 18),
    add_button(s3, Inches(0.7), Inches(6.55), Inches(2.35), Inches(0.5), u(r"\u041d\u0410\u0417\u0410\u0414"), STEEL, 18),
    add_button(s4, Inches(0.7), Inches(6.55), Inches(2.35), Inches(0.5), u(r"\u041d\u0410\u0417\u0410\u0414"), STEEL, 18),
    add_button(s5, Inches(0.7), Inches(6.55), Inches(2.35), Inches(0.5), u(r"\u041d\u0410\u0417\u0410\u0414"), STEEL, 18),
]

for button in back_to_s003:
    button.click_action.target_slide = s1


prs.save(OUT_FILE)
print(OUT_FILE)
