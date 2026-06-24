from pathlib import Path

from pptx import Presentation
from pptx.util import Inches


BASE_DIR = Path(__file__).resolve().parent
SRC = BASE_DIR / "S001-S007_live_preview_working_2026-06-24_v28.pptx"
OUT = BASE_DIR / "S001-S007_live_preview_working_2026-06-24_v29.pptx"

S006_INDEX = 16

LEFT_PANEL_X = Inches(0.7)
LEFT_PANEL_Y = Inches(1.45)
LEFT_PANEL_W = Inches(11.9)
HEAD_H = Inches(0.72)
BODY_X = LEFT_PANEL_X + Inches(0.28)
BODY_Y = Inches(2.45)
BODY_W = LEFT_PANEL_W - Inches(0.56)


def remove_shape(shape) -> None:
    shape._element.getparent().remove(shape._element)


def main() -> None:
    prs = Presentation(SRC)
    slide = prs.slides[S006_INDEX]

    for shape in list(slide.shapes):
        if shape.left >= Inches(6.2) and Inches(1.2) <= shape.top < Inches(5.9):
            remove_shape(shape)

    for shape in slide.shapes:
        text = getattr(shape, "text", "").strip()
        if text == "Квалификация стропальщика":
            shape.left = LEFT_PANEL_X
            shape.top = LEFT_PANEL_Y
            shape.width = LEFT_PANEL_W
            shape.height = HEAD_H
        elif text.startswith("2 разряд - простые грузы до 5 т."):
            shape.left = BODY_X
            shape.top = BODY_Y
            shape.width = BODY_W

    for shape in slide.shapes:
        if (
            not getattr(shape, "text", "").strip()
            and shape.left == LEFT_PANEL_X
            and shape.top == LEFT_PANEL_Y
        ):
            shape.width = LEFT_PANEL_W

    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    main()
