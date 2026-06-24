from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


BASE_DIR = Path(__file__).resolve().parent
SRC = BASE_DIR / "S001-S007_live_preview_working_2026-06-24_v29.pptx"
OUT = BASE_DIR / "S001-S007_live_preview_working_2026-06-24_v30.pptx"

S006_INDEX = 16

NAVY = RGBColor(18, 43, 66)
BLUE = RGBColor(38, 90, 137)
GREEN = RGBColor(66, 133, 95)
STEEL = RGBColor(83, 119, 149)
PALE = RGBColor(240, 244, 247)
ORANGE = RGBColor(222, 116, 54)
WHITE = RGBColor(255, 255, 255)
TEXT = RGBColor(33, 43, 54)
LIGHT_ORANGE = RGBColor(250, 242, 236)
LIGHT_BLUE = RGBColor(235, 242, 248)


def remove_shape(shape) -> None:
    shape._element.getparent().remove(shape._element)


def set_cell_text(cell, text, size, bold=False, color=TEXT, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE):
    cell.text = ""
    cell.fill.solid()
    tf = cell.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    p.space_after = Pt(0)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    cell.margin_left = Inches(0.07)
    cell.margin_right = Inches(0.07)
    cell.margin_top = Inches(0.03)
    cell.margin_bottom = Inches(0.03)


def set_fill(cell, color):
    cell.fill.solid()
    cell.fill.fore_color.rgb = color


def main() -> None:
    prs = Presentation(SRC)
    slide = prs.slides[S006_INDEX]

    for shape in list(slide.shapes):
        text = getattr(shape, "text", "").strip()
        if text == "Разряды 2-4":
            shape.text_frame.text = "Квалификация стропальщика"
        elif text == "Квалификация стропальщика":
            shape.text_frame.text = "Разряды 2-6"
        elif text.startswith("2 разряд - простые грузы до 5 т."):
            remove_shape(shape)

    table = slide.shapes.add_table(10, 4, Inches(0.82), Inches(2.08), Inches(11.52), Inches(3.98)).table

    col_widths = [1.28, 3.35, 4.25, 2.64]
    row_heights = [0.54, 0.34, 0.34, 0.34, 0.39, 0.39, 0.39, 0.39, 0.39, 0.42]

    for idx, width in enumerate(col_widths):
        table.columns[idx].width = Inches(width)
    for idx, height in enumerate(row_heights):
        table.rows[idx].height = Inches(height)

    headers = [
        "Разряд",
        "Допустимый\nвес груза",
        "Сложность груза\nи конструкций",
        "Длина лесных\nгрузов",
    ]
    for col, text in enumerate(headers):
        cell = table.cell(0, col)
        set_fill(cell, BLUE)
        set_cell_text(cell, text, 11, True, WHITE, PP_ALIGN.CENTER)

    grade_cells = {
        1: ("2 разряд", LIGHT_BLUE),
        2: ("3 разряд", PALE),
        4: ("4 разряд", PALE),
        7: ("5 разряд", LIGHT_ORANGE),
        9: ("6 разряд", LIGHT_ORANGE),
    }
    for row, (text, fill) in grade_cells.items():
        cell = table.cell(row, 0)
        set_fill(cell, fill)
        set_cell_text(cell, text, 10.5, True, TEXT, PP_ALIGN.CENTER)

    table.cell(2, 0).merge(table.cell(3, 0))
    table.cell(4, 0).merge(table.cell(6, 0))
    table.cell(7, 0).merge(table.cell(8, 0))

    weight_rows = {
        1: "До 5 тонн",
        2: "От 5 до 25 тонн (простые)",
        3: "До 5 тонн (средние)",
        4: "Свыше 25 тонн (простые)",
        5: "От 5 до 25 тонн (средние)",
        6: "До 5 тонн (сложные)",
        7: "Свыше 25 тонн (средние)",
        8: "От 5 до 50 тонн (сложные)",
        9: "Свыше 50 тонн",
    }
    for row, text in weight_rows.items():
        cell = table.cell(row, 1)
        set_fill(cell, WHITE)
        set_cell_text(cell, text, 10, False, TEXT, PP_ALIGN.LEFT)

    complexity_rows = {
        1: "Простые изделия и\nдетали",
        2: "Простые и средней\nсложности",
        4: "Средние, штучные,\nтребующие\nосторожности",
        7: "Сложные конструкции,\nстапельная сборка",
        9: "Особо\nответственные",
    }
    for row, text in complexity_rows.items():
        cell = table.cell(row, 2)
        set_fill(cell, WHITE)
        set_cell_text(cell, text, 9.5, row == 9, TEXT, PP_ALIGN.LEFT)

    table.cell(2, 2).merge(table.cell(3, 2))
    table.cell(4, 2).merge(table.cell(6, 2))
    table.cell(7, 2).merge(table.cell(8, 2))

    length_rows = {
        1: "До 3 метров",
        2: "До 3 метров",
        4: "Свыше 3 метров",
        7: "От 3 до 6 метров",
        9: "Свыше 6 метров",
    }
    for row, text in length_rows.items():
        cell = table.cell(row, 3)
        set_fill(cell, WHITE)
        set_cell_text(cell, text, 10, False, TEXT, PP_ALIGN.CENTER)

    table.cell(2, 3).merge(table.cell(3, 3))
    table.cell(4, 3).merge(table.cell(6, 3))
    table.cell(7, 3).merge(table.cell(8, 3))

    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    main()
