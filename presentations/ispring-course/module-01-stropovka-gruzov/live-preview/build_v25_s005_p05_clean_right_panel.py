from copy import deepcopy
from pathlib import Path
from zipfile import ZIP_DEFLATED, ZipFile
import xml.etree.ElementTree as ET

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


SRC = Path(__file__).resolve().parent / "S001-S007_live_preview_working_2026-06-23_v22.pptx"
OUT = Path(__file__).resolve().parent / "S001-S007_live_preview_working_2026-06-23_v25.pptx"

TEXT = RGBColor(33, 43, 54)
MUTED = RGBColor(95, 109, 123)
ORANGE = RGBColor(222, 116, 54)
WHITE = RGBColor(255, 255, 255)

P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
NS = {"p": P_NS, "a": A_NS}

MIN_WIDE_PANEL_CX = 3_000_000
HEADER_MAX_CY = 900_000
BODY_ADJ = "3050"
HEADER_ADJ = "19300"
TITLE_INSET = "160000"


def set_rounding(av_lst: ET.Element, value: str) -> None:
    for gd in list(av_lst.findall("a:gd", NS)):
        if gd.attrib.get("name") == "adj":
            av_lst.remove(gd)
    ET.SubElement(av_lst, f"{{{A_NS}}}gd", {"name": "adj", "fmla": f"val {value}"})


def fix_rounding(pptx_path: Path) -> None:
    tmp_path = pptx_path.with_suffix(".tmp")
    with ZipFile(pptx_path, "r") as zin, ZipFile(tmp_path, "w", compression=ZIP_DEFLATED) as zout:
        for info in zin.infolist():
            data = zin.read(info.filename)
            if info.filename.startswith("ppt/slides/slide") and info.filename.endswith(".xml"):
                root = ET.fromstring(data)
                changed = False
                for shape in root.findall(".//p:sp", NS):
                    sp_pr = shape.find("p:spPr", NS)
                    tx_body = shape.find("p:txBody", NS)
                    if sp_pr is None:
                        continue

                    geom = sp_pr.find("a:prstGeom", NS)
                    xfrm = sp_pr.find("a:xfrm", NS)
                    if geom is None or xfrm is None or geom.attrib.get("prst") != "roundRect":
                        continue

                    ext = xfrm.find("a:ext", NS)
                    if ext is None:
                        continue
                    cx = int(ext.attrib.get("cx", "0"))
                    cy = int(ext.attrib.get("cy", "0"))
                    if cx < MIN_WIDE_PANEL_CX:
                        continue

                    av_lst = geom.find("a:avLst", NS)
                    if av_lst is None:
                        av_lst = ET.SubElement(geom, f"{{{A_NS}}}avLst")

                    if cy <= HEADER_MAX_CY:
                        set_rounding(av_lst, HEADER_ADJ)
                        if tx_body is not None:
                            body_pr = tx_body.find("a:bodyPr", NS)
                            if body_pr is None:
                                body_pr = ET.SubElement(tx_body, f"{{{A_NS}}}bodyPr")
                            body_pr.set("lIns", TITLE_INSET)
                            body_pr.set("rIns", TITLE_INSET)
                        changed = True
                    else:
                        set_rounding(av_lst, BODY_ADJ)
                        changed = True

                if changed:
                    data = ET.tostring(root, encoding="utf-8", xml_declaration=True)

            zout.writestr(deepcopy(info), data)

    pptx_path.unlink()
    tmp_path.rename(pptx_path)


def build_text(box) -> None:
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    p.space_after = Pt(6)
    r = p.add_run()
    r.text = (
        "Главные требования к общепроизводственным СИЗ "
        "(спецодежде, обуви, перчаткам, очкам и каскам) закреплены в "
        "Техническом регламенте Таможенного союза ТР ТС 019/2011 "
        "«О безопасности средств индивидуальной защиты»."
    )
    r.font.size = Pt(12.5)
    r.font.color.rgb = TEXT

    p = tf.add_paragraph()
    p.alignment = PP_ALIGN.LEFT
    p.space_after = Pt(6)
    r = p.add_run()
    r.text = "СИЗ должны отвечать пяти ключевым критериям:"
    r.font.size = Pt(12.5)
    r.font.bold = True
    r.font.color.rgb = TEXT

    p = tf.add_paragraph()
    p.alignment = PP_ALIGN.LEFT
    p.space_after = Pt(6)
    r = p.add_run()
    r.text = (
        "безопасность самого изделия / защитная эффективность / "
        "эргономика и комфорт / стойкость к уходу / обязательная маркировка"
    )
    r.font.size = Pt(12.5)
    r.font.color.rgb = TEXT

    p = tf.add_paragraph()
    p.alignment = PP_ALIGN.LEFT
    p.space_after = Pt(0)
    r = p.add_run()
    r.text = (
        "На каждом СИЗ (или его ярлыке) должна быть понятная, несмываемая "
        "маркировка: наименование, размер, дата изготовления, срок годности, "
        "защитные свойства и единый знак обращения продукции на рынке (EAC)."
    )
    r.font.size = Pt(12.5)
    r.font.color.rgb = MUTED


def main() -> None:
    prs = Presentation(SRC)
    slide = prs.slides[15]

    header = slide.shapes[6]
    body_text = slide.shapes[7]
    build_text(body_text)

    panel_left = header.left
    panel_top = header.top
    panel_width = header.width
    panel_height = (body_text.top + body_text.height + Inches(0.18)) - panel_top

    outline = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        panel_left,
        panel_top,
        panel_width,
        panel_height,
    )
    outline.fill.solid()
    outline.fill.fore_color.rgb = WHITE
    outline.line.color.rgb = ORANGE
    outline.line.width = Pt(1.4)
    header._element.addprevious(outline._element)

    prs.save(OUT)
    fix_rounding(OUT)
    print(OUT)


if __name__ == "__main__":
    main()
